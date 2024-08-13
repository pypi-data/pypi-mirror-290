"""
Author: HECE - University of Liege, Pierre Archambeau
Date: 2024

Copyright (c) 2024 University of Liege. All rights reserved.

This script and its content are protected by copyright law. Unauthorized
copying or distribution of this file, via any medium, is strictly prohibited.
"""

try:
    import numpy as np
    from wx import dataview, TreeCtrl
    import wx
    import wx.propgrid as pg
    # from wxasync import AsyncBind
    from wx.core import VERTICAL, BoxSizer, Height, ListCtrl, StaticText, TextCtrl, Width
    from wx.glcanvas import GLCanvas, GLContext
    from wx.dataview import TreeListCtrl
    import wx.lib.ogl as ogl
    from PIL import Image, ImageOps
    from PIL.PngImagePlugin import PngInfo
    import io
    import json
    import glob
    import traceback
except ImportError as e:
    print(e)
    raise ImportError("Error importing wxPython, numpy, PIL, json, glob, traceback. Please check your installation.")

try:
    from osgeo import gdal
except ModuleNotFoundError:
    raise Exception("I can't find the 'gdal' package. You should get it from https://www.lfd.uci.edu/~gohlke/pythonlibs/")

try:
    from time import sleep
    from datetime import timedelta
    from multiprocessing import Pool
    from pathlib import Path
    from time import sleep
except ImportError as e:
    print(e)
    raise ImportError("Error importing time, datetime, multiprocessing, pathlib. Please check your installation.")

try:
    from OpenGL.GL import *
    from OpenGL.GLUT import *
except ImportError as e:
    msg=_('Error importing OpenGL library')
    msg+=_('   Python version : ' + sys.version)
    msg+=_('   Please check your version of opengl32.dll -- conflict may exist between different files present on your desktop')
    raise Exception(msg)

try:
    import matplotlib.pyplot as plt
    from matplotlib.widgets import Button as mplButton
    from matplotlib.ticker import FormatStrFormatter
    from os import scandir, listdir
    from os.path import exists, join, normpath
    from pptx import Presentation
    import threading
    from enum import Enum
    from typing import Literal, Union
    import logging
except ImportError as e:
    print(e)
    raise ImportError("Error importing matplotlib, os, threading, enum, typing, logging. Please check your installation.")

try:
    from .wolf_texture import genericImagetexture,imagetexture,Text_Image_Texture
    from .xyz_file import xyz_scandir, XYZFile
    from .mesh2d import wolf2dprev
    from .PyPalette import wolfpalette
    from .wolfresults_2D import Wolfresults_2D, views_2D
    from .PyTranslate import _
    from .PyVertex import cloud_vertices, getIfromRGB
    from .RatingCurve import SPWMIGaugingStations, SPWDCENNGaugingStations
    from .wolf_array import WOLF_ARRAY_MB, SelectionData, WolfArray, WolfArrayMB, CropDialog, header_wolf, WolfArrayMNAP, WOLF_ARRAY_FULL_SINGLE, WOLF_ARRAY_FULL_INTEGER8, WOLF_ARRAY_FULL_INTEGER16, WOLF_ARRAY_FULL_DOUBLE, WOLF_ARRAY_FULL_INTEGER
    from .PyParams import Wolf_Param, key_Param, Type_Param
    from .mesh2d.bc_manager import BcManager
    from .PyVertexvectors import *
    from .Results2DGPU import wolfres2DGPU
    from .PyCrosssections import crosssections, profile, Interpolator, Interpolators
    from .GraphNotebook import PlotNotebook
    from .lazviewer.laz_viewer import myviewer, read_laz, clip_data_xyz, xyz_laz_grids, choices_laz_colormap, Classification_LAZ
    from . import Lidar2002
    from .picc import Picc_data, Cadaster_data
    from .wolf_zi_db import ZI_Databse_Elt, PlansTerrier
    from .math_parser.calculator import Calculator
    from .wintab.wintab import Wintab
except ImportError as e:
    print(e)
    raise ImportError("Error importing wolf_texture, xyz_file, mesh2d, PyPalette, wolfresults_2D, PyTranslate, PyVertex, RatingCurve, wolf_array, PyParams, mesh2d.bc_manager, PyVertexvectors, Results2DGPU, PyCrosssections, GraphNotebook, lazviewer, picc, wolf_zi_db, math_parser.calculator, wintab. Please check your installation.")

try:
    from .hydrometry.kiwis_wolfgui import hydrometry_wolfgui
except ImportError as e:
    print(e)
    raise ImportError("Error importing hydrometry.kiwis_wolfgui. Please check your installation.")

try:
    from .pyshields import get_d_cr
    from .pyviews import WolfViews
    from .PyConfig import handle_configuration_dialog, WolfConfiguration, ConfigurationKeys
    from .GraphProfile import ProfileNotebook
    from .pybridges import Bridges, Bridge, Weirs, Weir
    from .tools_mpl import *
    from .wolf_tiles import Tiles
    from .lagrangian.particle_system_ui import Particle_system_to_draw as Particle_system
    from .opengl.py3d import Wolf_Viewer3D
    from .pyGui1D import GuiNotebook1D
except ImportError as e:
    print(e)
    raise ImportError("Error importing pyshields, pyviews, PyConfig, GraphProfile, pybridges, tools_mpl, wolf_tiles, lagrangian.particle_system_ui, opengl.py3d, pyGui1D. Please check your installation.")

ID_SELECTCS = 1000
ID_SORTALONG = 1001
ID_LOCMINMAX = 1002
ID_PLOTCS = 1003   #Manageactions ID for profile plots

LIST_1TO9 = [wx.WXK_NUMPAD1, wx.WXK_NUMPAD2, wx.WXK_NUMPAD3, wx.WXK_NUMPAD4, wx.WXK_NUMPAD5, wx.WXK_NUMPAD6, wx.WXK_NUMPAD7, wx.WXK_NUMPAD8, wx.WXK_NUMPAD9 ] + [ord(str(cur)) for cur in range(1,10)]

class draw_type(Enum):
    # FIXME: change this to be more robust -> Done !
    # Be careful with the enum name, it must be the same than the one used to create the tree list elements, but in lower case
    # see : self.treelist.AppendItem in __init__
    ARRAYS = 'arrays'
    BRIDGES= 'bridges'
    WEIRS = 'weirs'
    VECTORS = 'vectors'
    CLOUD = 'clouds'
    TRIANGULATION = 'triangulations'
    PARTICLE_SYSTEM = 'particle systems'
    CROSS_SECTIONS = 'cross_sections'
    OTHER = 'others'
    VIEWS = 'views'
    RES2D = 'wolf2d'
    WMSBACK = 'wms-background'
    WMSFORE = 'wms-foreground'
    TILES = 'tiles'


class DragdropFileTarget(wx.FileDropTarget):
    def __init__(self, window:"WolfMapViewer"):
        wx.FileDropTarget.__init__(self)
        self.window = window

    def OnDropFiles(self, x, y, filenames):

        def test_if_array(filename):

            ext = Path(filename).suffix

            if ext.lower() in ['.bin', '.npy', '.hbin', '.qxin','.qybin', '.top',
                               '.kbin', '.epsbin', '.tif', '.tiff', '.frot', '.topini_fine']:
                return True
            else:
                return False

        def test_if_arrayMB(filename):

            ext = Path(filename).suffix

            if ext.lower() in ['.hbinb', '.qxbinb','.qybinb', '.kbinb',
                               '.epsbinb', '.topini', '.frotini']:
                return True
            else:
                return False

        def test_if_vector(filename):
            ext = Path(filename).suffix

            if ext.lower() in ['.vec', '.vecz', '.shp', '.dxf']:
                return True
            else:
                return False

        def test_if_cloud(filename):
            ext = Path(filename).suffix

            if ext.lower() in ['.xyz']:
                return True
            else:
                return False

        for name in filenames:

            if Path(name).is_dir():
                for file in scandir(name):
                    if file.is_file():
                        self.OnDropFiles(x, y, [file.path])
                continue

            if test_if_array(name):
                ids = self.window.get_list_keys(draw_type.ARRAYS)
                id = Path(name).stem
                while id in ids:
                    id = id + '_1'

                try:
                    newobj = WolfArray(fname=name, mapviewer= self.window)
                    self.window.add_object('array', newobj = newobj, id  = id)
                except:
                    logging.error(_('Error while loading array : ') + name)

            elif test_if_arrayMB(name):
                ids = self.window.get_list_keys(draw_type.ARRAYS)
                id = Path(name).stem
                while id in ids:
                    id = id + '_1'

                try:
                    newobj = WolfArrayMB(fname=name, mapviewer= self.window)
                    self.window.add_object('array', newobj = newobj, id  = id)
                except:
                    logging.error(_('Error while loading array : ') + name)

            elif test_if_vector(name):
                ids = self.window.get_list_keys(draw_type.VECTORS)
                id = Path(name).stem
                while id in ids:
                    id = id + '_1'

                try:
                    newobj = Zones(filename=name, parent=self.window, mapviewer=self.window)
                    self.window.add_object('vector', newobj = newobj, id  = id)
                except:
                    logging.error(_('Error while loading vector : ') + name)

            elif test_if_cloud(name):
                ids = self.window.get_list_keys(draw_type.CLOUD)
                id = Path(name).stem
                while id in ids:
                    id = id + '_1'

                try:
                    newobj = cloud_vertices(fname=name, mapviewer=self.window)
                    self.window.add_object('cloud', newobj = newobj, id  = id)
                except:
                    logging.error(_('Error while loading cloud : ') + name)

        return True

class WolfMapViewer(wx.Frame):
    """
    Fenêtre de visualisation de données WOLF grâce aux WxWidgets
    """

    TIMER_ID = 100  # délai d'attente avant action

    mybc: list[BcManager]  # Gestionnaire de CL
    myarrays: list  # matrices ajoutées
    myvectors: list[Zones]  # zones vectorielles ajoutées
    myclouds: list[cloud_vertices]  # nuages de vertices
    mytri: list[Triangulation]  # triangulations
    myothers: list
    myviews:list[views_2D]
    mywmsback: list
    mywmsfore: list
    myres2D: list
    mytiles: list[Tiles]
    mypartsystems: list[Particle_system]
    myviewers3d:list[Wolf_Viewer3D]

    canvas: GLCanvas  # canvas OpenGL
    context: GLContext  # context OpenGL
    mytooltip: Wolf_Param  # Objet WOLF permettant l'analyse de ce qui est sous la souris
    treelist: TreeListCtrl  # Gestion des éléments sous forme d'arbre
    _lbl_selecteditem: StaticText
    leftbox: BoxSizer

    # DEPRECEATED
    # added: dict  # dictionnaire des éléments ajoutés

    active_vector: vector
    active_zone: zone
    active_zones: Zones
    active_array: WolfArray
    active_bc: BcManager
    active_view: WolfViews
    active_vertex: wolfvertex
    active_cs: crosssections
    active_tri: Triangulation
    active_tile: Tiles
    active_particle_system: Particle_system
    active_viewer3d: Wolf_Viewer3D

    def __init__(self,
                 wxparent = None,
                 title:str = _('Default Wolf Map Viewer'),
                 w:int=500,
                 h:int=500,
                 treewidth:int=200,
                 wolfparent=None,
                 wxlogging=None):

        """
        Create a Viewer for WOLF data/simulation

        :params wxparent: wx parent - set to None if main window
        :params title: title of the window
        :params w: width of the window in pixels
        :params h: height of the window in pixels
        :params treewidth: width of the tree list in pixels
        :params wolfparent: WOLF object parent -- see PyGui.py
        :params wxlogging: wx logging object

        """

        self._wxlogging = wxlogging
        self.action = None  # Action à entreprendre
        self.update_absolute_minmax = False  # Force la MAJ de la palette
        self.copyfrom = None  # aucun élément pointé par CTRL+C

        self.wolfparent = wolfparent

        self.regular = True  # Gestion de la taille de fenêtre d'affichage, y compris l'arbre de gestion
        self.sx = 1  # facteur d'échelle selon X = largeur en pixels/largeur réelle
        self.sy = 1  # facteur d'échelle selon Y = hauteur en pixels/hauteur réelle
        self.samescale = True  # force le même facteur d'échelle

        self.dynapar_dist = 1.

        # emprise initiale
        self.xmin = 0.
        self.ymin = 0.
        self.xmax = 40.
        self.ymax = 40.
        self.width = self.xmax - self.xmin  # largeur de la zone d'affichage en coordonnées réelles
        self.height = self.ymax - self.ymin  # hauteur de la zone d'affichage en coordonnées réelles
        self.canvaswidth = 100
        self.canvasheight = 100

        # position de la caméra
        self.mousex = self.width / 2.
        self.mousey = self.height / 2.
        self._last_mouse_pos = (0, 0, (0,0))

        self.bordersize = 0  # zone réservée au contour
        self.titlesize = 0  # zone réservée au titre
        self.treewidth = 200  # largeur de la zone d'arbre "treelist"

        self.backcolor = wx.Colour(255, 255, 255)  # couleur de fond
        self.mousedown = (0., 0.)  # position initiale du bouton position bas
        self.mouseup = (0., 0.)  # position initiale du bouton position haut
        self.oneclick = True  # détection d'un simple click ou d'un double-click
        self.move = False  # la souris est-elle en train de bouger?

        self.linked = False
        self.link_shareopsvect = True
        self.linkedList = None
        self.link_params = None

        self.forcemimic = True
        self.currently_readresults = False

        self.mylazdata = None
        # self.mylazdata_colors = None
        self.mylazgrid = None

        self.treewidth = treewidth
        super(WolfMapViewer, self).__init__(wxparent, title=title, size=(w + self.treewidth, h))

        self._dragdrop = DragdropFileTarget(self)
        self.SetDropTarget(self._dragdrop)

        # Gestion des menus
        self.popupmenu = wx.Menu()
        self.popupmenu.Bind(wx.EVT_MENU, self.OnPopupItemSelected)

        for text in [_('Save'), _('Save as'), _('Rename'), _('Duplicate'), _('Up'), _('Down'), _('Properties')]:
            item = self.popupmenu.Append(-1, text)

        self.menubar = wx.MenuBar()

        self.menuwolf2d = None
        self.menu_landmap = None
        self.menu2d_cache_setup = None
        self.menuparticlesystem = None
        self.menu2dGPU = None
        self.menuwalous = None
        self.timer_ps = None
        self.menusim2D  = None
        self.menulaz    = None
        self.menutiles  = None

        self.filemenu = wx.Menu()
        openitem = self.filemenu.Append(wx.ID_OPEN, _('Open project'), _('Open a complete project from file'))
        saveproject = self.filemenu.Append(wx.ID_ANY, _('Save project'), _('Save the current project to file'))
        self.filemenu.AppendSeparator()
        saveitem = self.filemenu.Append(wx.ID_SAVE, _('Save'), _('Save all checked arrays or vectors to files'))
        saveasitem = self.filemenu.Append(wx.ID_SAVEAS, _('Save as...'), _('Save all checked arrays or vectors to new files --> one file dialog per data'))
        savecanvas = self.filemenu.Append(wx.ID_ANY, _('Save to image...'), _('Save the canvas to image file on disk'))
        copycanvas = self.filemenu.Append(wx.ID_ANY, _('Copy image...'), _('Copy the canvas to image file to the clipboard'))

        self.filemenu.AppendSeparator()
        # --- GLTF
        self.menugltf = wx.Menu()
        self.filemenu.Append(wx.ID_ANY,_('Gltf2...'), self.menugltf)

        exportgltf = self.menugltf.Append(wx.ID_ANY, _('Export...'), _('Save data to gltf files'))
        importgltf = self.menugltf.Append(wx.ID_ANY, _('Import...'), _('Import data from gltf files'))
        compareitem = self.menugltf.Append(wx.ID_ANY, _('Compare...'), _('Create new frames to compare sculpting'))
        updategltf = self.menugltf.Append(wx.ID_ANY, _('Update...'), _('Update data from gltf files'))

        self.filemenu.AppendSeparator()

        # SIMULATION 2D

        sim2d = self.filemenu.Append(wx.ID_ANY, _('Create/Open multiblock model'), _('Create or open a multiblock model in a new viewer'))
        check2D = self.filemenu.Append(wx.ID_ANY, _('Check headers'), _('Check the header .txt files from an existing 2D simulation'))

        self.filemenu.AppendSeparator()

        # SIMULATION Hydrologique

        hydrol = self.filemenu.Append(wx.ID_ANY, _('Open hydrological model'), _('Hydrological simulation'))

        self.filemenu.AppendSeparator()

        # MULTIVIEWER

        compareitem = self.filemenu.Append(wx.ID_ANY, _('Set comparison'), _('Set comparison'))
        multiview = self.filemenu.Append(wx.ID_ANY, _('Multiviewer'), _('Multiviewer'))
        viewer3d = self.filemenu.Append(wx.ID_ANY, _('3D viewer'), _('3D viewer'))
        self.filemenu.AppendSeparator()


        # ---
        self.menucreateobj = wx.Menu()
        self.filemenu.Append(wx.ID_ANY,_('Create...'),self.menucreateobj)

        createarray = self.menucreateobj.Append(wx.ID_FILE6, _('Create array...'), _('New array (binary file - real)'))
        createarray2002 = self.menucreateobj.Append(wx.ID_ANY, _('Create array from Lidar 2002...'),
                                               _('Create array from Lidar 2002 (binary file - real)'))
        createarrayxyz = self.menucreateobj.Append(wx.ID_ANY, _('Create array from bathymetry file...'),
                                              _('Create array from XYZ (ascii file - real)'))
        createvector = self.menucreateobj.Append(wx.ID_FILE7, _('Create vectors...'), _('New vectors'))
        createview = self.menucreateobj.Append(wx.ID_ANY, _('Create view...'), _('New view'))
        createcloud = self.menucreateobj.Append(wx.ID_FILE8, _('Create cloud...'), _('New cloud'))
        createmanager2D = self.menucreateobj.Append(wx.ID_ANY, _('Create Wolf2D manager ...'), _('New manager 2D'))
        createscenario2D = self.menucreateobj.Append(wx.ID_ANY, _('Create scenarios manager ...'), _('New scenarios manager 2D'))
        createbcmanager2D = self.menucreateobj.Append(wx.ID_ANY, _('Create BC manager Wolf2D...'), _('New BC manager 2D'))
        createpartsystem = self.menucreateobj.Append(wx.ID_ANY, _('Create particle system...'), _('Create a particle system - Lagrangian view'))
        create1Dmodel = self.menucreateobj.Append(wx.ID_ANY, _('Create Wolf1D...'),('Create a 1D model using crossections, vectors and arrays...'))
        create_acceptability = self.menucreateobj.Append(wx.ID_ANY, _('Create acceptability manager...'), _('Create acceptability manager'))

        self.filemenu.AppendSeparator()


        # -----
        self.menuaddobj = wx.Menu()
        self.filemenu.Append(wx.ID_ANY,_('Add...'),self.menuaddobj)

        addarray = self.menuaddobj.Append(wx.ID_FILE1, _('Add array...'), _('Add array (binary file - real)'))
        addarraycrop = self.menuaddobj.Append(wx.ID_ANY, _('Add array and crop...'),
                                            _('Add array and crop (binary file - real)'))
        addvector = self.menuaddobj.Append(wx.ID_FILE2, _('Add vectors...'), _('Add vectors'))
        addtiles = self.menuaddobj.Append(wx.ID_ANY, _('Add tiles...'), _('Add tiles'))
        addtilescomp = self.menuaddobj.Append(wx.ID_ANY, _('Add tiles comparator...'), _('Add tiles comparator'))
        addtiles = self.menuaddobj.Append(wx.ID_ANY, _('Add tiles GPU...'), _('Add tiles from 2D GPU model -- 2 arrays will be added'))
        addcloud = self.menuaddobj.Append(wx.ID_FILE3, _('Add cloud...'), _('Add cloud'))
        addtri = self.menuaddobj.Append(wx.ID_ANY, _('Add triangulation...'), _('Add triangulation'))
        addprofiles = self.menuaddobj.Append(wx.ID_FILE4, _('Add cross sections...'), _('Add cross sections'))
        addres2D = self.menuaddobj.Append(wx.ID_ANY, _('Add Wolf2D results...'), _('Add Wolf 2D results'))
        addres2Dgpu = self.menuaddobj.Append(wx.ID_ANY, _('Add Wolf2D GPU results...'), _('Add Wolf 2D GPU results'))
        addpartsystem = self.menuaddobj.Append(wx.ID_ANY, _('Add particle system...'), _('Add a particle system - Lagrangian view'))
        addbridges = self.menuaddobj.Append(wx.ID_ANY, _('Add bridges...'), _('Add bridges from directory'))
        addweirs = self.menuaddobj.Append(wx.ID_ANY, _('Add weirs...'), _('Add bridges from directory'))
        addview = self.menuaddobj.Append(wx.ID_ANY, _('Add view...'), _('Add view from project file'))

        self.filemenu.AppendSeparator()
        addscan = self.filemenu.Append(wx.ID_FILE5, _('Recursive scan...'), _('Add recursively'))

        # Tools
        # ----------------

        self.tools_menu = wx.Menu()

        self.menu_contour_from_arrays = self.tools_menu.Append(wx.ID_ANY, _("Create contour from checked arrays..."), _("Create contour"))
        self.menu_calculator = self.tools_menu.Append(wx.ID_ANY, _("Calculator..."), _("Calculator"))
        self.calculator = None

        # Cross sections
        # ----------------

        self.cs_menu = wx.Menu()
        self.link_cs_zones = self.cs_menu.Append(wx.ID_ANY, _("Link cross sections to active zones"),
                                                   _("Link cross section"))
        self.sortalong = self.cs_menu.Append(ID_SORTALONG, _("Sort along..."),
                                               _("Sort cross sections along support vector"))
        self.select_cs = self.cs_menu.Append(ID_SELECTCS, _("Pick one cross section"), _("Select cross section"),
                                               kind=wx.ITEM_CHECK)
        self.menumanagebanks = self.cs_menu.Append(wx.ID_ANY, _("Manage banks..."), _("Manage banks"))
        self.menucreatenewbanks = self.cs_menu.Append(wx.ID_ANY, _("Create banks from vertices..."),
                                                        _("Manage banks"))
        self.renamecs = self.cs_menu.Append(wx.ID_ANY, _("Rename cross sections..."), _("Rename"))
        self.menutrianglecs = self.cs_menu.Append(wx.ID_ANY, _("Triangulate cross sections..."), _("Triangulate"))
        self.menuexportgltfonebyone = self.cs_menu.Append(wx.ID_ANY, _("Export cross sections to gltf..."),
                                                            _("Export gltf"))
        self.menupontgltfonebyone = self.cs_menu.Append(wx.ID_ANY, _("Create bridge and export gltf..."),
                                                          _("Bridge gltf"))
        # self.menuimport3dfaces_from_DXF = self.toolsmenu.Append(wx.ID_ANY, _("Import triangulation..."), _("DXF"))
        self.menuinteractptri = self.cs_menu.Append(wx.ID_ANY, _("Interpolate on active triangulation..."), _("InterpolateTri"))
        self.menucomparecloud = self.cs_menu.Append(wx.ID_ANY, _("Compare cloud to array..."), _("Comparison"))
        self.menucomparetri = self.cs_menu.Append(wx.ID_ANY, _("Compare triangles to array..."), _("Comparison"))

        #Profile plots
        #The action for plotting cross section's profile is initialised.
        self.plot_cs = self.cs_menu.Append(ID_PLOTCS, _("Plot cross section"),_("Plot cross section"),kind=wx.ITEM_CHECK)

        self.menuviewerinterpcs = None
        self.menuinterpcs = None

        # COLORMAP  menu
        self.minmaxmenu = wx.Menu()
        self.locminmax = self.minmaxmenu.Append(ID_LOCMINMAX, _("Local minmax"), _("Adapt colormap on current zoom"),
                                                kind=wx.ITEM_CHECK)
        paluniform= self.minmaxmenu.Append(wx.ID_ANY, _("Compute and apply unique colormap on all..."),
                                           _("Unique colormap"))
        paluniform_fomfile= self.minmaxmenu.Append(wx.ID_ANY, _("Load and apply unique colormap on all..."),
                                           _("Unique colormap"))
        paluniform_inparts= self.minmaxmenu.Append(wx.ID_ANY, _("Force uniform in parts on all..."),
                                           _("Uniform in parts"))
        pallinear= self.minmaxmenu.Append(wx.ID_ANY, _("Force linear interpolation on all..."),
                                           _("Linear colormap"))

        self.analyzemenu = wx.Menu()
        plotvect = self.analyzemenu.Append(wx.ID_ANY, _("Plot active vector..."),
                                           _("Plot the active vector and linked arrays"))
        plotpoly = self.analyzemenu.Append(wx.ID_ANY, _("Plot active polygons..."),
                                           _("Plot the active polygons and linked arrays"))

        masksimul = self.analyzemenu.Append(wx.ID_ANY, _("Load and apply mask (nap)..."),
                                           _("Apply mask from sim2D"))
        filterinund = self.analyzemenu.Append(wx.ID_ANY, _("Filter inundation arrays..."),
                                           _("Filter arrays"))
        exporttif = self.analyzemenu.Append(wx.ID_ANY, _("Export arrays as Geotif..."),
                                           _("Export arrays as Geotif"))
        exportshape = self.analyzemenu.Append(wx.ID_ANY, _("Export arrays as Shapefile..."),
                                           _("Export arrays as Shapefile"))
        plotqvect = self.analyzemenu.Append(wx.ID_ANY, _("Integrate Q along active vector..."),
                                           _("Integrate Q along the active vector"))
        plotqvect = self.analyzemenu.Append(wx.ID_ANY, _("Integrate Q along active zone..."),
                                           _("Integrate Q along the active zone"))
        plothselect = self.analyzemenu.Append(wx.ID_ANY, _("Plot stats unknown (selected nodes)..."),
                                           _("Compute stats and plot on the selected nodes"))
        plothvector = self.analyzemenu.Append(wx.ID_ANY, _("Plot stats unknown (inside active vector)..."),
                                           _("Compute stats and plot on nodes inside the active vector"))
        plothzone = self.analyzemenu.Append(wx.ID_ANY, _("Plot stats unknown (inside active zone)..."),
                                           _("Compute stats and plot on nodes inside the active zone"))

        self.filemenu.AppendSeparator()
        menuquit = self.filemenu.Append(wx.ID_EXIT, _('&Quit\tCTRL+Q'), _('Quit application'))

        # If one uses the accelerator key then it is tied to the
        # wx.ID_EXIT. Moreover the accelerator key will be shadowed by
        # EVT_CHAR_HOOK if one is not careful.  Note that using
        # accelerators on anything else then wx.EVT_MENU is reported
        # hackish at best on the WWW.

        accel_tbl = wx.AcceleratorTable([(wx.ACCEL_CTRL,  ord('Q'), menuquit.GetId() )])
        self.SetAcceleratorTable(accel_tbl)

        # Gestion des outils --> Utile pour ManageActions
        self.tools = {}
        curtool = self.tools[ID_SELECTCS] = {}
        curtool['menu'] = self.select_cs
        curtool['name'] = 'Select nearest profile'

        curtool = self.tools[ID_PLOTCS] = {}
        curtool['menu'] = self.plot_cs
        curtool['name'] = 'Plot cross section'

        self.mybc = []

        self.active_vector = None
        self.active_zone = None
        self.active_zones = None
        self.active_vertex = None
        self.active_array = None
        self.active_bc = None
        self.active_tri = None
        self.active_cloud = None
        self.active_view = None
        self.active_cs = None
        self.active_profile = None
        self.active_res2d = None
        self.active_particle_system = None
        self.active_viewer3d = None
        self.active_landmap:PlansTerrier = None
        self.selected_treeitem = None

        curtool = self.tools[ID_SORTALONG] = {}
        curtool['menu'] = self.sortalong
        curtool['name'] = 'Sort along vector'

        curtool = self.tools[ID_LOCMINMAX] = {}
        curtool['menu'] = self.locminmax
        curtool['name'] = None

        self.menubar.Append(self.filemenu, _('&File'))

        # Help
        self.helpmenu = wx.Menu()
        self.helpmenu.Append(wx.ID_ANY, _('Shortcuts'), _('Shortcuts'))
        self.helpmenu.Append(wx.ID_ANY, _('Show logs/informations'), _('Logs'))
        self.helpmenu.Append(wx.ID_ANY, _('Show values'), _('Data/Values'))
        self.helpmenu.Append(wx.ID_ANY, _('About'), _('About'))
        self.helpmenu.Append(wx.ID_ANY, _('Check for updates'), _('Update?'))

        self.menubar.Append(self.helpmenu, _('&Help'))

        # ajout du menu pour les données LAZ
        self.menu_laz()

        self.menubar.Append(self.tools_menu, _('&Tools'))
        self.menubar.Append(self.cs_menu, _('&Cross sections'))

        self.menubar.Append(self.minmaxmenu, _('&Colormap'))
        self.menubar.Append(self.analyzemenu, _('&Analyze'))
        self.SetMenuBar(self.menubar)
        self.Bind(wx.EVT_MENU, self.OnMenubar)
        self.Bind(wx.EVT_MENU_HIGHLIGHT, self.OnMenuHighlight)

        # Ajout du conteneur OpenGL
        self.canvas = GLCanvas(self)
        self.canvas.SetDropTarget(self._dragdrop)

        self.context = GLContext(self.canvas)
        self.mybackisloaded = False
        self.myfrontisloaded = False

        # ajout d'une liste en arbre des objets
        self.treelist = TreeListCtrl(self, style= wx.dataview.TL_CHECKBOX | wx.LC_EDIT_LABELS | wx.TR_FULL_ROW_HIGHLIGHT)
        self._lbl_selecteditem = StaticText(self, style=wx.ALIGN_CENTER_HORIZONTAL)
        self.selected_object = None

        self.root = self.treelist.GetRootItem()
        self.treelist.AppendColumn(_('Objects to plot'))
        self.myitemsarray = self.treelist.AppendItem(self.root, _("Arrays"))
        self.myitemsvector = self.treelist.AppendItem(self.root, _("Vectors"))
        self.myitemscloud = self.treelist.AppendItem(self.root, _("Clouds"))
        self.myitemstri = self.treelist.AppendItem(self.root, _("Triangulations"))
        self.myitemsres2d = self.treelist.AppendItem(self.root, _("Wolf2D"))
        self.myitemsps = self.treelist.AppendItem(self.root, _("Particle systems"))
        self.myitemsothers = self.treelist.AppendItem(self.root, _("Others"))
        self.myitemsviews = self.treelist.AppendItem(self.root, _("Views"))
        self.myitemswmsback = self.treelist.AppendItem(self.root, _("WMS-background"))
        self.myitemswmsfore = self.treelist.AppendItem(self.root, _("WMS-foreground"))

        width, height = self.GetClientSize()
        self.bordersize = int((w - width + self.treewidth) / 2)
        self.titlesize = h - height - self.bordersize
        self.SetSize(w + 2 * self.bordersize + self.treewidth, h + self.bordersize + self.titlesize)

        # dimensionnement et positionnement de la fenêtre OpenGL
        self.canvas.SetSize(width - self.treewidth, height)
        self.canvas.SetPosition((self.treewidth, 0))

        self.setbounds()

        # dimensionnement et positionnement de l'arbre
        self.leftbox = BoxSizer(orient=wx.VERTICAL)
        self.leftbox.Add(self.treelist, 1, wx.LEFT)
        self.leftbox.Add(self._lbl_selecteditem, 0, wx.LEFT)
        self.treelist.SetSize(self.treewidth, height)


        self.CreateStatusBar(1)

        self.SetSizer(self.leftbox)

        # self.treelist.SetPosition((0,0))

        # fenêtre ToolTip
        self.mytooltip = Wolf_Param(self, _("Data/Results"), to_read=False, withbuttons=False, toolbar=False, DestroyAtClosing=False)
        self.mytooltip.SetSize(300, 400)
        self.mytooltip.prop.SetDescBoxHeight(20) # Hauteur de la zone de description
        self.mytooltip.Show(True)
        self._oldpos_tooltip = None

        #Notebooks
        self.notebookcs = None
        self.notebookprof = None
        self.notebookbanks = None

        #Axes
        self.myaxcs = None
        self.myaxprof = None

        #Figures
        self.myfigcs = None
        self.myfigprof = None

        self.cloudmenu=None
        self._configuration = None

        self.compare_results = None

        self.InitUI()

        # self._wintab = Wintab(self.GetHandle())

        # if self._wintab:

        #     import win32gui
        #     import win32con

        #     # self.oldWndProc = win32gui.SetWindowLong(self.GetHandle(), win32con.GWL_WNDPROC, self.MyWndProc)

    # def MyWndProc(self, hWnd, msg, wParam, lParam):
    #     import win32con

    #     # Intercept a specific Windows message (for example, WM_KEYDOWN)
    #     # if msg == 0x7FF0:
    #     #     key_code = wParam
    #     #     print(f"Key pressed: {key_code}")

    #     #     # Process the message or do something custom
    #     #     if key_code == win32con.VK_ESCAPE:
    #     #         print("Escape key pressed, intercepting the event.")

    #             # # You can return 0 to indicate the message has been processed
    #             # return 0
    #     # print(msg)
    #     return 0

    @property
    def wxlogging(self):
        return self._wxlogging

    @wxlogging.setter
    def wxlogging(self, value):
        self._wxlogging = value

    def check_logging(self):
        """ Check if logging window is shown """

        if self._wxlogging is None:
            logging.info(_('No logging window'))
            return

        self._wxlogging.Show()

    def check_tooltip(self):
        """ Check if tooltip window is shown """

        if self.mytooltip is None:
            logging.info(_('No tooltip window'))
            return

        self.mytooltip.Show()


    def open_hydrological_model(self):
        """ Open a hydrological model """

        from .PyGui import HydrologyModel

        newview = HydrologyModel(splash = False)

    def create_2D_MB_model(self):
        """ Create a 2D model """

        from .PyGui import Wolf2DModel

        newview = Wolf2DModel(splash = False)

    def check_2D_MB_headers(self):
        """ Check headers of a 2D simulation without opening viewer"""

        # Check 2D simulation
        dlg = wx.FileDialog(self, _("Choose 2D simulation file"), wildcard="all (*.*)|*.*", style=wx.FD_OPEN)
        if dlg.ShowModal() == wx.ID_CANCEL:
            dlg.Destroy()
            return

        filename = dlg.GetPath()
        dlg.Destroy()

        from .mesh2d.wolf2dprev import prev_sim2D

        sim = prev_sim2D(filename)
        sim.verify_files()


    def get_mapviewer(self):
        """ Retourne une instance WolfMapViewer """
        return self

    def do_quit(self):
        pass

    def create_cloud_menu(self):

        if self.cloudmenu is None:
            self.cloudmenu = wx.Menu()
            self.menubar.Append(self.cloudmenu, _('Cloud'))

            interpcloudonarray = self.cloudmenu.Append(wx.ID_ANY, _("Interpolate active cloud on active array..."),
                                           _("Interpolation"))

    def get_choices_arrays(self):
        """Boîte de dialogue permettant de choisir une ou plusieurs matrices parmi celles chargées"""

        dlg = wx.MultiChoiceDialog(self,_('Choose one or multiple arrays'),
                                    _('Choose'),choices=[cur.idx for cur in self.myarrays])
        ret = dlg.ShowModal()

        if ret == wx.ID_CANCEL:
            return None

        mychoices = self.myarrays[dlg.GetSelections()]
        return mychoices

    def menu_tiles(self):
        """ Menu for tiles """
        if self.menutiles is None:
            self.menutiles = wx.Menu()
            self.menubar.Append(self.menutiles, _('&Tiles'))

            picktiles = self.menutiles.Append(wx.ID_ANY, _('Pick a tile and load data'), _('Right click to pick a tile'))
            data_active_polygon_tiles = self.menutiles.Append(wx.ID_ANY, _('Select data within the active polygonal area'), _('Select data within the bouding box of the active polygonal area'))
            data_tmpvec_tiles = self.menutiles.Append(wx.ID_ANY, _('Select data within a temporary polygonal area'), _('Right click to add points + Enter'))

            self.Bind(wx.EVT_MENU, self.pick_tile, picktiles)
            self.Bind(wx.EVT_MENU, self.create_data_from_tiles_activevec, data_active_polygon_tiles)
            self.Bind(wx.EVT_MENU, self.create_data_from_tiles_tmpvec, data_tmpvec_tiles)

    def pick_tile(self, event: wx.Event):

        if self.active_tile is None:
            logging.warning(_('No active tile -- Please load data first'))
            return

        self.action = 'select active tile'
        logging.info(_('Select active tile'))

    def create_data_from_tiles_activevec(self, event: wx.Event):

        if self.active_tile is None:
            logging.warning(_('No active tile -- Please load data first'))
            return

        if self.active_vector is None:
            logging.warning(_('No active vector -- Please activate a vector first'))
            return

        self._create_data_from_tiles_common()

    def _create_data_from_tiles_common(self):

        from .wolf_vrt import create_vrt, crop_vrt

        dirdata = self.active_tile.linked_data_dir

        glob_vrt = glob.glob(join(dirdata,'*.vrt'))

        if len(glob_vrt) == 0:
            file_vrt = r'tmp.vrt'
            create_vrt(dirdata, fout=file_vrt)
        else:
            glob_vrt = glob_vrt[0]

        dlg = wx.FileDialog(None, _('Choose filename'), wildcard='tif (*.tif)|*.tif', defaultDir=dirdata, defaultFile='{}_crop.tif'.format(self.active_vector.myname), style=wx.FD_SAVE)
        ret = dlg.ShowModal()
        if ret == wx.ID_CANCEL:
            dlg.Destroy()
            return
        fout = dlg.GetPath()
        dlg.Destroy()

        bbox = self.active_vector.get_bounds_xx_yy()

        crop_vrt(glob_vrt, bbox, fout=fout)
        logging.info(_('File {} created').format(fout))

        dlg = wx.MessageDialog(self, _('Do you want to load the created file ?'), _('Load file'), wx.YES_NO | wx.ICON_QUESTION)
        ret = dlg.ShowModal()
        if ret == wx.ID_CANCEL:
            dlg.Destroy()
            return
        elif ret == wx.ID_YES:
            self.add_object('array', filename=fout, id = self.active_vector.myname)

    def create_data_from_tiles_tmpvec(self, event: wx.Event):

        if self.active_tile is None:
            logging.warning(_('No active tile -- Please load data first'))
            return

        self.start_action('create polygon - tiles', _('Extract data from tiles inside polygon'))

        self.active_vector = vector()
        self.active_vector.myname = 'crop_tiles'
        self.active_vector.add_vertex(wolfvertex(0.,0.))

    def menu_laz(self):
        """ Menu for LAZ Data """
        if self.menulaz is None:

            self.menulaz = wx.Menu()
            self.menulazdata = wx.Menu()
            self.menulazgrid = wx.Menu()

            self.menubar.Append(self.menulaz, _('&LAZ'))
            self.menulaz.AppendSubMenu(self.menulazdata, _('LAZ data'))
            self.menulaz.AppendSubMenu(self.menulazgrid, _('LAZ grid'))

            readlaz = self.menulazdata.Append(wx.ID_ANY, _('Initialize from npz'), _('LAZ from Numpy array'))
            readlaz_gridinfo = self.menulazgrid.Append(wx.ID_ANY, _('Initialize from GridInfos'), _('LAZ from gridinfos - subgridding of LAZ files'), kind=wx.ITEM_CHECK)

            bridgelaz = self.menulazdata.Append(wx.ID_ANY, _('Create cloud points from bridges'), _('LAZ Bridge'))
            buildinglaz = self.menulazdata.Append(wx.ID_ANY, _('Create cloud points from buildings'), _('LAZ Buildings'))

            croplaz = self.menulaz.Append(wx.ID_ANY, _('Clip LAZ grid on current zoom'), _('Select LAZ data based on the visible screen extent'),)
            viewlaz = self.menulaz.Append(wx.ID_ANY, _('Create LAZ viewer'), _('Create a LAZ Viewer based on loaded data'))
            aroundlaz = self.menulaz.Append(wx.ID_ANY, _('Plot LAZ around active vector'), _('Display a Matplotlib plot with the LAZ values around the active vector/polyline'),)
            pick_aroundlaz = self.menulaz.Append(wx.ID_ANY, _('Plot LAZ around temporary vector'), _('Display a Matplotlib plot with the LAZ values around a temporary vector/polyline -- Right clicks to add points + Enter'),)
            updatecolors_laz = self.menulaz.Append(wx.ID_ANY, _('Change colors - Classification'), _('Change color map associated to the current classification'),)
            fillarray_laz = self.menulaz.Append(wx.ID_ANY, _('Fill active array from LAZ data'), _('Fill an array from the LAZ data'),)
            selectarray_laz = self.menulaz.Append(wx.ID_ANY, _('Select cells in array from LAZ data'), _('Select nodes in active array from the LAZ data'),)

    def menu_wolf2d(self):

        if self.menuwolf2d is None:
            self.menuwolf2d = wx.Menu()

            self.menu2d_curentview = self.menuwolf2d.Append(wx.ID_ANY, _("Change current view"), _("Current view"))
            self.menu2d_lastres = self.menuwolf2d.Append(wx.ID_ANY, _("Read last result"), _("Current view"))
            self.menu2d_epsilon = self.menuwolf2d.Append(wx.ID_ANY, _("Set epsilon water depth"), _("Set the epsilon used in the mask"))

            self.menu_filter_independent = self.menuwolf2d.Append(wx.ID_ANY, _("Filter independent"), _("Filter independent"), kind=wx.ITEM_CHECK)

            # self.menu2d_bc = self.menuwolf2d.Append(wx.ID_ANY, _("Manage boundary conditions..."), _("BC manager"))
            self.menu2d_video = self.menuwolf2d.Append(wx.ID_ANY, _("Create video..."), _("Video/Movie"))

            self.menuwolf2d.AppendSeparator()

            self.menu2d_dangermap = self.menuwolf2d.Append(wx.ID_ANY, _("Danger map"), _("Compute the danger map"))
            self.menu2d_dangermaph = self.menuwolf2d.Append(wx.ID_ANY, _("Danger map - only h"), _("Compute the danger map"))

            self.menuwolf2d.AppendSeparator()

            self.menubar.Append(self.menuwolf2d, _('Results 2D'))

            self.menuwolf2d.Bind(wx.EVT_MENU, self.Onmenuwolf2d)

    def menu_walous(self):

        if self.menuwalous is None:
            self.menuwalous = wx.Menu()

            self.menuwalous_crop = self.menuwalous.Append(wx.ID_ANY, _("Crop on active array"), _("Crop active array"))
            self.menuwalous_cropscreen = self.menuwalous.Append(wx.ID_ANY, _("Crop on screen"), _("Crop screen"))
            self.menuwalous_map = self.menuwalous.Append(wx.ID_ANY, _("Map active array"), _("Map active array"))
            self.menuwalous_legend = self.menuwalous.Append(wx.ID_ANY, _("Legend"), _("Legend"))
            self._walous_filepath = None
            self._walous_layer = None
            self._walous_map = None

            self.menubar.Append(self.menuwalous, _('Walous'))

            self.menuwalous.Bind(wx.EVT_MENU, self.Onmenuwalous)


    def get_canvas_bounds(self, gridsize:float = None):
        """
        Retourne les limites de la zone d'affichage

        :return: [xmin, ymin, xmax, ymax]

        """

        if gridsize is None:

            return [self.xmin, self.ymin, self.xmax, self.ymax]

        else:

            xmin = float(np.rint(self.xmin / gridsize) * gridsize)
            ymin = float(np.rint(self.ymin / gridsize) * gridsize)
            xmax = float(np.rint(self.xmax / gridsize) * gridsize)
            ymax = float(np.rint(self.ymax / gridsize) * gridsize)

            return [xmin, ymin, xmax, ymax]

    def Onmenuwalous(self, event: wx.MenuEvent):

        id = event.GetId()
        item = self.menubar.FindItemById(event.GetId())

        if item is None:
            return

        itemlabel = item.ItemLabel

        if itemlabel in [_("Crop on active array"), _("Crop on screen")]:


            if itemlabel == _("Crop on screen"):

                bounds = self.get_canvas_bounds(gridsize=1.)

                def_outdrir = ''
                spatial_res = 1.

                if self.active_array is not None:
                    spatial_res = self.active_array.dx

                dlg = wx.TextEntryDialog(None,_("Spatial resolution [m] ?"), value = str(spatial_res))

                dlg.ShowModal()
                try:
                    spatial_res = float(dlg.GetValue())
                    dlg.Destroy()
                except:
                    dlg.Destroy()
                    logging.warning(_("Bad value -- Rety"))
                    return

            else:

                if self.active_array is None:
                    logging.warning(_('No active array -- Please activate data first'))
                    return

                bounds = self.active_array.get_bounds()
                def_outdrir = Path(self.active_array.filename).parent
                spatial_res = self.active_array.dx

            from .pywalous import Walous_data, WALOUS2MANNING_MAJ_NIV1, WALOUS2MANNING_MAJ_NIV2, update_palette_walous

            if self._walous_filepath is None:
                dlg = wx.FileDialog(self, _("Choose the Walous shape file"), wildcard="Shapefile (*.shp)|*.shp|all (*.*)|*.*", style=wx.FD_OPEN)
                if dlg.ShowModal() == wx.ID_CANCEL:
                    dlg.Destroy()
                    return

                self._walous_filepath = Path(dlg.GetPath())
                dlg.Destroy()


            dlg = wx.FileDialog(self, _("Choose the output file"), wildcard="Geotif (*.tif)|*.tif|all (*.*)|*.*", style=wx.FD_SAVE, defaultDir=str(def_outdrir))
            if dlg.ShowModal() == wx.ID_CANCEL:
                dlg.Destroy()
                return

            output = Path(dlg.GetPath())
            dlg.Destroy()

            # choix de la couche entre MAJ_NIV1 et MAJ_NIV2
            dlg = wx.SingleChoiceDialog(None, _("Choose a layer"), "Choices", ['MAJ_NIV1', 'MAJ_NIV2'])
            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            self._walous_layer = dlg.GetStringSelection()

            locwalous = Walous_data(self._walous_filepath.parent, self._walous_filepath.name)
            ret = locwalous.rasterize(bounds=bounds,
                                      layer=self._walous_layer,
                                      fn_out=output,
                                      pixel_size=spatial_res)

            if isinstance(ret, int):
                logging.error(_('Error {}').format(ret))
                return

            if Path(output).exists():
                logging.info(_('File {} created').format(output))
            else:
                logging.error(_('File {} not created').format(output))
                return

            dlg = wx.MessageDialog(self, _('Do you want to load the created file ?'), _('Load file'), wx.YES_NO | wx.ICON_QUESTION)
            ret = dlg.ShowModal()

            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            elif ret == wx.ID_YES:
                walousarray = WolfArray(fname=output)
                update_palette_walous(self._walous_layer, walousarray.mypal)
                self.add_object('array', newobj=walousarray, id = 'walous_crop')
                dlg.Destroy()

        elif itemlabel == _("Legend"):

            from .pywalous import WalousLegend

            newlegend = WalousLegend(self)
            newlegend.Show()

        elif itemlabel == _("Map active array"):

            from .pywalous import DlgMapWalous, WALOUS2MANNING_MAJ_NIV1, WALOUS2MANNING_MAJ_NIV2

            if self.active_array is None:
                logging.warning(_('No active array -- Please activate data first'))
                return

            vals = self.active_array.get_unique_values()

            if self._walous_layer is None:

                if vals[0] > 10:
                    self._walous_layer = 'MAJ_NIV2'
                else:
                    self._walous_layer = 'MAJ_NIV1'

            dlg = DlgMapWalous(self, which=self._walous_layer)

            ret = dlg.ShowModal()

            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            mapvals = dlg.get_mapping()
            dlg.Destroy()

            if mapvals == -1:
                logging.error(_('Bad values -- retry'))
                return

            self.active_array.map_values(mapvals)

            self.active_array.reset_plot()

    def Onmenuwolf2d(self, event: wx.MenuEvent):

        id = event.GetId()
        item = self.menubar.FindItemById(event.GetId())

        if item is None:
            return

        itemlabel = item.ItemLabel


        if itemlabel == _("Read last result"):

            self.read_last_result()

        elif itemlabel == _("Change current view"):

            # Change view for results

            autoscale = False
            choices = [cur.value for cur in views_2D]
            dlg = wx.SingleChoiceDialog(None, _("Pick a view"), "Choices", choices)
            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            method = dlg.GetStringSelection()

            method = list(views_2D)[choices.index(method)]

            dlg.Destroy()

            diamsize = None
            if method == views_2D.SHIELDS_NUMBER :

                if self.active_res2d is not None:
                    sediment_diam = self.active_res2d.sediment_diameter
                    sediment_density = self.active_res2d.sediment_density
                elif self.compare_results is not None:
                    sediment_diam = 0.001
                    sediment_density = 2.650
                else:
                    logging.warning(_('No active 2D result or comparison !'))
                    return

                dlg = wx.TextEntryDialog(None,_("Diameter grain size [m] ?"), value = str(sediment_diam))
                ret = dlg.ShowModal()
                if ret == wx.ID_CANCEL:
                    dlg.Destroy()
                    return
                try:
                    diamsize = float(dlg.GetValue())
                except:
                    dlg.Destroy()
                    logging.warning(_("Bad value -- Rety"))
                    return

                dlg = wx.TextEntryDialog(None,_("Density grain [-] ?"), value = str(sediment_density))
                ret = dlg.ShowModal()
                if ret == wx.ID_CANCEL:
                    dlg.Destroy()
                    return
                try:
                    density = float(dlg.GetValue())
                except:
                    dlg.Destroy()
                    logging.warning(_("Bad value -- Rety"))
                    return

            if len(self.myres2D)>1:

                dlg = wx.MessageDialog(None, _('Apply to all results?'), style=wx.YES_NO)
                ret = dlg.ShowModal()
                if ret == wx.ID_NO:
                    if diamsize is not None:
                        self.active_res2d.sediment_diameter = diamsize
                        self.active_res2d.sediment_density = density
                        self.active_res2d.load_default_colormap('shields_cst')

                    self.active_res2d.set_currentview(method, force_wx = True, force_updatepal = True)
                else:
                    for curarray in self.iterator_over_objects(draw_type.RES2D):
                        curarray:Wolfresults_2D
                        if diamsize is not None:
                            curarray.sediment_diameter = diamsize
                            curarray.sediment_density  = density
                            curarray.load_default_colormap('shields_cst')

                        curarray.set_currentview(method, force_wx = True, force_updatepal = True)

            else:
                if self.active_res2d is not None:
                    if diamsize is not None:
                        self.active_res2d.sediment_diameter = diamsize
                        self.active_res2d.sediment_density = density
                        self.active_res2d.load_default_colormap('shields_cst')
                    self.active_res2d.set_currentview(method, force_wx = True, force_updatepal = True)

            if self.compare_results is not None:
                # update compare results
                if diamsize is not None:
                    self.compare_results.set_shields_param(diamsize, density)
                self.compare_results.update_type_result(method)

        elif itemlabel == _("Set epsilon water depth"):

            dlg = wx.TextEntryDialog(self, _('Enter an epsilon [m]'),value='0.0')

            ret = dlg.ShowModal()

            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            try:
                neweps = float(dlg.GetValue())
                dlg.Destroy()
            except:
                logging.error(_('Bad value -- retry !'))
                dlg.Destroy()
                return

            for curmodel in self.iterator_over_objects(draw_type.RES2D):
                curmodel: Wolfresults_2D
                curmodel.epsilon = neweps
                curmodel._epsilon_default = neweps
                curmodel.read_oneresult(curmodel.current_result)
                curmodel.set_currentview()

        elif itemlabel == _("Filter independent"):

            self.menu_filter_independent.IsChecked = not self.menu_filter_independent.IsChecked

            for curmodel in self.iterator_over_objects(draw_type.RES2D):
                curmodel: Wolfresults_2D
                curmodel.to_filter_independent = not self.menu_filter_independent.IsChecked

        # elif itemlabel == _("Manage boundary conditions..."):

        #     if self.active_res2d is not None:
        #         self.active_res2d.myparams.editing_bc(self.myres2D)

        elif itemlabel ==_("Create video..."):
            if self.active_res2d is not None:
                self.create_video()

        elif itemlabel == _("Danger map - only h"):
            if self.active_res2d is None:
                logging.warning(_('No active 2D result !'))
                return

            with wx.NumberEntryDialog(None, _('Danger map'), _('From step'), _('Danger map'), 1, 1, self.active_res2d.get_nbresults()) as dlg:

                if dlg.ShowModal() == wx.ID_CANCEL:
                    return

                start_step = dlg.GetValue()

            with wx.NumberEntryDialog(None, _('Danger map'), _('To step'), _('Danger map'), self.active_res2d.get_nbresults(), start_step, self.active_res2d.get_nbresults()) as dlg:

                if dlg.ShowModal() == wx.ID_CANCEL:
                    return

                end_step = dlg.GetValue()

            with wx.NumberEntryDialog(None, _('Danger map'), _('Every'), _('Danger map'), 1, 1, 60) as dlg:

                if dlg.ShowModal() == wx.ID_CANCEL:
                    return

                every = dlg.GetValue()

            danger_map = self.active_res2d.danger_map_only_h(start_step-1, end_step-1, every)

            with wx.DirDialog(None, _('Choose a directory'), style=wx.DD_DEFAULT_STYLE) as dlg:

                if dlg.ShowModal() == wx.ID_CANCEL:
                    return

                outdir = dlg.GetPath()

            danger_map.write_all(Path(outdir) / 'danger_h.tif')

        elif itemlabel == _("Danger map"):
            if self.active_res2d is None:
                logging.warning(_('No active 2D result !'))
                return

            with wx.NumberEntryDialog(None, _('Danger map'), _('From step'), _('Danger map'), 1, 1, self.active_res2d.get_nbresults()) as dlg:

                if dlg.ShowModal() == wx.ID_CANCEL:
                    return

                start_step = dlg.GetValue()

            with wx.NumberEntryDialog(None, _('Danger map'), _('To step'), _('Danger map'), self.active_res2d.get_nbresults(), start_step, self.active_res2d.get_nbresults()) as dlg:

                if dlg.ShowModal() == wx.ID_CANCEL:
                    return

                end_step = dlg.GetValue()

            with wx.NumberEntryDialog(None, _('Danger map'), _('Every'), _('Danger map'), 1, 1, 60) as dlg:

                if dlg.ShowModal() == wx.ID_CANCEL:
                    return

                every = dlg.GetValue()

            danger_maps = self.active_res2d.danger_map(start_step-1, end_step-1, every)

            with wx.DirDialog(None, _('Choose a directory'), style=wx.DD_DEFAULT_STYLE) as dlg:

                if dlg.ShowModal() == wx.ID_CANCEL:
                    return

                outdir = dlg.GetPath()

            names = ['danger_h.tif', 'danger_u.tif', 'danger_q.tif', 'danger_Z.tif']
            for name, danger_map in zip(names, danger_maps):
                danger_map.write_all(Path(outdir) / name)

        elif itemlabel == _("Setup cache..."):

            if self.active_res2d is None:
                logging.warning(_('No active 2D result !'))
                return

            dlg = wx.MessageDialog(None, _('Cache only water depth results ?'), style=wx.YES_NO)
            ret = dlg.ShowModal()
            if ret == wx.ID_NO:
                only_h = False
            else:
                only_h = True
            dlg.Destroy()

            dlg = wx.MessageDialog(None, _('Cache all results ?'), style=wx.YES_NO)
            ret = dlg.ShowModal()
            if ret == wx.ID_NO:

                dlg_start = wx.SingleChoiceDialog(None, _('Choosing the start index'),
                                                  _('Choices'),
                                                  [str(cur) for cur in range(1,self.active_res2d.get_nbresults()+1)])
                ret = dlg_start.ShowModal()
                if ret == wx.ID_CANCEL:
                    dlg_start.Destroy()
                    return

                start_idx = int(dlg_start.GetStringSelection())
                dlg_start.Destroy()

                dlg_end   = wx.SingleChoiceDialog(None, _('Choosing the end index'),
                                                  _('Choices'),
                                                  [str(cur) for cur in range(start_idx + 1,self.active_res2d.get_nbresults()+1)])

                ret = dlg_end.ShowModal()
                if ret == wx.ID_CANCEL:
                    dlg_end.Destroy()
                    return

                dlg_end.Destroy()

                end_idx = int(dlg_end.GetStringSelection())

                logging.info(_('Caching from {} to {} - Be patient !').format(start_idx, end_idx))
                self.active_res2d.setup_cache(start_idx = start_idx-1, end_idx = end_idx-1, only_h=only_h)
                logging.info(_('Caching done !'))
            else:
                logging.info(_('Caching all results - Be patient !'))
                self.active_res2d.setup_cache(only_h=only_h)
                logging.info(_('Caching done !'))

            dlg.Destroy()

        elif itemlabel == _("Clear cache..."):

            if self.active_res2d is None:
                logging.warning(_('No active 2D result !'))
                return

            self.active_res2d.clear_cache()
            logging.info(_('Cache cleared !'))

        elif itemlabel == _("Show tiles..."):

            if self.active_res2d is None:
                logging.warning(_('No active 2D result !'))
                return

            self.active_res2d.show_tiles()


    def menu_2dgpu(self):

        if self.menuwolf2d is not None:
            if self.menu2d_cache_setup  is None:
                self.menu2d_cache_setup = self.menuwolf2d.Append(wx.ID_ANY, _("Setup cache..."), _("Set up cache for 2D GPU model"))
                self.menu2d_cache_reset = self.menuwolf2d.Append(wx.ID_ANY, _("Clear cache..."), _("Clear cache for 2D GPU model"))
                self.menu2d_show_tiles = self.menuwolf2d.Append(wx.ID_ANY, _("Show tiles..."), _("Show a grid of tiles for 2D GPU model"))

    def menu_landmaps(self):

        if self.menu_landmap is None:
            self.menu_landmap = wx.Menu()
            self.menubar.Append(self.menu_landmap, _('&Landmap'))

            self.menupick_landmap_full = self.menu_landmap.Append(wx.ID_ANY, _("Pick landmap full..."), _("Pick landmap full resolution"))
            self.menupick_landmap_low = self.menu_landmap.Append(wx.ID_ANY, _("Pick landmap low..."), _("Pick landmap low resolution"))
            self.menu_landmap.AppendSeparator()

            self.menu_colortransparent_landmap = self.menu_landmap.Append(wx.ID_ANY, _("Transparent color "), _("Change transparent color associated to the landmap"))
            self.menu_tolerance_landmap = self.menu_landmap.Append(wx.ID_ANY, _("Set tolerance"), _("Set tolerance for the transparent color landmap"))
            self.menu_color_landmap = self.menu_landmap.Append(wx.ID_ANY, _("Change colors"), _("Change color map associated to the landmap"))

            self.menu_landmap.Bind(wx.EVT_MENU, self.pick_landmap_full, self.menupick_landmap_full)
            self.menu_landmap.Bind(wx.EVT_MENU, self.pick_landmap_low, self.menupick_landmap_low)
            self.menu_landmap.Bind(wx.EVT_MENU, self.change_colors_landmap, self.menu_color_landmap)
            self.menu_landmap.Bind(wx.EVT_MENU, self.change_transparent_color_landmap, self.menu_colortransparent_landmap)
            self.menu_landmap.Bind(wx.EVT_MENU, self.set_tolerance_landmap, self.menu_tolerance_landmap)

    def change_transparent_color_landmap(self, event: wx.Event):

        if self.active_landmap is None:
            logging.warning(_('No active landmap -- Please load data first'))

        data = wx.ColourData()
        data.SetColour(self.active_landmap.transparent_color)
        with wx.ColourDialog(self, data) as dlg:
            if dlg.ShowModal() == wx.ID_OK:
                data = dlg.GetColourData()
                color = data.GetColour()
                self.active_landmap.set_transparent_color([color.Red(), color.Green(), color.Blue()])

    def set_tolerance_landmap(self, event: wx.Event):

        if self.active_landmap is None:
            logging.warning(_('No active landmap -- Please load data first'))

        dlg = wx.TextEntryDialog(self, _('Set the tolerance for the transparent color'), _('Tolerance'), str(self.active_landmap.tolerance))
        ret = dlg.ShowModal()
        if ret == wx.ID_CANCEL:
            dlg.Destroy()
            return
        tol = int(dlg.GetValue())
        tol = max(0, tol)
        dlg.Destroy()

        self.active_landmap.set_tolerance(tol)


    def change_colors_landmap(self, event: wx.Event):

        if self.active_landmap is None:
            logging.warning(_('No active landmap -- Please load data first'))

        data = wx.ColourData()
        data.SetColour(self.active_landmap.color)
        with wx.ColourDialog(self, data) as dlg:
            if dlg.ShowModal() == wx.ID_OK:
                data = dlg.GetColourData()
                color = data.GetColour()
                self.active_landmap.set_color(color)

    def pick_landmap_full(self, event: wx.Event):

        self.action = 'pick landmap full'
        logging.info(_('Pick landmap - Full resolution'))

    def pick_landmap_low(self, event: wx.Event):

        self.action = 'pick landmap low'
        logging.info(_('Pick landmap - Low resolution'))

    def menu_particlesystem(self):
        if self.menuparticlesystem is None:

            self.menuparticlesystem = wx.Menu()
            self.menuparticlesystem_load = wx.Menu()

            self.menuparticlesystem.Append(wx.ID_ANY, _("Set..."), _("Set arrays as the domain/uv of the particle system -- Must be a 2D array - Mask will be used to separate water and land"))
            self.menuparticlesystem.Append(wx.ID_ANY, _("Set emitter from selected nodes"), _("Set the selected nodes as emitters of the particle system"))
            self.menuparticlesystem.AppendSubMenu(self.menuparticlesystem_load, _("Load..."),  _('Load data for the particle system in the UI'))

            self.menuparticlesystem_load.Append(wx.ID_ANY, _("Load domain..."), _("Loading the domain in the UI"))
            self.menuparticlesystem_load.Append(wx.ID_ANY, _("Load ~domain..."), _("Loading the negative of the domain in the UI"))
            self.menuparticlesystem_load.Append(wx.ID_ANY, _("Load emitters..."), _("Loading the emitters in the UI"))
            self.menuparticlesystem_load.Append(wx.ID_ANY, _("Load uv..."), _("Loading the UV velocity field in the UI"))
            self.menuparticlesystem_load.Append(wx.ID_ANY, _("Load uv norm..."), _("Loading the norm of the velocity field in the UI"))
            # self.menuparticlesystem.Append(wx.ID_ANY, _("Set emitters..."), _("Set active zones as the emitters of the particle system -- Each checked zone will be used as an emitter"))
            # self.menuparticlesystem.Append(wx.ID_ANY, _("Set emitter..."), _("Set only the active vector as an emitters of the particle system"))
            # self.menuparticlesystem.Append(wx.ID_ANY, _("Set uv..."), _("Choose U and V arrays for the particle system -- Must be 2D arrays"))
            self.menuparticlesystem.AppendSeparator()
            self.menuparticlesystem.Append(wx.ID_ANY, _("Check"), _("Check if the particle system is ready to be computed"))
            self.menuparticlesystem.Append(wx.ID_ANY, _("Bake"), _("Compute the particle system"))
            self.menuparticlesystem.Append(wx.ID_ANY, _("Reset"), _("Clear all results but keep the particle system settings"))
            self.menuparticlesystem.AppendSeparator()
            # self.menuparticlesystem.AppendSeparator()
            self.menuparticlesystem.Append(wx.ID_ANY, _("Start"), _("Run all steps"))
            self.menuparticlesystem.Append(wx.ID_ANY, _("Stop"), _("Stop the current animation"))
            self.menuparticlesystem.Append(wx.ID_ANY, _("Resume"), _("Resume animation"))

            self.timer_ps = wx.Timer(self)

            self.menuparticlesystem.Bind(wx.EVT_MENU, self.action_menu_particlesystem)
            self.Bind(wx.EVT_TIMER, self.update_particlesystem, self.timer_ps)

            self.menubar.Append(self.menuparticlesystem, _('Particle system'))

    def action_menu_particlesystem(self, event: wx.Event):
        """ Action to perform when the timer is triggered """

        if self.active_particle_system is not None:

            itemlabel = self.menuparticlesystem.FindItemById(event.GetId()).GetItemLabelText()


            if itemlabel == _("Start"):

                if self.active_particle_system is not None:
                    self.active_particle_system.current_step = 0
                    self.active_particle_system.current_step_idx = 0
                    self.timer_ps.Start(1000. / self.active_particle_system.fps)

            elif itemlabel == _("Stop"):

                self.timer_ps.Stop()

            elif itemlabel == _("Resume"):

                self.timer_ps.Start(1000. / self.active_particle_system.fps)

            elif itemlabel == _("Load domain..."):
                domain = self.active_particle_system.get_domain(output_type='wolf')
                self.add_object('array', id=domain.idx, newobj=domain, ToCheck=True)
                self.Refresh()

            elif itemlabel == _("Load ~domain..."):
                domain:WolfArray = self.active_particle_system.get_domain(output_type='wolf')
                domain.idx = domain.idx + '_neg'
                domain.mask_reset()

                ones = np.where(domain.array.data == 1)
                domain.array[:,:] = 1
                domain.array[ones] = 0

                domain.mask_data(domain.nullvalue)
                self.add_object('array', id=domain.idx, newobj=domain, ToCheck=True)
                self.Refresh()

            elif itemlabel == _("Load emitters..."):
                emitters = self.active_particle_system.get_emitters(output_type='wolf')
                self.add_object('vector', id=emitters.idx, newobj=emitters, ToCheck=True)
                self.Refresh()

            elif itemlabel == _("Load uv..."):
                u = self.active_particle_system.get_u(output_type='wolf')
                v = self.active_particle_system.get_v(output_type='wolf')
                self.add_object('array', id=u.idx, newobj=u, ToCheck=True)
                self.add_object('array', id=v.idx, newobj=v, ToCheck=True)
                self.Refresh()

            elif itemlabel == _("Load uv norm..."):
                uvnorm = self.active_particle_system.get_uv_absolute(output_type='wolf')
                self.add_object('array', id=uvnorm.idx, newobj=uvnorm, ToCheck=True)
                self.Refresh()

            elif itemlabel == _("Bake"):
                check, msg = self.active_particle_system.bake()

                if not check:
                    dlg = wx.MessageDialog(self, msg, _('Error'), wx.OK | wx.ICON_ERROR)
                    dlg.ShowModal()
                    dlg.Destroy()
                    return
            elif itemlabel == _("Reset"):
                self.active_particle_system.reset()

            elif itemlabel == _("Set..."):
                from .ui.wolf_multiselection_collapsiblepane import Wolf_MultipleSelection

                setter = Wolf_MultipleSelection(self,
                                                title=_("Set particle system"),
                                                message=_("Choose arrays/emitters for the particle system"),
                                                values_dict={'domain': self.get_list_keys(draw_type.ARRAYS),
                                                             'u': self.get_list_keys(draw_type.ARRAYS),
                                                             'v': self.get_list_keys(draw_type.ARRAYS),
                                                             'emitters': self.get_list_keys(draw_type.VECTORS)},
                                                info='Set : \n - domain (1 value)\n - u and v (multiple values)\n - emitters (multiple values)',
                                                styles=[wx.LB_SINGLE, wx.LB_EXTENDED, wx.LB_EXTENDED, wx.LB_EXTENDED],
                                                max_choices=[1, None, None, None],
                                                delete_if_transfer = [True, False, False, True],
                                                destroyOK=False)
                setter.ShowModal()
                ret_dict = setter.get_values()
                setter.Destroy()

                if 'domain' in ret_dict:
                    if len(ret_dict['domain']) == 1:
                        domain = self.getobj_from_id(ret_dict['domain'][0])
                        self.active_particle_system.set_domain(domain)
                if 'u' in ret_dict and 'v' in ret_dict:
                    if len(ret_dict['u']) >0:
                        assert len(ret_dict['u']) == len(ret_dict['v']), _('Please select the same number of u and v arrays')

                        time = 0.
                        for u,v in zip(ret_dict['u'], ret_dict['v']):
                            u = self.getobj_from_id(u)
                            v = self.getobj_from_id(v)
                            u:WolfArray
                            v:WolfArray
                            assert u.array.shape == v.array.shape, _('Please select arrays with the same shape')
                            assert u.origx == v.origx and u.origy == v.origy, _('Please select arrays with the same origin')
                            assert u.dx == v.dx and u.dy == v.dy, _('Please select arrays with the same resolution')
                            self.active_particle_system.set_uv((u, v),
                                                            (u.origx, u.origy, u.dx, u.dy),
                                                            time = time)
                            time += 1.

                if 'emitters' in ret_dict:
                    if len(ret_dict['emitters'])>0:
                        emitters = [self.getobj_from_id(cur) for cur in ret_dict['emitters']]
                        self.active_particle_system.set_emitters(emitters)

                if self.active_particle_system._ui is not None:
                    self.active_particle_system.show_properties()

            elif itemlabel == _("Set emitter from selected nodes"):
                if self.active_array is None:
                    logging.warning(_('No active array -- Please activate an array first'))
                    return
                if len(self.active_array.SelectionData.myselection) == 0 and len(self.active_array.SelectionData.selections) ==0:
                    logging.warning(_('No selection -- Please select some nodes first'))
                    return

                from .lagrangian.emitter import Emitter

                newemitters=[]
                if len(self.active_array.SelectionData.myselection) > 0:
                    indices = [self.active_array.get_ij_from_xy(cur[0], cur[1]) for cur in self.active_array.SelectionData.myselection]
                    newemitters = [Emitter(indices,
                                         header = (self.active_array.origx, self.active_array.origy, self.active_array.dx, self.active_array.dy))]

                if len(self.active_array.SelectionData.selections) > 0:

                    for cursel in self.active_array.SelectionData.selections.values():
                        indices = [self.active_array.get_ij_from_xy(cur[0], cur[1]) for cur in cursel['select']]
                        newemitters += [Emitter(indices, header = (self.active_array.origx, self.active_array.origy, self.active_array.dx, self.active_array.dy))]

                self.active_particle_system.set_emitters(newemitters)

                if self.active_particle_system._ui is not None:
                    self.active_particle_system.show_properties()

            # elif itemlabel == _("Set emitters..."):

            #     if self.active_zones is None:
            #         logging.warning(_('No active zones -- Please activate zones first'))
            #         return

            #     self.active_particle_system.set_emitters(self.active_zones)

            # elif itemlabel == _("Set emitter..."):

            #     if self.active_vector is None:
            #         logging.warning(_('No active vector -- Please activate a vector first'))
            #         return

            #     self.active_particle_system.set_emitter(self.active_vector)

            # elif itemlabel == _("Set uv..."):

            #     list_arrays = self.multiple_choice_object(draw_type.ARRAYS, message=_('Choose U and V arrays for the particle system -- first == u ; second == v'), titel='UV choice' )

            #     if len(list_arrays) != 2:
            #         logging.error(_('Please select two arrays and ONLY two arrays'))
            #         return

            #     self.active_particle_system.set_uv(tuple(list_arrays))

            elif itemlabel == _("Check"):
                check, msg = self.active_particle_system.check()

                if not check:
                    dlg = wx.MessageDialog(self, msg, _('Error'), wx.OK | wx.ICON_ERROR)
                    dlg.ShowModal()
                    dlg.Destroy()
                    return
                else:
                    dlg = wx.MessageDialog(self, _('All is fine -- You can bake you system !'), _('Chesk particle system'), wx.OK | wx.ICON_INFORMATION)
                    dlg.ShowModal()
                    dlg.Destroy()
                    return

    def update_particlesystem(self, event: wx.Event):
        """ Animation of the particle system """

        if self.active_particle_system is not None:

            nb = self.active_particle_system.nb_steps
            self.active_particle_system.current_step_idx += 1
            self.Paint()
            self._update_mytooltip()

            if self.active_particle_system.current_step_idx == nb-1:
                self.timer_ps.Stop()

    def menu_sim2D(self):
        """ Menu for 2D simulations """

        if self.menusim2D is None:
            self.menusim2D = wx.Menu()
            self.menubar.Append(self.menusim2D, _('Tools 2D'))

            menu2d_options = self.menusim2D.Append(wx.ID_ANY, _("Parameters..."), _("Parameters"))
            menu2d_zbin2hbin = self.menusim2D.Append(wx.ID_ANY, _("Convert zbin to hbin"), _("Convert zbin to hbin"))
            menu2d_hbin2zbin = self.menusim2D.Append(wx.ID_ANY, _("Convert hbin to zbin"), _("Convert hbin to zbin"))
            menu2D_zbinb2hbinb = self.menusim2D.Append(wx.ID_ANY, _("Convert zbinb to hbinb"), _("Convert zbinb to hbinb"))
            menu2d_hbinb2zbinb = self.menusim2D.Append(wx.ID_ANY, _("Convert hbinb to zbinb"), _("Convert hbinb to zbinb"))
            menu2d_forcemask = self.menusim2D.Append(wx.ID_ANY, _("Reset mask of all arrays"), _("Reset mask"))

            # update = self.menusim2D.Append(wx.ID_ANY, _('Update model from current mask'), _('Update model'))
            # updateblocfile = self.menusim2D.Append(wx.ID_ANY, _('Update .bloc file'), _('Update bloc'))
            # updatefreesurface = self.menusim2D.Append(wx.ID_ANY, _('Update free surface elevation - IC'), _('Update free surface elevation'))
            # updaterough = self.menusim2D.Append(wx.ID_ANY, _('Update roughness coeff'), _('Update roughness coefficient'))
            # updateic = self.menusim2D.Append(wx.ID_ANY, _('Update IC reading mode'), _('Update IC'))
            # menu2d_tft_ic = self.menusim2D.Append(wx.ID_ANY,_("Transfer initial conditions..."),_("Transfer IC"))

            self.menusim2D.Bind(wx.EVT_MENU, self.Onmenusim2D)

    def Onmenusim2D(self, event: wx.MenuEvent):
        """ Action to perform whern menu 2D entry is selected """

        id = event.GetId()
        item = self.menubar.FindItemById(event.GetId())

        if item is None:
            return

        itemlabel = item.ItemLabel

        from .PyGui import Wolf2DModel

        if not isinstance(self.wolfparent, Wolf2DModel):
            logging.error(_('This is not a 2D model'))
            return

        self.wolfparent:Wolf2DModel

        if itemlabel == _('Update .bloc file'):

            msg = _('If you continue the .bloc file will be relpaced !')+'\n'
            msg += '\n'
            msg += _('Continue ?')+'\n'

            dlg = wx.MessageDialog(self,msg,caption = _('Attention'), style = wx.YES_NO)
            ret = dlg.ShowModal()
            dlg.Destroy()
            if ret == wx.ID_NO:
                return

            self.wolfparent.write_bloc_file()

        elif itemlabel == _('Reset mask of all arrays'):

            self.wolfparent.sim.force_mask()

        elif itemlabel == _('Convert zbin to hbin'):

            if self.wolfparent.sim._zbin is not None:
                self.wolfparent.sim.zbin2hbin()
                self.wolfparent.sim.hbin.reset_plot()

        elif itemlabel == _('Convert hbin to zbin'):

            if self.wolfparent.sim._hbin is not None:
                self.wolfparent.sim.hbin2zbin()
                self.wolfparent.sim.zbin.reset_plot()

        elif itemlabel == _('Convert zbinb to hbinb'):

            if self.wolfparent.sim._zbinb is not None:
                self.wolfparent.sim.zbinb2hbinb()
                self.wolfparent.sim.hbinb.reset_plot()

        elif itemlabel == _('Convert hbinb to zbinb'):

            if self.wolfparent.sim._hbinb is not None:
                self.wolfparent.sim.hbinb2zbinb()
                self.wolfparent.sim.zbinb.reset_plot()

        elif itemlabel == _("Transfer initial conditions..."):

            if self.active_array is not None:
                from .PyGui import Wolf2DModel
                if isinstance(self.wolfparent,Wolf2DModel):
                    self.wolfparent.transfer_ic(self.active_vector)

        elif itemlabel == _("Parameters..."):
            self.wolfparent.show_properties()

        elif itemlabel == _('Update free surface elevation - IC'):

            if len(self.active_array.SelectionData.myselection)==0:

                msg = _('There is none selected nodes in the active array !')+'\n'
                msg += '\n'
                msg += _('Please select the desired zone and retry !')+'\n'

                logging.warning(msg)
                return

            self.wolfparent.extend_freesurface_elevation(self.active_array.SelectionData.myselection)

        elif itemlabel== _('Update roughness coeff'):

            if len(self.active_array.SelectionData.myselection)==0:

                msg = _('There is none selected nodes in the active array !')+'\n'
                msg += '\n'
                msg += _('Please select the desired zone and retry !')+'\n'

                logging.warning(msg)
                return

            self.wolfparent.extend_roughness(self.active_array.SelectionData.myselection)

        # elif itemlabel == _('Update IC reading mode'):

        #     self.wolfparent.set_type_ic()

        elif itemlabel == _('Update model from current mask'):

            if type(self.active_array) not in [WolfArray]:
                msg = _('Please select a mono-block array !')+'\n'
                dlg=wx.MessageBox(msg,style=wx.OK)
                return

            msg = _('If you continue, the mask of all arrays will be replaced by the current mask !')+'\n'
            msg += _('The external contour in the .bloc file will also be relpaced.')+'\n'
            msg += '\n'
            msg += _('Continue ?')+'\n'

            dlg = wx.MessageDialog(self,msg,caption = _('Attention'), style = wx.YES_NO)
            ret = dlg.ShowModal()
            dlg.Destroy()
            if ret == wx.ID_NO:
                return

            with wx.lib.busy.BusyInfo(_('Updating 2D model')):
                wait = wx.BusyCursor()

                sux,suy,cont,interior = self.active_array.suxsuy_contour(self.wolfparent.filenamegen,True)

                self.wolfparent.mimic_mask(self.active_array)
                self.wolfparent.replace_external_contour(cont,interior)

                del wait

            self.wolfparent.extend_bed_elevation()

    def get_configuration(self) -> Union[WolfConfiguration, None]:
        """ Get global configuration parameters """

        # At this point, I'm not too sure about
        # which window/frame does what. So to be on
        # the safe side, I make sure that the configuration
        # menu is active only on the "first" window.
        # Moreover, I try to go up the frame/window
        # hierarchy to get the configuration (which will therefore
        # be treated as a singleton)
        if self.wolfparent:
            return self.wolfparent.get_configuration()
        else:
            return None

    @property
    def bkg_color(self):
        """ Return the background color from configs """
        config = self.get_configuration()
        if config is None:
            return [255.,255.,255.,255.]
        else:
            return config[ConfigurationKeys.COLOR_BACKGROUND]

    @property
    def ticks_size(self) -> float:
        """ Return the ticks spacing from configs """

        config = self.get_configuration()
        if config is None:
            return 100.
        else:
            return config[ConfigurationKeys.TICKS_SIZE]

    @property
    def ticks_bounds(self) -> bool:
        """ Return the ticks bounds from configs """

        config = self.get_configuration()
        if config is None:
            return True
        else:
            return config[ConfigurationKeys.TICKS_BOUNDS]


    def GlobalOptionsDialog(self, event):
        handle_configuration_dialog(self, self.get_configuration())

    # def import_3dfaces(self):

    #     dlg = wx.FileDialog(None, _('Choose filename'),
    #                         wildcard='dxf (*.dxf)|*.dxf|gltf (*.gltf)|*.gltf|gltf binary (*.glb)|*.glb|All (*.*)|*.*', style=wx.FD_OPEN)
    #     ret = dlg.ShowModal()
    #     if ret == wx.ID_CANCEL:
    #         dlg.Destroy()
    #         return

    #     fn = dlg.GetPath()
    #     dlg.Destroy()

    #     mytri = Triangulation(plotted=True,mapviewer=self)

    #     if fn.endswith('.dxf'):
    #         mytri.import_dxf(fn)
    #     elif fn.endswith('.gltf') or fn.endswith('.glb'):
    #         mytri.import_from_gltf(fn)

    #     self.add_object('triangulation',newobj=mytri,id=fn)
    #     self.active_tri = mytri

    def triangulate_cs(self):

        msg = ''
        if self.active_zones is None:
            msg += _(' The active zones is None. Please activate the desired object !\n')
        if self.active_cs is None:
            msg += _(' The is no cross section. Please active the desired object or load file!')

        if msg != '':
            dlg = wx.MessageBox(msg, 'Required action')
            return

        dlg = wx.NumberEntryDialog(None, _('What is the desired size [cm] ?'), 'ds', 'ds size', 100, 1, 10000)
        ret = dlg.ShowModal()
        if ret == wx.ID_CANCEL:
            dlg.Destroy()
            return

        ds = float(dlg.GetValue()) / 100.
        dlg.Destroy()

        self.myinterp = Interpolators(self.active_zones, self.active_cs, ds)

        self.add_object('vector', newobj=self.myinterp.myzones, ToCheck=False, id='Interp_mesh')

        if self.menuviewerinterpcs is None:
            self.menuviewerinterpcs = self.cs_menu.Append(wx.ID_ANY, _("New cloud Viewer..."),
                                                            _("Cloud viewer Interpolate"))
        if self.menuinterpcs is None:
            self.menuinterpcs = self.cs_menu.Append(wx.ID_ANY, _("Interpolate on active array..."), _("Interpolate"))

    def interpolate_cloud(self):
        """
        Interpolation d'un nuage de point sur une matrice

        Il est possible d'utiliser une autre valeur que la coordonnées Z des vertices
        """
        if self.active_cloud is not None and self.active_array is not None:

            keyvalue='z'
            if self.active_cloud.header:
                choices = list(self.active_cloud.myvertices[0].keys())
                dlg = wx.SingleChoiceDialog(None, "Pick the value to interpolate", "Choices", choices)
                ret = dlg.ShowModal()
                if ret == wx.ID_CANCEL:
                    dlg.Destroy()
                    return

                keyvalue = dlg.GetStringSelection()
                dlg.Destroy()

            choices = ["nearest", "linear", "cubic"]
            dlg = wx.SingleChoiceDialog(None, "Pick an interpolate method", "Choices", choices)
            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            method = dlg.GetStringSelection()
            dlg.Destroy()

            self.active_cloud.interp_on_array(self.active_array,keyvalue,method)

    def interpolate_cs(self):
        if self.active_array is not None and self.myinterp is not None:

            choices = ["nearest", "linear", "cubic"]
            dlg = wx.SingleChoiceDialog(None, "Pick an interpolate method", "Choices", choices)
            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            method = dlg.GetStringSelection()
            dlg.Destroy()

            self.myinterp.interp_on_array(self.active_array, method)

    def interpolate_triangulation(self):

        if self.active_array is not None and self.active_tri is not None:

            self.active_array.interpolate_on_triangulation(self.active_tri.pts, self.active_tri.tri, )

    def compare_cloud2array(self):
        """
        Compare the active cloud points to the active array

        """

        if self.active_array is None :
            logging.warning(_('No active array -- Please activate an array first'))
            return

        if self.active_cloud is None:
            logging.warning(_('No active cloud -- Please activate a cloud first'))
            return

        self.active_array.compare_cloud(self.active_cloud)

    def compare_tri2array(self):

        if self.active_array is not None and self.active_tri is not None:

            self.active_array.compare_tri(self.active_tri)

    def copy_canvasogl(self, mpl=True, ds=0., figsizes=[10.,10.], palette:wolfpalette = None):
        """

        Generate image based on UI context and copy to the Clipboard

        Args:
            mpl (bool, optional): Using Matplolib as renderer. Defaults to True.
            ds (float, optional): Ticks size. Defaults to 0..
            figsizes : fig size in inches
        """
        if wx.TheClipboard.Open():
            self.Paint()

            if self.SetCurrentContext():

                # Récupération  du buffer OpenGL
                glPixelStorei(GL_PACK_ALIGNMENT, 1)
                data = glReadPixels(0, 0, self.canvaswidth, self.canvasheight, GL_RGBA, GL_UNSIGNED_BYTE)
                # Création d'une image sur base du buffer
                myimage: Image.Image
                myimage = Image.frombuffer("RGBA", (self.canvaswidth, self.canvasheight), data)
                # On tranpose car OpenGL travaille avec comme référence le coin inférieur gauche
                myimage = myimage.transpose(1)

                metadata = PngInfo()
                metadata.add_text('xmin', str(self.xmin))
                metadata.add_text('ymin', str(self.ymin))
                metadata.add_text('xmax', str(self.xmax))
                metadata.add_text('ymax', str(self.ymax))

                if mpl:
                    if ds == 0.:
                        ds = self.ticks_size
                        # dlg = wx.NumberEntryDialog(self,
                        #                         _("xmin : {:.3f} \nxmax : {:.3f} \nymin : {:.3f} \nymax : {:.3f} \n\n  dx : {:.3f}\n  dy : {:.3f}").format(
                        #                             self.xmin, self.xmax, self.ymin, self.ymax, self.xmax - self.xmin,
                        #                             self.ymax - self.ymin),
                        #                         _("Interval [m]"), _("Ticks interval ?"), 500, 1, 10000)
                        # ret = dlg.ShowModal()

                        # if ret == wx.ID_CANCEL:
                        #     dlg.Destroy()
                        #     return

                        # ds = float(dlg.GetValue())
                        # dlg.Destroy()

                    # Création d'une graphique Matplotlib
                    extent = (self.xmin, self.xmax, self.ymin, self.ymax)
                    fig, ax = plt.subplots(1, 1)

                    w, h = [self.width, self.height]

                    # neww = figsizes[0]/figsizes[1] *h
                    neww = figsizes[0]
                    newh = h/w * figsizes[0]

                    fig.set_size_inches(neww, newh)

                    pos = ax.imshow(myimage,
                                    origin='upper',
                                    extent=extent)
                    # fig.colorbar(pos,ax=ax)

                    x1 = np.ceil((self.xmin // ds) * ds)
                    if x1 < self.xmin:
                        x1 += ds
                    x2 = int((self.xmax // ds) * ds)
                    if x2 > self.xmax:
                        x2 -= ds
                    y1 = np.ceil((self.ymin // ds) * ds)
                    if y1 < self.ymin:
                        y1 += ds
                    y2 = int((self.ymax // ds) * ds)
                    if y2 > self.ymax:
                        y2 -= ds

                    x_label_list = np.linspace(x1, x2, int((x2 - x1) / ds) + 1, True)
                    if self.ticks_bounds:
                        x_label_list = np.insert(x_label_list, 0, self.xmin)
                        x_label_list = np.insert(x_label_list, -1, self.xmax)
                        x_label_list = np.unique(x_label_list)

                    y_label_list = np.linspace(y1, y2, int((y2 - y1) / ds) + 1, True)
                    if self.ticks_bounds:
                        y_label_list = np.insert(y_label_list, 0, self.ymin)
                        y_label_list = np.insert(y_label_list, -1, self.ymax)
                        y_label_list = np.unique(y_label_list)

                    ax.set_xticks(x_label_list)
                    ax.set_yticks(y_label_list)

                    ax.set_xticklabels(plt.FormatStrFormatter('%.1f').format_ticks(x_label_list), fontsize=14, rotation=30)
                    ax.set_yticklabels(plt.FormatStrFormatter('%.1f').format_ticks(y_label_list), fontsize=14)

                    ax.set_xlabel('X [m]')
                    ax.set_ylabel('Y [m]')

                    #création de 2 buffers
                    buf = io.BytesIO()
                    #sauvegarde de la figure au format png
                    fig.tight_layout()
                    fig.savefig(buf, format='png')

                    #déplacement au début du buffer
                    buf.seek(0)
                    #lecture du buffer et conversion en image avec PIL
                    im = Image.open(buf)

                    if palette is None:
                        if self.active_array is not None:
                            palette = self.active_array.mypal
                        elif self.active_res2d is not None:
                            palette = self.active_res2d.mypal

                    if palette is not None:
                        bufpal = io.BytesIO()
                        palette.export_image(bufpal,'v')
                        bufpal.seek(0)

                        #lecture du buffer et conversion en image avec PIL
                        impal = Image.open(bufpal)
                        impal = impal.resize((int(impal.size[0]*im.size[1]*.8/impal.size[1]),int(im.size[1]*.8)))

                        imnew = Image.new('RGB',(im.size[0]+impal.size[0], im.size[1]), (255,255,255))

                        # On colle l'image du buffer et la palette pour ne former qu'une seul image à copier dans le clipboard
                        imnew.paste(im.convert('RGB'),(0,0))
                        imnew.paste(impal.convert('RGB'),(im.size[0]-10, int((im.size[1]-impal.size[1])/3)))
                        im=imnew

                    #création d'un objet bitmap wx
                    wxbitmap = wx.Bitmap().FromBuffer(im.width,im.height,im.tobytes())

                    # objet wx exportable via le clipboard
                    dataobj = wx.BitmapDataObject()
                    dataobj.SetBitmap(wxbitmap)

                    wx.TheClipboard.SetData(dataobj)
                    wx.TheClipboard.Close()

                    buf.close()

                    return fig, ax, im

                else:
                    """ Création d'un objet bitmap wx sur base du canvas
                    et copie dans le clipboard
                    """
                    # wxbitmap = wx.Bitmap().FromBuffer(myimage.width,myimage.height,myimage.tobytes())
                    wxbitmap = wx.Bitmap().FromBufferRGBA(myimage.width,myimage.height,myimage.tobytes())

                    # objet wx exportable via le clipboard
                    dataobj = wx.BitmapDataObject()
                    dataobj.SetBitmap(wxbitmap)

                    wx.TheClipboard.SetData(dataobj)
                    wx.TheClipboard.Close()

                    return myimage

        else:
            wx.MessageBox("Can't open the clipboard", "Error")

    def display_canvasogl(self,mpl =True, ds=0., fig: Figure = None, ax: Axes = None, clear = True, redraw =True, palette=False, title=''):
        """
        This method takes a matplotlib figure and axe and,
        returns a clear screenshot of the information displayed in the wolfpy GUI.
        """
        self.Paint()
        myax = ax
        if redraw:
            if clear:
                myax.clear()


        if self.SetCurrentContext():
            glPixelStorei(GL_PACK_ALIGNMENT, 1)
            data = glReadPixels(0,0,self.canvaswidth, self.canvasheight, GL_RGBA,GL_UNSIGNED_BYTE)
            myimage: Image.Image
            myimage =  Image.frombuffer("RGBA",(self.canvaswidth,self.canvasheight),data)
            myimage = myimage.transpose(1)

            if mpl:
                if ds ==0.:
                    ds = self.ticks_size

                extent = (self.xmin, self.xmax, self.ymin, self.ymax)

                myax.imshow(myimage, origin ='upper', extent=extent)

                x1 = np.ceil((self.xmin//ds)*ds)
                if x1 < self.xmin:
                    x1 += ds
                x2 = int((self.xmax//ds)*ds)
                if x2 >self.xmax:
                    x2 -= ds
                y1 = np.ceil((self.ymin//ds)*ds)
                if y1 < self.ymin:
                    y1 += ds
                y2 = int((self.ymax // ds) * ds)
                if y2 > self.ymax:
                    y2 -= ds

                x_label_list = np.linspace(x1,x2, int((x2-x1)/ds) +1, True)
                if self.ticks_bounds:
                    x_label_list = np.insert(x_label_list,0,self.xmin)
                    x_label_list = np.insert(x_label_list,-1, self.xmax)
                    x_label_list = np.unique(x_label_list)

                y_label_list = np.linspace(y1, y2, int((y2 - y1) / ds) + 1, True)
                if self.ticks_bounds:
                    y_label_list = np.insert(y_label_list, 0, self.ymin)
                    y_label_list = np.insert(y_label_list, -1, self.ymax)
                    y_label_list = np.unique(y_label_list)

                myax.set_xticks(x_label_list)
                myax.set_yticks(y_label_list)

                myax.set_xticklabels(FormatStrFormatter('%.1f').format_ticks(x_label_list), fontsize = 'xx-small', rotation = 90)
                myax.set_yticklabels(FormatStrFormatter('%.1f').format_ticks(y_label_list), fontsize='xx-small')
                myax.xaxis.set_ticks_position('top')
                myax.xaxis.set_label_position('top')

                myax.set_xlabel('X - coordinates ($m$)')
                myax.set_ylabel('Y - coordinates ($m$)')
                myax.xaxis.set_ticks_position('bottom')
                myax.xaxis.set_label_position('bottom')

                if title!='':
                    myax.set_title(title)

                fig.canvas.draw()
                fig.canvas.flush_events()

        else:
            logging.warning( "Can't open the clipboard", "Error")

    def get_mpl_plot(self, center = [0., 0.], width = 500., height = 500., title='', toshow=True) -> tuple[Figure, Axes]:
        """
        Récupère un graphique matplotlib sur base de la fenêtre OpenGL et de la palette de la matrice active
        """
        self.zoom_on(center=center, width=width, height= height, canvas_height=self.canvasheight, forceupdate=True)


        fig,axes = plt.subplots(1,2, gridspec_kw={'width_ratios': [20, 1]})
        self.display_canvasogl(fig=fig,ax=axes[0])

        if self.active_array is not None:
            self.active_array.mypal.export_image(None, h_or_v='v', figax=(fig,axes[1]))
        elif self.active_res2d is not None:
            self.active_res2d.mypal.export_image(None, h_or_v='v', figax=(fig,axes[1]))

        axes[0].xaxis.set_ticks_position('bottom')
        axes[0].xaxis.set_label_position('bottom')

        fig.set_size_inches(12,10)

        fontsize(axes[0], 12)
        fontsize(axes[1], 12)

        if title!='':
            axes[0].set_title(title)

        fig.tight_layout()
        if toshow:
            fig.show()

        return fig, axes

    def create_video(self, fn:str = '', framerate:int = 0, start_step:int = 0, end_step:int = 0, every:int = 0):
        """
        Création d'une vidéo sur base des résultats
        """
        try:
            import cv2
        except:
            logging.error(_('Please install opencv-python'))
            return

        if fn=='':
            dlg = wx.FileDialog(parent   = None,
                                message  = _('Choose file name'),
                                wildcard = 'AVI video file (*.avi)|*.avi',
                                style    = wx.FD_SAVE)
            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return
            fn = dlg.GetPath()
            dlg.Destroy()

        if not fn.endswith('.avi'):
            fn+='.avi'

        if framerate<1:
            dlg = wx.NumberEntryDialog(None, _("Frame rate [nb_images/second]"), _('Frame rate'), _('Frame rate'), 24, 1, 100)

            ret = dlg.ShowModal()

            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            framerate = int(dlg.GetValue())
            dlg.Destroy()


        times,steps = self.active_res2d.get_times_steps()

        el_time = str(timedelta(seconds=int(times[self.active_res2d.current_result])))

        fig:Figure
        fig, ax = self.get_mpl_plot([self.mousex, self.mousey],
                                    self.width,
                                    self.height,
                                    title=_('Current time {:0>8} s'.format(el_time)),
                                    toshow=False)

        video = cv2.VideoWriter(fn, cv2.VideoWriter_fourcc(*'XVID'), framerate, fig.canvas.get_width_height())

        nb = self.active_res2d.get_nbresults()
        cid = max(self.active_res2d.current_result,1)

        self.active_res2d.mypal.automatic=False

        if start_step==0:
            dlg = wx.NumberEntryDialog(None, _("First step"), _('From'), _('from'), cid, 1, nb)

            ret = dlg.ShowModal()

            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            start_step = int(dlg.GetValue())
            dlg.Destroy()

        if end_step==0:
            dlg = wx.NumberEntryDialog(None, _("Final step"), _('To'), _('To'), nb, 1, nb)

            ret = dlg.ShowModal()

            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            end_step = int(dlg.GetValue())
            dlg.Destroy()

        if every==0:
            dlg = wx.NumberEntryDialog(None, _("Interval"), _('Interval'), _('Interval'), 1, 1, end_step-start_step)

            ret = dlg.ShowModal()

            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            interval = int(dlg.GetValue())
            dlg.Destroy()

        self.read_one_result(start_step)

        for curmodel in self.iterator_over_objects(draw_type.RES2D):
            curmodel: Wolfresults_2D
            curmodel.step_interval_results = interval

        for idx in tqdm(range(start_step, end_step, interval)):

            image = Image.frombytes('RGB', fig.canvas.get_width_height(),fig.canvas.tostring_rgb())

            video.write(cv2.cvtColor(np.asarray(image),cv2.COLOR_RGB2BGR))

            self.simul_next_step()

            el_time = str(timedelta(seconds=int(times[self.active_res2d.current_result])))

            self.display_canvasogl(fig=fig,
                                   ax=ax[0],
                                   title=_('Current time {:0>8} s'.format(el_time)))

        for curmodel in self.iterator_over_objects(draw_type.RES2D):
            curmodel: Wolfresults_2D
            curmodel.step_interval_results = 1


    def get_canvas_as_image(self) -> Image.Image:
        """
        Récupère la fenêtre OpenGL sous forme d'image
        """

        self.Paint()

        if self.SetCurrentContext():
            glPixelStorei(GL_PACK_ALIGNMENT, 1)
            data = glReadPixels(0, 0, self.canvaswidth, self.canvasheight, GL_RGBA, GL_UNSIGNED_BYTE)
            myimage: Image.Image
            myimage = Image.frombuffer("RGBA", (self.canvaswidth, self.canvasheight), data)
            myimage = myimage.transpose(1)

            return myimage

    def save_canvasogl(self, fn:str='', mpl:bool=True, ds:float=0.):
        """
        Sauvegarde de la fenêtre d'affichage dans un fichier

        fn : File name (.png file)
        """

        if not fn.endswith('.png'):
            fn += '.png'

        self.Paint()

        if fn == '':
            dlg = wx.FileDialog(None, _('Choose file name'), wildcard='PNG (*.png)|*.png|JPG (*.jpg)|*.jpg',
                                style=wx.FD_SAVE)
            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return
            fn = dlg.GetPath()
            dlg.Destroy()

        if self.SetCurrentContext():
            glPixelStorei(GL_PACK_ALIGNMENT, 1)
            data = glReadPixels(0, 0, self.canvaswidth, self.canvasheight, GL_RGBA, GL_UNSIGNED_BYTE)
            myimage: Image.Image
            myimage = Image.frombuffer("RGBA", (self.canvaswidth, self.canvasheight), data)
            myimage = myimage.transpose(1)

            metadata = PngInfo()
            metadata.add_text('xmin', str(self.xmin))
            metadata.add_text('ymin', str(self.ymin))
            metadata.add_text('xmax', str(self.xmax))
            metadata.add_text('ymax', str(self.ymax))

            if mpl:
                if ds == 0.:
                    dlg = wx.NumberEntryDialog(self,
                                               _("xmin : {:.3f} \nxmax : {:.3f} \nymin : {:.3f} \nymax : {:.3f} \n\n  dx : {:.3f}\n  dy : {:.3f}").format(
                                                   self.xmin, self.xmax, self.ymin, self.ymax, self.xmax - self.xmin,
                                                   self.ymax - self.ymin),
                                               _("Interval [m]"), _("Ticks interval ?"), 500, 1, 10000)
                    ret = dlg.ShowModal()

                    if ret == wx.ID_CANCEL:
                        dlg.Destroy()
                        return

                    ds = float(dlg.GetValue())
                    dlg.Destroy()

                extent = (self.xmin, self.xmax, self.ymin, self.ymax)
                fig, ax = plt.subplots(1, 1)
                pos = ax.imshow(myimage, origin='upper',
                          extent=extent)
                # fig.colorbar(pos,ax=ax)

                x1 = np.ceil((self.xmin // ds) * ds)
                if x1 < self.xmin:
                    x1 += ds
                x2 = int((self.xmax // ds) * ds)
                if x2 > self.xmax:
                    x2 -= ds
                y1 = np.ceil((self.ymin // ds) * ds)
                if y1 < self.ymin:
                    y1 += ds
                y2 = int((self.ymax // ds) * ds)
                if y2 > self.ymax:
                    y2 -= ds

                x_label_list = np.linspace(x1, x2, int((x2 - x1) / ds) + 1, True)
                x_label_list = np.insert(x_label_list, 0, self.xmin)
                x_label_list = np.insert(x_label_list, -1, self.xmax)
                x_label_list = np.unique(x_label_list)

                y_label_list = np.linspace(y1, y2, int((y2 - y1) / ds) + 1, True)
                y_label_list = np.insert(y_label_list, 0, self.ymin)
                y_label_list = np.insert(y_label_list, -1, self.ymax)
                y_label_list = np.unique(y_label_list)

                ax.set_xticks(x_label_list)
                ax.set_yticks(y_label_list)

                ax.set_xticklabels(plt.FormatStrFormatter('%.1f').format_ticks(x_label_list), fontsize=8, rotation=30)
                ax.set_yticklabels(plt.FormatStrFormatter('%.1f').format_ticks(y_label_list), fontsize=8)

                ax.set_xlabel('X [m]')
                ax.set_ylabel('Y [m]')

                plt.savefig(fn, dpi=300)
            else:
                myimage.save(fn, pnginfo=metadata)

            return fn, ds
        else:
            raise NameError(
                'Opengl setcurrent -- maybe a conflict with an existing opengl32.dll file - please rename the opengl32.dll in the libs directory and retry')

    def reporting(self, dir=''):
        if dir == '':
            dlg = wx.DirDialog(None, "Choose directory to store reporting", style=wx.FD_SAVE)
            ret = dlg.ShowModal()

            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            dir = dlg.GetPath()
            dlg.Destroy()

        myppt = Presentation(__file__)
        slide = myppt.slides.add_slide(0)

        for curzone in self.myzones:
            for curvec in curzone.myvectors:
                curvec: vector
                if curvec.nbvertices > 1:
                    oldwidth = curvec.myprop.width
                    curvec.myprop.width = 4
                    myname = curvec.myname

                    self.Activate_vector(curvec)

                    if self.linked:
                        for curview in self.linkedList:
                            title = curview.GetTitle()
                            curview.zoomon_activevector()
                            fn = path.join(dir, title + '_' + myname + '.png')
                            curview.save_canvasogl(fn)
                    else:
                        self.zoomon_activevector()
                        fn = path.join(dir, myname + '.png')
                        self.save_canvasogl(fn)

                        fn = path.join(dir, 'palette_v_' + myname + '.png')
                        self.active_array.mypal.export_image(fn, 'v')
                        fn = path.join(dir, 'palette_h_' + myname + '.png')
                        self.active_array.mypal.export_image(fn, 'h')

                    curvec.myprop.width = oldwidth

    def InitUI(self):

        self.Bind(wx.EVT_SIZE, self.OnSize)
        self.Bind(wx.EVT_CLOSE, self.OnClose)

        # self.canvas.Bind(wx.EVT_CONTEXT_MENU, self.OnShowPopup)
        self.canvas.Bind(wx.EVT_PAINT, self.OnPaint)

        self.treelist.Bind(wx.EVT_CHAR_HOOK, self.OnHotKey)
        self.treelist.Bind(dataview.EVT_TREELIST_SELECTION_CHANGED,self.onselectitem)
        self.canvas.Bind(wx.EVT_CHAR_HOOK, self.OnHotKey)

        # self.treelist.Bind(wx.EVT_KEY_DOWN, self.OnKeyDown)
        # self.canvas.Bind(wx.EVT_KEY_DOWN, self.OnKeyDown)

        self.canvas.Bind(wx.EVT_BUTTON, self.OnButton)
        self.canvas.Bind(wx.EVT_RIGHT_DCLICK, self.OnRDClick)
        self.canvas.Bind(wx.EVT_LEFT_DCLICK, self.OnLDClick)
        self.canvas.Bind(wx.EVT_LEFT_DOWN, self.OnLDown)
        self.canvas.Bind(wx.EVT_MIDDLE_DOWN, self.OnLDown)
        self.canvas.Bind(wx.EVT_RIGHT_DOWN, self.OnRightDown)
        self.canvas.Bind(wx.EVT_RIGHT_UP, self.OnRightUp)
        self.canvas.Bind(wx.EVT_MOTION, self.OnMotion)
        self.canvas.Bind(wx.EVT_LEAVE_WINDOW, self.OnLeave)
        self.canvas.Bind(wx.EVT_MOUSEWHEEL, self.OnButton)

        self.treelist.Bind(dataview.EVT_TREELIST_ITEM_CHECKED, self.OnCheckItem)
        self.treelist.Bind(dataview.EVT_TREELIST_ITEM_ACTIVATED, self.OnActivateTreeElem)
        self.treelist.Bind(dataview.EVT_TREELIST_ITEM_CONTEXT_MENU, self.OntreeRight)
        # dispo dans wxpython 4.1 self.Bind(wx.EVT_GESTURE_ZOOM,self.OnZoomGesture)

        self.Centre()

        self.mybc = []
        self.myarrays = []
        self.mypartsystems = []
        self.myvectors = []
        self.mytiles = []
        self.myclouds = []
        self.mytri = []
        self.myothers = []
        self.myviews = []
        self.mywmsback = []
        self.mywmsfore = []
        self.myres2D = []
        self.myviewers3d = []

        # liste des éléments modifiable dans l'arbre
        self.all_lists = [self.myarrays, self.myvectors, self.myclouds, self.mytri, self.myothers, self.myviews, self.myres2D, self.mytiles, self.mypartsystems, self.myviewers3d]

        if self.get_configuration() is not None:
            # see PyGui.py if necessary

            self.menu_options = wx.Menu()
            self.menubar.Append(self.menu_options, _('Options'))
            self.option_global = self.menu_options.Append(wx.ID_ANY,_("Global"),_("Modify global options"))
            self.Bind(wx.EVT_MENU, self.GlobalOptionsDialog, self.option_global)

        self.Show(True)

    def OnSize(self, e):
        """
        Redimensionnement de la fenêtre
        """
        if self.regular:
            # retrouve la taille de la fenêtre
            width, height = self.GetClientSize()
            # enlève la barre d'arbre
            width -= self.treewidth
            # définit la taille de la fenêtre graphique OpenGL et sa position (à droite de l'arbre)
            self.canvas.SetSize(width, height)
            self.canvas.SetPosition((self.treewidth, 0))
            # calcule les limites visibles sur base de la taille de la fenêtre et des coefficients sx sy
            self.setbounds()
            # fixe la taille de l'arbre (notamment la hauteur)
            # self.treelist.SetSize(self.treewidth,height)
            e.Skip()

    def ManageActions(self, id):
        """
        Gestion des actions via les menus

        TODO : A généraliser?
        """
        curmenu = self.tools[id]['menu']

        if curmenu.IsCheckable():
            if not curmenu.IsChecked():
                curmenu.Check(False)
                self.action = None

                if id == ID_LOCMINMAX:
                    self.update_absolute_minmax = True
            else:
                curmenu.Check()
                if not self.tools[id]['name'] is None:
                    self.action = self.tools[id]['name']

        else:
            if id == ID_SORTALONG:
                # Tri le long d'un vecteur
                if not self.active_cs is None and not self.active_vector is None:
                    self.active_cs: crosssections
                    self.active_vector: vector
                    self.active_cs.sort_along(self.active_vector.asshapely_ls(), self.active_vector.myname, False)
                else:
                    msg = ''
                    if self.active_cs is None:
                        msg += _('Please select the active cross sections \n')
                    if self.active_vector is None:
                        msg += _('Please select the active supprt vector')
                    mydiag = wx.MessageDialog(self, msg, _('Sort along'))
                    mydiag.ShowModal()

    def center_view_on(self, cx, cy):
        """
        Center the view on the point of (map) coordinates (x,y)
        """

        self.mousex, self.mousey = cx, cy

        # retrouve la taille de la fenêtre OpenGL
        width, height = self.canvas.GetSize()

        # calcule la taille selon X et Y en coordonnées réelles
        width = width / self.sx
        height = height / self.sy

        # retrouve les bornes min et max sur base de la valeur centrale qui est censée ne pas bouger
        self.xmin = self.mousex - width / 2.
        self.xmax = self.xmin + width
        self.ymin = self.mousey - height / 2.
        self.ymax = self.ymin + height

    def setbounds(self,updatescale=True):
        """
        Calcule les limites visibles de la fenêtrte graphique sur base des
        facteurs d'échelle courants
        """

        if updatescale:
            self.updatescalefactors()

            # retrouve la taille de la fenêtre OpenGL
            width, height = self.canvas.GetSize()
            self.canvaswidth = width
            self.canvasheight = height

            # calcule la taille selon X et Y en coordonnées réelles
            width = width / self.sx
            height = height / self.sy

            # retrouve les bornes min et max sur base de la valeur centrale qui est censée ne pas bouger
            self.xmin = self.mousex - width / 2.
            self.xmax = self.xmin + width
            self.ymin = self.mousey - height / 2.
            self.ymax = self.ymin + height

            self.width = width
            self.height = height

            self.mousex = self.xmin + width / 2.
            self.mousey = self.ymin + height / 2.

            self.updatescalefactors()

        else:
            # retrouve les bornes min et max sur base de la valeur centrale qui est censée ne pas bouger
            self.xmin = self.mousex - self.width / 2.
            self.xmax = self.xmin + self.width
            self.ymin = self.mousey - self.height / 2.
            self.ymax = self.ymin + self.height

        self.mybackisloaded = False
        self.myfrontisloaded = False

        self.Refresh()
        self.mimicme()

    def setsizecanvas(self,width,height):
        """ Redimensionne la fenêtre graphique """
        self.canvas.SetClientSize(width, height)

    def updatescalefactors(self):
        """ Mise à jour des facteurs d'échelle
            This one updates the scale factors based on the relative sizes
            of the GLCanvas and the footprint that should fit in it.
        """

        width, height = self.canvas.GetSize()

        self.sx = 1
        self.sy = 1
        if self.width > 0 and width >0 :
            self.sx = float(width) / self.width
        if self.height > 0 and height > 0 :
            self.sy = float(height) / self.height

        self.sx = min(self.sx, self.sy)
        self.sy = self.sx

    def add_viewer_and_link(self):
        """ Ajout d'une nouvelle fenêtre de visualisation et liaison avec la fenêtre courante """

        dlg = wx.TextEntryDialog(self, _('Enter a caption for the new window'))

        ret = dlg.ShowModal()

        if ret == wx.ID_CANCEL:
            dlg.Destroy()
            return

        newcap = dlg.GetValue()
        dlg.Destroy()
        newview = WolfMapViewer(None, newcap, w=600, h=600, wxlogging=self.wxlogging, wolfparent=self.wolfparent)
        newview.add_grid()
        newview.add_WMS()

        if self.linkedList is None:
            self.linkedList = [self]

        self.linkedList.append(newview)

        for curview in self.linkedList:
            curview.linked = True
            curview.linkedList = self.linkedList
            curview.link_shareopsvect = False

        logging.info(_('New viewer added and linked'))

    def add_grid(self):
        """ Ajout d'une grille """

        mygrid = Grid(1000.)
        self.add_object('vector', newobj=mygrid, ToCheck=False, id='Grid')

    def add_WMS(self):
        """ Ajout de couches WMS """
        xmin = 0
        xmax = 0
        ymin = 0
        ymax = 0
        orthos = {'IMAGERIE': {'1971': 'ORTHO_1971', '1994-2000': 'ORTHO_1994_2000',
                               '2006-2007': 'ORTHO_2006_2007',
                               '2009-2010': 'ORTHO_2009_2010',
                               '2012-2013': 'ORTHO_2012_2013',
                               '2015': 'ORTHO_2015', '2016': 'ORTHO_2016', '2017': 'ORTHO_2017',
                               '2018': 'ORTHO_2018', '2019': 'ORTHO_2019', '2020': 'ORTHO_2020',
                               '2021': 'ORTHO_2021', '2022 printemps': 'ORTHO_2022_PRINTEMPS', '2022 été': 'ORTHO_2022_ETE',
                               '2023 été': 'ORTHO_2023_ETE',
                               }}
        data_2021 = {'EAU': {'IDW': 'ZONES_INONDEES_IDW',
                             'Emprise': 'ZONES_INONDEES',
                             'Emprise wo Alea': 'ZONES_INONDEES_wo_alea'}}

        for idx, (k, item) in enumerate(orthos.items()):
            for kdx, (m, subitem) in enumerate(item.items()):
                self.add_object(which='wmsback',
                                newobj=imagetexture('PPNC', m, k, subitem,
                                                    self, xmin, xmax, ymin, ymax, -99999, 1024),
                                ToCheck=False, id='PPNC ' + m)

        for idx, (k, item) in enumerate(data_2021.items()):
            for kdx, (m, subitem) in enumerate(item.items()):
                self.add_object(which='wmsback',
                                newobj=imagetexture('PPNC', m, k, subitem,
                                                    self, xmin, xmax, ymin, ymax, -99999, 1024),
                                ToCheck=False, id='Data 2021 ' + m)

        self.add_object(which='wmsback',
                        newobj=imagetexture('PPNC', 'Orthos France', 'OI.OrthoimageCoverage.HR', '',
                                            self, xmin, xmax, ymin, ymax, -99999, 1024, France=True, epsg='EPSG:27563'),
                        ToCheck=False, id='Orthos France')

        forelist = {'EAU': {'Aqualim': 'RES_LIMNI_DGARNE', 'Alea': 'ALEA_INOND', 'Lidaxes': 'LIDAXES'},
                    'LIMITES': {'Secteurs Statistiques': 'LIMITES_QS_STATBEL'},
                    'INSPIRE': {'Limites administratives': 'AU_wms'},
                    'PLAN_REGLEMENT': {'Plan Percellaire': 'CADMAP_2021_PARCELLES'}}

        for idx, (k, item) in enumerate(forelist.items()):
            for kdx, (m, subitem) in enumerate(item.items()):
                self.add_object(which='wmsfore',
                                newobj=imagetexture('PPNC', m, k, subitem,
                                                    self, xmin, xmax, ymin, ymax, -99999, 1024),
                                ToCheck=False, id=m)

    def set_compare(self, ListArrays:list[WolfArray]=None, share_colormap:bool=True):
        """
        Comparison of 2 arrays

        :param ListArrays: List of 2 arrays to compare
        :param share_colormap: Share the colormap between the 2 arrays
        """

        # assert len(ListArrays) == 2, _('List of arrays must contain 2 and only 2 arrays - Here, you have provided {} arrays'.format(len(ListArrays)))

        # Création de 3 fenêtres de visualisation basées sur la classe "WolfMapViewer"
        first = self
        second = WolfMapViewer(None, 'Comparison', w=600, h=600, wxlogging=self.wxlogging, wolfparent = self.wolfparent)
        third  = WolfMapViewer(None, 'Difference', w=600, h=600, wxlogging=self.wxlogging, wolfparent = self.wolfparent)

        second.add_grid()
        third.add_grid()
        second.add_WMS()
        third.add_WMS()

        # Création d'une liste contenant les 3 instances d'objet "WolfMapViewer"
        mylist:list[WolfMapViewer] = []
        mylist.append(first)
        mylist.append(second)
        mylist.append(third)

        # On indique que les objets sont liés en activant le Booléen et en pointant la liste précédente
        for curlist in mylist:
            curlist.linked = True
            curlist.linkedList = mylist

        if ListArrays is not None:
            if len(ListArrays) == 2:
                mnt = ListArrays[0]
                mns = ListArrays[1]

                if not mnt.is_like(mns):
                    logging.warning(_('The 2 arrays must have the same shape - Here, the 2 arrays have different shapes'))
                    return

            else:
                logging.warning(_('List of arrays must contain 2 and only 2 arrays - Here, you have provided {} arrays'.format(len(ListArrays))))
                return
        else:
            logging.warning(_('You must fill the List of arrays with 2 and only 2 arrays - Here, the list is void'))
            return

        mns: WolfArray
        mnt: WolfArray
        diff: WolfArray

        # Recherche d'un masque union des masques partiels
        mns.mask_union(mnt)

        # Création du différentiel -- Les opérateurs mathématiques sont surchargés
        diff = mns - mnt

        # on attribue une matrice par interface graphique
        mnt.change_gui(first)
        mns.change_gui(second)
        diff.change_gui(third)

        path = os.path.dirname(__file__)
        fn = join(path, 'models\\diff16.pal')

        # on partage la palette de couleurs
        if share_colormap:
            mns.add_crosslinked_array(mnt)
            mns.share_palette()

        # on dissocie la palette de la différence
        diff.mypal = wolfpalette()
        if isinstance(diff, WolfArrayMB):
            diff.link_palette()

        diff.mypal.readfile(fn)
        diff.mypal.automatic = False
        diff.myops.palauto.SetValue(0)

        mnt.mypal.automatic = False
        mnt.myops.palauto.SetValue(0)
        if not share_colormap:
            mns.mypal.automatic = False
            mns.myops.palauto.SetValue(0)
            mns.mypal.updatefrompalette(mnt.mypal)

        # Ajout des matrices dans les fenêtres de visualisation
        first.add_object('array', newobj=mnt, ToCheck=True, id='source')
        second.add_object('array', newobj=mns, ToCheck=True, id='comp')
        third.add_object('array', newobj=diff, ToCheck=True, id='diff=comp-source')

        # Partage des vecteurs de la classe d'opérations
        mnt.myops.myzones = mns.myops.myzones
        diff.myops.myzones = mns.myops.myzones

        first.active_array = mnt
        second.active_array = mns
        third.active_array = diff

        mnt.reset_plot()
        mns.reset_plot()
        diff.reset_plot()

    def set_compare_all(self, ListArrays=None, names:list[str] = None):
        """ Comparison of 2 or 3 arrays """

        assert len(ListArrays) == 2 or len(ListArrays) == 3, _('List of arrays must contain 2 or 3 arrays - Here, you have provided {} arrays'.format(len(ListArrays)))
        if names is not None:
            assert len(names) == len(ListArrays)-1, _('List of names must contain the number of names as arrays minus one - Here, you have provided {} names for {} arrays'.format(len(names), len(ListArrays)))
        else:
            names = ['comp1', 'comp2']

        # Création de 3 fenêtres de visualisation basées sur la classe "WolfMapViewer"
        first = self
        second = WolfMapViewer(None, 'Comparison {}'.format(names[0]), w=600, h=600, wxlogging=self.wxlogging, wolfparent=self.wolfparent)
        third  = WolfMapViewer(None, 'Difference {}'.format(names[0]), w=600, h=600, wxlogging=self.wxlogging, wolfparent=self.wolfparent)
        if len(ListArrays) == 3:
            fourth = WolfMapViewer(None, 'Comparison {}'.format(names[1]), w=600, h=600, wxlogging=self.wxlogging, wolfparent=self.wolfparent)
            fifth  = WolfMapViewer(None, 'Difference {}'.format(names[1]), w=600, h=600, wxlogging=self.wxlogging, wolfparent=self.wolfparent)

        # Création d'une liste contenant les multiples instances d'objet "WolfMapViewer"
        list = []
        list.append(first)
        list.append(second)
        list.append(third)
        if len(ListArrays) == 3:
            list.append(fourth)
            list.append(fifth)

        for curview in list:
            if curview is not self:
                curview.add_grid()
                curview.add_grid()

        # On indique que les objets sont liés en actiavt le Booléen et en pointant la liste précédente
        for curview in list:
            curview.linked = True
            curview.linkedList = list

        comp2 = None
        if ListArrays is not None:
            if len(ListArrays) == 2:
                src = ListArrays[0]
                comp1 = ListArrays[1]
            elif len(ListArrays) == 3:
                src = ListArrays[0]
                comp1 = ListArrays[1]
                comp2 = ListArrays[2]
            else:
                return
        else:
            return

        src: WolfArray
        comp1: WolfArray
        diff1: WolfArray
        comp2: WolfArray
        diff2: WolfArray

        # Création du différentiel -- Les opérateurs mathématiques sont surchargés
        diff1 = comp1 - src

        comp1.copy_mask(src, True)
        diff1.copy_mask(src, True)

        src.change_gui(first)
        comp1.change_gui(second)
        diff1.change_gui(third)

        src.mypal.automatic = False
        comp1.mypal.automatic = False
        src.myops.palauto.SetValue(0)
        comp1.myops.palauto.SetValue(0)

        src.mypal.isopop(src.array, src.nbnotnull)
        comp1.mypal.updatefrompalette(src.mypal)

        # Ajout des matrices dans les fenêtres de visualisation
        first.add_object('array', newobj=src, ToCheck=True, id='source')
        second.add_object('array', newobj=comp1, ToCheck=True, id='comp')
        third.add_object('array', newobj=diff1, ToCheck=True, id='diff=comp-source')

        comp1.myops.myzones = src.myops.myzones
        diff1.myops.myzones = src.myops.myzones

        first.active_array = src
        second.active_array = comp1
        third.active_array = diff1

        if comp2 is not None:
            diff2 = comp2 - src
            comp2.copy_mask(src, True)
            diff2.copy_mask(src, True)

            comp2.change_gui(fourth)
            diff2.change_gui(fifth)

            comp2.mypal.automatic = False
            comp2.myops.palauto.SetValue(0)

            comp2.mypal.updatefrompalette(src.mypal)

            # Ajout des matrices dans les fenêtres de visualisation
            fourth.add_object('array', newobj=comp2, ToCheck=True, id='comp2')
            fifth.add_object('array', newobj=diff2, ToCheck=True, id='diff2=comp2-source')

            comp2.myops.myzones = src.myops.myzones
            diff2.myops.myzones = src.myops.myzones

            fourth.active_array = comp2
            fifth.active_array = diff2

    def set_blender_sculpting(self):
        """
        Mise en place de la structure nécessaire pour comparer la donnée de base avec la donnée sculptée sous Blender

        La donnée de base est la matrice contenue dans la fenêtre actuelle

        Fenêtres additionnelles :
            - information sur les volumes de déblai/remblai et bilan
            - matrice sculptée
            - différentiel entre scultage - source
            - gradient
            - laplacien
            - masque de modification
        """
        myframe = wx.Frame(None, title=_('Excavation and backfill'))
        sizergen = wx.BoxSizer(wx.VERTICAL)
        sizer1 = wx.BoxSizer(wx.HORIZONTAL)
        sizer2 = wx.BoxSizer(wx.HORIZONTAL)
        sizer3 = wx.BoxSizer(wx.HORIZONTAL)
        sizergen.Add(sizer1)
        sizergen.Add(sizer2)
        sizergen.Add(sizer3)

        labexc = wx.StaticText(myframe, label=_('Excavation : '))
        labback = wx.StaticText(myframe, label=_('Backfill   : '))
        labbal = wx.StaticText(myframe, label=_('Balance   : '))
        sizer1.Add(labexc)
        sizer2.Add(labback)
        sizer3.Add(labbal)

        font = wx.Font(18, wx.DECORATIVE, wx.NORMAL, wx.NORMAL)

        Exc = wx.StaticText(myframe, label=' [m³]')
        Back = wx.StaticText(myframe, label=' [m³]')
        Bal = wx.StaticText(myframe, label=' [m³]')

        labexc.SetFont(font)
        labback.SetFont(font)
        labbal.SetFont(font)
        Exc.SetFont(font)
        Back.SetFont(font)
        Bal.SetFont(font)

        sizer1.Add(Exc)
        sizer2.Add(Back)
        sizer3.Add(Bal)

        myframe.SetSizer(sizergen)
        myframe.Layout()
        myframe.Centre(wx.BOTH)
        myframe.Show()

        if self.link_params is None:
            self.link_params = {}

        self.link_params['ExcavationBackfill'] = myframe
        self.link_params['Excavation'] = Exc
        self.link_params['Backfill'] = Back
        self.link_params['Balance'] = Bal

        # Création de fenêtres de visualisation basées sur la classe "WolfMapViewer"
        first = self
        second = WolfMapViewer(None, 'Sculpting', w=600, h=600, wxlogging=self.wxlogging, wolfparent=self.wolfparent)
        third  = WolfMapViewer(None, 'Difference', w=600, h=600, wxlogging=self.wxlogging, wolfparent=self.wolfparent)
        fourth = WolfMapViewer(None, 'Gradient', w=600, h=600, wxlogging=self.wxlogging, wolfparent=self.wolfparent)
        fifth  = WolfMapViewer(None, 'Laplace', w=600, h=600, wxlogging=self.wxlogging, wolfparent=self.wolfparent)
        sixth  = WolfMapViewer(None, 'Unitary Mask', w=600, h=600, wxlogging=self.wxlogging, wolfparent=self.wolfparent)

        # Création d'une liste contenant les 3 instances d'objet "WolfMapViewer"
        list = []
        list.append(first)
        list.append(second)
        list.append(third)
        list.append(fourth)
        list.append(fifth)
        list.append(sixth)

        for curlist in list:
            curlist.add_grid()
            curlist.add_WMS()

        # On indique que les objets sont liés en actiavt le Booléen et en pointant la liste précédente
        for curlist in list:
            curlist.linked = True
            curlist.linkedList = list

        source: WolfArray
        sourcenew: WolfArray
        diff: WolfArray
        grad: WolfArray
        lap: WolfArray
        unimask: WolfArray

        source = self.active_array
        sourcenew = WolfArray(mold=source)

        # Création du différentiel -- Les opérateurs mathématiques sont surchargés
        diff = source - source
        grad = source.get_gradient_norm()
        lap = source.get_laplace()
        unimask = WolfArray(mold=diff)

        np.divide(diff.array.data, abs(diff.array.data), out=unimask.array.data, where=diff.array.data != 0.)

        grad.copy_mask(source, True)
        lap.copy_mask(source, True)
        diff.copy_mask(source, True)
        unimask.copy_mask(source, True)

        sourcenew.change_gui(second)
        diff.change_gui(third)
        grad.change_gui(fourth)
        lap.change_gui(fifth)
        unimask.change_gui(sixth)

        path = os.path.dirname(__file__)
        fn=join(path,'models\\diff16.pal')

        if exists(fn):
            diff.mypal.readfile(fn)
            diff.mypal.automatic=False
            diff.myops.palauto.SetValue(0)

        fn=join(path,'models\\diff3.pal')
        if exists(fn):
            unimask.mypal.readfile(fn)
            unimask.mypal.automatic=False
            unimask.myops.palauto.SetValue(0)

        # Ajout des matrices dans les fenêtres de visualisation
        second.add_object('array', newobj=sourcenew, ToCheck=True, id='source_new')
        third.add_object('array', newobj=diff, ToCheck=True, id='diff=comp-source')
        fourth.add_object('array', newobj=grad, ToCheck=True, id='gradient')
        fifth.add_object('array', newobj=lap, ToCheck=True, id='laplace')
        sixth.add_object('array', newobj=unimask, ToCheck=True, id='unimask')

        #pointage des vecteurs attachés à chaque matrice dans chaque GUI de façon à c que les modifications se répercutent  partout
        sourcenew.myops.myzones = source.myops.myzones
        diff.myops.myzones = source.myops.myzones
        grad.myops.myzones = source.myops.myzones
        lap.myops.myzones = source.myops.myzones
        unimask.myops.myzones = source.myops.myzones

        second.active_array = sourcenew
        third.active_array = diff
        fourth.active_array = grad
        fifth.active_array = lap
        sixth.active_array = unimask

        self.mimicme()

    def update_blender_sculpting(self):
        """ Mise à jour des fenêtres de visualisation pour la comparaison avec Blender """
        if not self.linked:
            return
        if len(self.linkedList) != 6:
            return

        # Création de fenêtres de visualisation basées sur la classe "WolfMapViewer"
        first = self.linkedList[0]
        second = self.linkedList[1]
        third = self.linkedList[2]
        fourth = self.linkedList[3]
        fifth = self.linkedList[4]
        sixth = self.linkedList[5]

        source = first.active_array
        sourcenew = second.active_array
        diff = third.active_array
        grad = fourth.active_array
        lap = fifth.active_array
        unimask = sixth.active_array

        fn = ''
        if self.link_params is not None:
            if 'gltf file' in self.link_params.keys():
                fn = self.link_params['gltf file']
                fnpos = self.link_params['gltf pos']

        if fn == '':
            for curgui in self.linkedList:
                if curgui.link_params is not None:
                    if 'gltf file' in curgui.link_params.keys():
                        fn = self.link_params['gltf file']
                        fnpos = self.link_params['gltf pos']
                        break

        with wx.lib.busy.BusyInfo(_('Importing gltf/glb')):
            wait = wx.BusyCursor()
            sourcenew.import_from_gltf(fn, fnpos, 'scipy')
            del wait

        with wx.lib.busy.BusyInfo(_('Update plots')):
            # Création du différentiel -- Les opérateurs mathématiques sont surchargés
            diff.array = (sourcenew - source).array
            grad.array = sourcenew.get_gradient_norm().array
            lap.array = sourcenew.get_laplace().array
            np.divide(diff.array.data, abs(diff.array.data), out=unimask.array.data, where=diff.array.data != 0.)

            diff.copy_mask(sourcenew, True)
            lap.copy_mask(sourcenew, True)
            grad.copy_mask(sourcenew, True)
            unimask.copy_mask(sourcenew, True)

            first.Paint()
            second.Paint()
            third.Paint()
            fourth.Paint()
            fifth.Paint()
            sixth.Paint()

            Exc: wx.StaticText
            Back: wx.StaticText
            Bal: wx.StaticText
            if not 'ExcavationBackfill' in self.link_params.keys():
                for curgui in self.linkedList:
                    if curgui.link_params is not None:
                        if 'ExcavationBackfill' in curgui.link_params.keys():
                            myframe = curgui.link_params['ExcavationBackfill']
                            Exc = curgui.link_params['Excavation']
                            Back = curgui.link_params['Backfill']
                            Bal = curgui.link_params['Balance']
            else:
                myframe = self.link_params['ExcavationBackfill']
                Exc = self.link_params['Excavation']
                Back = self.link_params['Backfill']
                Bal = self.link_params['Balance']

            Exc.SetLabel("{:.2f}".format(np.sum(diff.array[diff.array < 0.])) + ' [m³]')
            Back.SetLabel("{:.2f}".format(np.sum(diff.array[diff.array > 0.])) + ' [m³]')
            Bal.SetLabel("{:.2f}".format(np.sum(diff.array)) + ' [m³]')

    def zoomon_activevector(self, size:float=500., forceupdate:bool=True):
        """
        Zoom on active vector

        :param size: size of the zoomed window
        :param forceupdate: force the update of the window
        """

        if self.active_vector is None:
            logging.warning(_('No active vector'))
            return

        curvec = self.active_vector
        if curvec.xmin == -99999:
            curvec.find_minmax()

        bounds = [curvec.xmin, curvec.xmax, curvec.ymin, curvec.ymax]

        dx = bounds[1] - bounds[0]
        dy = bounds[3] - bounds[2]

        self.mousex = bounds[0] + dx / 2.
        self.mousey = bounds[2] + dy / 2.
        self.width = max(size, dx)
        self.height = max(size, dy)

        self.updatescalefactors()
        self.setbounds()
        self.mimicme()

        if forceupdate:
            self.update()
            if self.linked:
                for cur in self.linkedList:
                    if cur is not self:
                        cur.update()


    def zoomon_active_vertex(self, size:float = 20, forceupdate:bool = True):
        """
        Zoom on active vertex.

        :param size: size of the zoomed window
        :param forceupdate: force the update of the window
        """
        if self.active_vector is None:
            logging.warning(_('No active vector'))
            return
        curvec = self.active_vector
        if curvec.xmin == -99999:
            curvec.find_minmax()

        if self.active_vector is None:
            logging.warning(_('No active vector'))
            return

        grid = self.active_zones.xls
        row = grid.GetGridCursorRow()

        x = float(grid.GetCellValue(row, 0))
        y = float(grid.GetCellValue(row, 1))
        z = float(grid.GetCellValue(row, 2))
        curvert = wolfvertex(x, y, z)
        self.mousex = curvert.x
        self.mousey = curvert.y
        self.width = size
        self.height = size

        self.updatescalefactors()
        self.setbounds()
        self.mimicme()

        if forceupdate:
            self.update()
            if self.linked:
                for cur in self.linkedList:
                    if cur is not self:
                        cur.update()

    def zoom_on_id(self, id:str, drawtype:draw_type = draw_type.ARRAYS, forceupdate=True, canvas_height=1024):
        """
        Zoom on id

        :param id: id of the object to zoom on
        :param drawtype: type of object to zoom on - Different types elements can have the same id

        """

        if drawtype not in [draw_type.ARRAYS, draw_type.VECTORS]:
            logging.warning(_('Draw type must be either ARRAYS or VECTORS'))
            return

        obj = self.get_obj_from_id(id, drawtype)

        if obj is None:
            logging.warning(_('No object found with id {} and drawtype {}'.format(id, drawtype)))
            return

        if drawtype == draw_type.ARRAYS:
            self.zoom_on_array(obj, forceupdate=forceupdate, canvas_height=canvas_height)
        elif drawtype == draw_type.VECTORS:
            self.zoom_on_vector(obj, forceupdate=forceupdate, canvas_height=canvas_height)

    def zoom_on_array(self, array:WolfArray, forceupdate=True, canvas_height=1024):
        """ Zoom on array """

        if array.xmin == -99999:
            array.find_minmax()

        bounds = array.get_bounds()

        center = [(bounds[0][1] + bounds[0][0]) / 2., (bounds[1][1] + bounds[1][0]) / 2.]
        width  = bounds[0][1] - bounds[0][0]
        height = bounds[1][1] - bounds[1][0]

        self.zoom_on({'center':center, 'width':width, 'height':height}, forceupdate=forceupdate, canvas_height=canvas_height)

    def zoom_on_vector(self, vector:vector, forceupdate=True, canvas_height=1024):
        """ Zoom on vector """

        if vector.xmin == -99999:
            vector.find_minmax()

        bounds = vector.get_bounds_xx_yy()

        center = [(bounds[0][1] + bounds[0][0]) / 2., (bounds[1][1] + bounds[1][0]) / 2.]
        width  = bounds[0][1] - bounds[0][0]
        height = bounds[1][1] - bounds[1][0]

        self.zoom_on({'center':center, 'width':width, 'height':height}, forceupdate=forceupdate, canvas_height=canvas_height)

    def create_Zones_from_arrays(self, arrays:list[WolfArray], id:str = None, add_legend:bool=True) -> Zones:
        """
        Create a Zones instance from list of WolfArrays

        One zone per array.

        One vector per zone with the masked contour.

        :param arrays: list of WolfArrays
        :param id: id of the Zones instance
        :param add_legend: add legend to the vector -- centroid of the contour

        """

        # création de l'instance de Zones
        new_zones = Zones(idx = 'contour' if id is None else id.lower(), mapviewer=self)

        for curarray in arrays:

            if isinstance(curarray, WolfArray):

                curarray.nullify_border(1)

                new_zone = zone(name = curarray.idx)
                new_zones.add_zone(new_zone, forceparent=True)

                sux, sux, curvect, interior = curarray.suxsuy_contour()
                new_zone.add_vector(curvect, forceparent=True)

                curvect.set_legend_to_centroid(curarray.idx)
                curvect.myprop.width = 2

                rectvect = vector(name = 'rect_boundary')
                new_zone.add_vector(rectvect, forceparent=True)

                bounds = curarray.get_bounds()

                rectvect.add_vertex(wolfvertex(bounds[0][0], bounds[1][0]))
                rectvect.add_vertex(wolfvertex(bounds[0][1], bounds[1][0]))
                rectvect.add_vertex(wolfvertex(bounds[0][1], bounds[1][1]))
                rectvect.add_vertex(wolfvertex(bounds[0][0], bounds[1][1]))
                rectvect.close_force()

                rectvect.myprop.color = getIfromRGB([255,0,0])
                rectvect.myprop.width = 2

                logging.info(_('{} treated'.format(curarray.idx)))
            else:
                logging.warning(_('All elements in the list must be of type WolfArray'))

        new_zones.find_minmax(update=True)

        return new_zones


    def zoom_on(self, zoom_dict = None, width = 500, height = 500, center = None, xll = None, yll = None, forceupdate=True, canvas_height=1024):
        """
        Zoom on a specific area

        It is possible to zoom on a specific area by providing the zoom parameters in :
          - a dictionnary
          - width and height of the zoomed window and the lower left corner coordinates
          - width and height of the zoomed window and the center coordinates

        :param zoom_dict: dictionnary containing the zoom parameters - possible keys : 'width', 'height', 'center', 'xmin', 'ymin', 'xmax', 'ymax'
        :param width: width of the zoomed window [m]
        :param height: height of the zoomed window [m]
        :param center: center of the zoomed window [m] - tuple (x,y)
        :param xll: lower left X coordinate of the zoomed window [m]
        :param yll: lower left Y coordinate of the zoomed window [m]
        :param forceupdate: force the update of the window
        :param canvas_height: height of the canvas [pixels]


        Examples :

          - zoom_on(zoom_dict = {'center':(500,500), 'width':1000, 'height':1000})
          - zoom_on(width = 1000, height = 1000, xll = 500, yll = 500)
          - zoom_on(width = 1000, height = 1000, center = (500,500))
        """
        if zoom_dict is not None:
            width  = 99999
            height = 99999
            xll  = 99999
            yll = 99999
            xmax  = 99999
            ymax = 99999
            if 'center' in zoom_dict.keys():
                center = zoom_dict['center']
            if 'width' in zoom_dict.keys():
                width = zoom_dict['width']
            if 'height' in zoom_dict.keys():
                height = zoom_dict['height']
            if 'xmin' in zoom_dict.keys():
                xll = zoom_dict['xmin']
            if 'ymin' in zoom_dict.keys():
                yll = zoom_dict['ymin']
            if 'xmax' in zoom_dict.keys():
                xmax = zoom_dict['xmax']
            if 'ymax' in zoom_dict.keys():
                ymax = zoom_dict['ymax']

            if width == 99999:
                width = xmax-xll
            if height == 99999:
                height = ymax-yll

        if center is not None and len(center)==2:
            self.mousex = center[0]
            self.mousey = center[1]
            self.width = width
            self.height = height
        elif (xll is not None) and (yll is not None):
            self.mousex = xll + width/2
            self.mousey = yll + height/2
            self.width = width
            self.height = height

        # fixe la taille de la fenêtre
        v_height = canvas_height
        v_width = int(v_height*(float(width)/float(height)))

        self.SetClientSize(v_width + self.treewidth, v_height)

        self.updatescalefactors()
        self.mimicme()

        if forceupdate:
            self.update()
            if self.linked:
                for cur in self.linkedList:
                    if cur is not self:
                        cur.update()

    def zoom_on_active_profile(self, size:float=500., forceupdate:bool=True):
        """ Zoom on active profile """

        curvec = self.active_profile
        if curvec.xmin == -99999:
            curvec.find_minmax()

        bounds = [curvec.xmin, curvec.xmax, curvec.ymin, curvec.ymax]

        dx = bounds[1] - bounds[0]
        dy = bounds[3] - bounds[2]

        self.mousex = bounds[0] + dx / 2.
        self.mousey = bounds[2] + dy / 2.
        self.width = max(size, dx)
        self.height = max(size, dy)

        self.updatescalefactors()
        self.setbounds()
        self.mimicme()

        if forceupdate:
            self.update()
            if self.linked:
                for cur in self.linkedList:
                    if cur is not self:
                        cur.update()

    def read_project(self, fn):
        """
        Projet WOLF GUI

        Fichier de paramètres contenant les types et chemins d'accès aux données à ajouter

        A compléter...

        """
        myproject = Wolf_Param(None, filename=fn, toShow=False)

        mykeys = ['cross_sections', 'vector', 'array']

        with wx.lib.busy.BusyInfo(_('Opening project')):
            wait = wx.BusyCursor()
            if 'which' in myproject.myparams.keys():
                which = myproject.myparams['which']['action'][key_Param.VALUE]
                if which == 'compare':
                    ListCompare = []
                    if 'array' in myproject.myparams.keys():
                        for curid, curname in zip(myproject.myparams['array'].keys(), myproject.myparams['array'].values()):
                            ListCompare.append(WolfArray(normpath(curname[key_Param.VALUE])))

                    self.set_compare(ListCompare)
                    return

            if 'cross_sections' in myproject.myparams.keys():
                for curid, curname in zip(myproject.myparams['cross_sections'].keys(),
                                        myproject.myparams['cross_sections'].values()):
                    if curid != 'format' and curid != 'dirlaz':
                        mycs = crosssections(curname[key_Param.VALUE],
                                            format=myproject.myparams['cross_sections']['format'][key_Param.VALUE],
                                            dirlaz=myproject.myparams['cross_sections']['dirlaz'][key_Param.VALUE])

                        self.add_object('cross_sections', newobj=mycs, id=curid)


            if myproject.get_group('tiles') is not None:

                curid = myproject.get_param('tiles', 'id')
                curfile = myproject.get_param('tiles', 'tiles_file')
                curdatadir = myproject.get_param('tiles', 'data_dir')
                curcompdir = myproject.get_param('tiles', 'comp_dir')

                if exists(curfile):
                    mytiles = Tiles(filename= curfile, parent=self, linked_data_dir=curdatadir)
                    mytiles.set_comp_dir(curcompdir)
                    self.add_object('tiles', newobj=mytiles, id=curid)
                else:
                    logging.info(_('Bad parameter in project file - tiles : ')+ curfile)

            if myproject.get_group('laz_grid') is not None:

                curdatadir = myproject.get_param('laz_grid', 'data_dir')

                self.init_laz_from_gridinfos(curdatadir)

            if 'vector' in myproject.myparams.keys():
                for curid, curname in zip(myproject.myparams['vector'].keys(), myproject.myparams['vector'].values()):
                    if exists(curname[key_Param.VALUE]):
                        myvec = Zones(curname[key_Param.VALUE], parent=self)
                        self.add_object('vector', newobj=myvec, id=curid)
                    else:
                        logging.info(_('Bad parameter in project file - vector : ')+ curname[key_Param.VALUE])

            if 'array' in myproject.myparams.keys():
                for curid, curname in zip(myproject.myparams['array'].keys(), myproject.myparams['array'].values()):

                    if exists(curname[key_Param.VALUE]):
                        curarray = WolfArray(curname[key_Param.VALUE])
                        self.add_object('array', newobj=curarray, id=curid)
                    else:
                        logging.info(_('Bad parameter in project file - array : ')+ curname[key_Param.VALUE])

            if 'cloud' in myproject.myparams.keys():
                for curid, curname in zip(myproject.myparams['cloud'].keys(), myproject.myparams['cloud'].values()):
                    if exists(curname[key_Param.VALUE]):
                        mycloud = cloud_vertices(curname[key_Param.VALUE])
                        self.add_object('cloud', newobj=mycloud, id=curid)
                    else:
                        logging.info(_('Bad parameter in project file - cloud : ')+ curname[key_Param.VALUE])

            if 'wolf2d' in myproject.myparams.keys():
                for curid, curname in zip(myproject.myparams['wolf2d'].keys(), myproject.myparams['wolf2d'].values()):
                    if exists(curname[key_Param.VALUE]):
                        curwolf = Wolfresults_2D(curname[key_Param.VALUE])
                        self.add_object('res2d', newobj=curwolf, id=curid)
                    else:
                        logging.info(_('Bad parameter in project file - wolf2d : ')+ curname[key_Param.VALUE])

                self.menu_wolf2d()

            if 'gpu2d' in myproject.myparams.keys():
                for curid, curname in zip(myproject.myparams['gpu2d'].keys(), myproject.myparams['gpu2d'].values()):
                    if exists(curname[key_Param.VALUE]):

                        if 'simul_gpu_results' in curname[key_Param.VALUE]:
                            curwolf = wolfres2DGPU(Path(curname[key_Param.VALUE]))
                        else:
                            if exists(join(curname[key_Param.VALUE], 'simul_gpu_results')):
                                curwolf = wolfres2DGPU(Path(join(curname[key_Param.VALUE], 'simul_gpu_results')))
                            else:
                                logging.info(_('Bad directory : ')+ curname[key_Param.VALUE])

                        self.add_object('res2d', newobj=curwolf, id=curid)
                    else:
                        logging.info(_('Bad directory : ')+ curname[key_Param.VALUE])

                self.menu_wolf2d()
                self.menu_2dgpu()

            if 'palette' in myproject.myparams.keys():
                self.project_pal = {}
                for curid, curname in zip(myproject.myparams['palette'].keys(), myproject.myparams['palette'].values()):
                    if exists(curname[key_Param.VALUE]):
                        mypal = wolfpalette(None, '')
                        mypal.readfile(curname[key_Param.VALUE])
                        mypal.automatic = False

                        self.project_pal[curid] = mypal
                    else:
                        logging.info(_('Bad parameter in project file - palette : ')+ curname[key_Param.VALUE])

            if 'palette-array' in myproject.myparams.keys():
                curarray: WolfArray
                if self.project_pal is not None:
                    for curid, curname in zip(myproject.myparams['palette-array'].keys(),
                                            myproject.myparams['palette-array'].values()):
                        if curname[key_Param.VALUE] in self.project_pal.keys():
                            curarray = self.getobj_from_id(curid)
                            if curarray is not None:
                                mypal:wolfpalette
                                mypal = self.project_pal[curname[key_Param.VALUE]]
                                curarray.mypal = mypal
                                if mypal.automatic:
                                    curarray.myops.palauto.SetValue(1)
                                else:
                                    curarray.myops.palauto.SetValue(0)
                                curarray.updatepalette(0)
                                curarray.delete_lists()
                            else:
                                logging.warning(_('Bad parameter in project file - palette-array : ')+ curid)

            if 'cross_sections_link' in myproject.myparams.keys():
                if 'linkzones' in myproject.myparams['cross_sections_link'].keys():
                    idx = myproject.myparams['cross_sections_link']['linkzones'][key_Param.VALUE]

                    for curzones in self.iterator_over_objects(draw_type.VECTORS):
                        curzones: Zones
                        if curzones.idx == idx:
                            self.active_cs.link_external_zones(curzones)

                    zonename = ''
                    vecname = ''

                    if 'sortzone' in myproject.myparams['cross_sections_link'].keys():
                        zonename = myproject.myparams['cross_sections_link']['sortzone'][key_Param.VALUE]
                    if 'sortname' in myproject.myparams['cross_sections_link'].keys():
                        vecname = myproject.myparams['cross_sections_link']['sortname'][key_Param.VALUE]

                    if zonename != '' and vecname != '':
                        names = [cur.myname for cur in curzones.myzones]
                        idx = names.index(zonename)
                        curzone = curzones.myzones[idx]
                        names = [cur.myname for cur in curzone.myvectors]
                        idx = names.index(vecname)
                        curvec = curzone.myvectors[idx]

                        if curvec is not None:
                            curvec: vector
                            self.active_cs.sort_along(curvec.asshapely_ls(), curvec.myname, False)

            if 'vector_array_link' in myproject.myparams.keys():
                for curid, curname in zip(myproject.myparams['vector_array_link'].keys(), myproject.myparams['vector_array_link'].values()):

                    locvec = None
                    locarray = None
                    for curvec in self.myvectors:
                        if curvec.idx == curname[key_Param.VALUE].lower():
                            locvec=curvec
                            break

                    for curarray in self.myarrays:
                        if curarray.idx == curid.lower():
                            locarray=curarray
                            break

                    if locvec is not None and locarray is not None:

                        locarray.linkedvec = locvec.myzones[0].myvectors[0]

                    else:

                        logging.warning(_('Bad vec-array association in project file !'))
                        logging.warning(curid)
                        logging.warning(curname[key_Param.VALUE])
            del wait

    def save_project(self, fn):
        myproject = Wolf_Param(None, toShow=False)

        mykeys = ['cross_sections', 'vector', 'array', 'wolf2d']
        for curkey in mykeys:
            myproject[curkey] = {}
        """
        # myproject.myparams['which']={}
        # myproject.myparams['which']['action']={}
        # myproject.myparams['which']['action'][key_Param.VALUE]

        # mycs = self.active_cs
        # if mycs is not None:
        #     myproject.myparams['cross_sections']={}
        #     myproject.myparams['cross_sections']['mycs']={}
        #     myproject.myparams['cross_sections']['mycs'][key_Param.VALUE]=mycs.filename

        #     myproject.myparams['vector']={}
        #     myproject.myparams['vector']['river']={}
        ##     myproject.myparams['vector']['river'][key_Param.VALUE]=self.added[draw_type.VECTORS.value][0].filename

        # if 'array' in myproject.myparams.key():
        #     for curid,curname in zip(myproject.myparams['array'].keys(),myproject.myparams['array'].values()):
        #         curarray=WolfArray(curname[key_Param.VALUE])
        #         self.add_object('array',newobj=curarray,id=curid)
        """
        # matrices
        try:
            for curel in self.iterator_over_objects(draw_type.ARRAYS):
                myproject['array'][curel.idx] = {}
                myproject['array'][curel.idx][key_Param.VALUE] = curel.filename

        except:
            pass

        # résultats 2D
        try:
            for curel in self.iterator_over_objects(draw_type.RES2D):
                myproject['wolf2d'][curel.idx] = {}
                myproject['wolf2d'][curel.idx][key_Param.VALUE] = curel.filename

        except:
            pass

        # vecteurs
        try:
            for curel in self.iterator_over_objects(draw_type.VECTORS):
                myproject['vector'][curel.idx] = {}
                myproject['vector'][curel.idx][key_Param.VALUE] = curel.filename

        except:
            pass

    def plot_laz_around_active_vec(self):
        if self.active_vector is None:
            return

        if self.mylazgrid is None:
            return

        dlg = wx.NumberEntryDialog(None, _('Enter the size of the window around the active vector [cm]'), _('Window size'),_('Window size'), 500, 0, 2000)
        ret = dlg.ShowModal()
        if ret != wx.ID_OK:
            dlg.Destroy()
            return

        value = dlg.GetValue()/100.
        dlg.Destroy()

        fig,ax = self.mylazgrid.plot_laz(self.active_vector.asshapely_ls(), length_buffer=value, show=False)

        if self.active_array is not None:
            copy_vec = vector()
            copy_vec.myvertices = self.active_vector.myvertices.copy()
            copy_vec.split(abs(self.active_array.dx)/2., False)
            copy_vec.get_values_on_vertices(self.active_array)
            s,z = copy_vec.get_sz()
            notmasked = np.where(z != -99999.)
            ax.plot(s[notmasked], z[notmasked], c='black', linewidth=2.0)

        fig.show()

    def clip_laz_gridded(self):
        """ Clip laz grid on current zoom """

        if self.mylazgrid is None:
            return

        curbounds = [[self.xmin, self.xmin + self.width], [self.ymin, self.ymin + self.height]]

        self.mylazdata = self.mylazgrid.scan(curbounds)

        logging.info(_('Clip LAZ grid on current zoom {}-{}  {}-{}').format(curbounds[0][0],curbounds[0][1],curbounds[1][0],curbounds[1][1]))

    def select_active_array_from_laz(self, array:WolfArray = None, used_codes:list = None, chunk_size:float = 500.):
        """ select some nodes from laz data

        :param array: array to fill
        :param used_codes: codes to use
        """
        if self.mylazgrid is None:
            return

        if array is None:
            logging.error(_('No array'))
            return

        if used_codes is None:
            keycode = [key for key,val in self.mylazgrid.colors.classification.items()]
            names = [val[0] for key,val in self.mylazgrid.colors.classification.items()]

            with wx.MultiChoiceDialog(None, _('Choose the codes to use'), _('Codes'), names) as dlg:
                if dlg.ShowModal() == wx.ID_OK:
                    used_codes = dlg.GetSelections()
                    used_codes = [float(keycode[cur]) for cur in used_codes]
                else:
                    return

        curbounds = array.get_bounds()

        # align bounds on chunk_size
        curbounds[0][0] = curbounds[0][0] - curbounds[0][0] % chunk_size
        curbounds[0][1] = curbounds[0][1] + chunk_size - curbounds[0][1] % chunk_size
        curbounds[1][0] = curbounds[1][0] - curbounds[1][0] % chunk_size
        curbounds[1][1] = curbounds[1][1] + chunk_size - curbounds[1][1] % chunk_size

        chunck_x = np.arange(curbounds[0][0], curbounds[0][1], chunk_size)
        chunck_y = np.arange(curbounds[1][0], curbounds[1][1], chunk_size)

        for curx in tqdm(chunck_x, 'Chunks'):
            for cury in chunck_y:
                curbounds = [[curx, curx + chunk_size], [cury, cury + chunk_size]]

                logging.info(_('Scan {}-{}  {}-{}').format(curbounds[0][0],curbounds[0][1],curbounds[1][0],curbounds[1][1]))
                self.mylazdata = self.mylazgrid.scan(curbounds)
                # logging.info(_('Scan done'))

                data = {}
                for curcode in used_codes:
                    data[curcode] = self.mylazdata[self.mylazdata[:, 3] == curcode]

                for curdata in data.values():

                    if curdata.shape[0] == 0:
                        continue

                    i,j = array.get_ij_from_xy(curdata[:, 0], curdata[:, 1]) #= np.float32(self.mylazdata[:, 2])

                    keys = np.vstack((i,j)).T

                    # unique keys
                    keys = np.unique(keys, axis=0)

                    array.SelectionData._add_nodes_to_selectionij(keys, verif = False)

        array.SelectionData.update_nb_nodes_selection()
        self.Paint()

        logging.info(_('Selection done'))

    def fill_active_array_from_laz(self, array:WolfArray = None, used_codes:list = [], operator:int = -1, chunk_size:float = 500.):
        """ Fill active array with laz data

        :param array: array to fill
        :param used_codes: codes to use
        :param operator: operator to use
        """

        if self.mylazgrid is None:
            return

        if array is None:
            logging.error(_('No array'))
            return

        if len(used_codes) == 0 :
            keycode = [key for key,val in self.mylazgrid.colors.classification.items()]
            names = [val[0] for key,val in self.mylazgrid.colors.classification.items()]

            with wx.MultiChoiceDialog(None, _('Choose the codes to use'), _('Codes'), names) as dlg:
                if dlg.ShowModal() == wx.ID_OK:
                    data = {}
                    used_codes = dlg.GetSelections()
                    used_codes = [float(keycode[cur]) for cur in used_codes]
                else:
                    return

        if operator == -1:
            with wx.SingleChoiceDialog(None, _('Choose the operator'), _('Operator'), ['max', 'percentile 95', 'percentile 5', 'min', 'mean', 'median', 'sum']) as dlg:
                if dlg.ShowModal() == wx.ID_OK:
                    if dlg.GetStringSelection() == 'max':
                        operator = np.max
                    elif dlg.GetStringSelection() == 'min':
                        operator = np.min
                    elif dlg.GetStringSelection() == 'mean':
                        operator = np.mean
                    elif dlg.GetStringSelection() == 'median':
                        operator = np.median
                    elif dlg.GetStringSelection() == 'sum':
                        operator = np.sum
                    elif dlg.GetStringSelection() == 'percentile 95':
                        operator = lambda x: np.percentile(x, 95)
                    elif dlg.GetStringSelection() == 'percentile 5':
                        operator = lambda x: np.percentile(x, 5)
                else:
                    return

        with wx.NumberEntryDialog(None, _('Minimum number of points to operate'), _('Minimum'), _('Minimum points'), 1, 1, 20) as dlg:
            if dlg.ShowModal() == wx.ID_OK:
                minpoints = dlg.GetValue()
            else:
                return

        bounds = array.get_bounds()

        # align bounds on chunk_size
        bounds[0][0] = bounds[0][0] - bounds[0][0] % chunk_size
        bounds[0][1] = bounds[0][1] + chunk_size - bounds[0][1] % chunk_size
        bounds[1][0] = bounds[1][0] - bounds[1][0] % chunk_size
        bounds[1][1] = bounds[1][1] + chunk_size - bounds[1][1] % chunk_size

        chunks_x = np.arange(bounds[0][0], bounds[0][1], chunk_size)
        chunks_y = np.arange(bounds[1][0], bounds[1][1], chunk_size)

        for curx in tqdm(chunks_x, 'Chunks'):
            for cury in chunks_y:

                curbounds = [[curx, curx + chunk_size], [cury, cury + chunk_size]]

                logging.info(_('Scan {}-{}  {}-{}').format(curbounds[0][0],curbounds[0][1],curbounds[1][0],curbounds[1][1]))
                self.mylazdata = self.mylazgrid.scan(curbounds)
                # logging.info(_('Scan done'))

                if len(self.mylazdata) == 0:
                    continue

                # Test codes
                data = {}
                for curcode in used_codes:
                    data[curcode] = self.mylazdata[self.mylazdata[:, 3] == curcode]

                # Treat data for each code
                for curdata in data.values():

                    if curdata.shape[0] == 0:
                        continue
                    else:
                        logging.info(_('Code {} : {} points'.format(curdata[0,3], curdata.shape[0])))

                    # get i,j from x,y
                    i,j = array.get_ij_from_xy(curdata[:, 0], curdata[:, 1]) #= np.float32(self.mylazdata[:, 2])

                    # keep only valid points -- inside the array
                    used = np.where((i >=0) & (i < array.nbx) & (j >=0) & (j < array.nby))[0]

                    if len(used) == 0:
                        continue

                    i = i[used]
                    j = j[used]
                    z = curdata[used, 2]

                    # create a key array
                    keys = np.vstack((i,j)).T
                    # find unique keys
                    keys = np.unique(keys, axis=0)

                    # create a ijz array
                    ijz = np.vstack((i, j, z)).T

                    # sort ijz array according to keys
                    #
                    # the most important indice is the last one enumerated in lexsort
                    # see : https://numpy.org/doc/stable/reference/generated/numpy.lexsort.html
                    ijz = ijz[np.lexsort((ijz[:,1], ijz[:,0]))]

                    # find first element of each key
                    idx = np.where(np.abs(np.diff(ijz[:,0])) + np.abs(np.diff(ijz[:,1])) != 0)[0]

                    # add last element
                    idx = np.concatenate((idx, [ijz.shape[0]]))

                    assert len(idx) == keys.shape[0], 'Error in filling'

                    logging.info(_('Cells to fill : {}'.format(len(idx))))

                    # apply operator
                    vals = {}
                    start_ii = 0
                    for ii, key in enumerate(keys):
                        end_ii = idx[ii]+1

                        if end_ii - start_ii >= minpoints:
                            vals[(key[0], key[1])] = operator(ijz[start_ii:end_ii,2])

                        start_ii = end_ii

                    if len(vals) > 0:
                        # create a new ijz array
                        newijz = np.asarray([[key[0], key[1], val] for key, val in vals.items()], dtype = np.float32)

                        array.fillin_from_ijz(newijz)

        array.reset_plot()
        self.Paint()

        logging.info(_('Filling done'))

    def init_laz_from_numpy(self, fn=None):
        """ Read LAZ data stored in numpy array"""

        if fn is None:
            filternpz = "npz (*.npz)|*.npz|LAZ (*.laz)|*.laz|LAS (*.las)|*.las|all (*.*)|*.*"
            dlg = wx.FileDialog(None, _('Choose LAS file'), wildcard=filternpz)
            ret = dlg.ShowModal()
            if ret != wx.ID_OK:
                return

            fn = dlg.GetPath()

        self.mylazdata = read_laz(fn)
        # self.mylazdata_colors = Classification_LAZ()
        # self.mylazdata_colors.init_2023()

        if self.linked:
            if len(self.linkedList) > 0:
                for curframe in self.linkedList:
                    curframe.mylazdata = self.mylazdata

    def init_laz_from_gridinfos(self, dirlaz:str = None):

        if dirlaz is None:
            dlg = wx.DirDialog(None, _('Choose directory where LAZ data/gridinfo are stored'))
            ret = dlg.ShowModal()
            if ret != wx.ID_OK:
                return

            dirlaz = dlg.GetPath()

        self.mylazgrid = xyz_laz_grids(dirlaz)

        dlg = wx.SingleChoiceDialog(None, _('Choose the classification'), _('Classification'), ['SPW-Geofit 2023', 'SPW 2013-2014'], wx.CHOICEDLG_STYLE)
        ret = dlg.ShowModal()
        if ret != wx.ID_OK:
            dlg.Destroy()
            return

        if dlg.GetStringSelection() == 'SPW 2013-2014':
            self.mylazgrid.colors.init_2013()
        else:
            self.mylazgrid.colors.init_2023()

        dlg.Destroy()

        if self.linked:
            if len(self.linkedList) > 0:
                for curframe in self.linkedList:
                    curframe.mylazgrid = self.mylazgrid


    def managebanks(self):
        if self.notebookbanks is None:
            self.notebookbanks = PlotNotebook(self)
            self.mypagebanks = self.notebookbanks.add(_("Manager banks interpolator"), "ManagerInterp")

        msg = ''
        if self.active_cs is None:
            msg += _(' The is no cross section. Please activate the desired object !')

        if msg != '':
            dlg = wx.MessageBox(msg, 'Required action')
            return

        if self.active_cs.linked_zones is None:
            msg += _(' The active zones is None. Please link the desired object to the cross sections !\n')
        # if self.active_zone is None:
        #     msg+=_(' The active zone is None. Please activate the desired object !\n')

        if msg != '':
            dlg = wx.MessageBox(msg, 'Required action')
            return

        self.mypagebanks.pointing(self, self.active_cs, self.active_vector)
        self.notebookbanks.Show(True)

    def _set_fn_fnpos_gltf(self):
        """
        Définition du nom de fichier GLTF/GLB à lire pour réaliser la comparaison
        Utilisation d'une fenêtre de dialogue WX

        Cette fonction n'est a priori appelée que depuis set_fn_fnpos_gltf
        """
        dlg = wx.FileDialog(None, _('Choose filename'),
                            wildcard='glb (*.glb)|*.glb|gltf2 (*.gltf)|*.gltf|All (*.*)|*.*', style=wx.FD_OPEN)
        ret = dlg.ShowModal()
        if ret == wx.ID_CANCEL:
            dlg.Destroy()
            return

        fn = dlg.GetPath()
        dlg.Destroy()

        dlg = wx.FileDialog(None, _('Choose pos filename'), wildcard='pos (*.pos)|*.pos|All (*.*)|*.*',
                            style=wx.FD_OPEN)
        ret = dlg.ShowModal()
        if ret == wx.ID_CANCEL:
            dlg.Destroy()
            return

        fnpos = dlg.GetPath()
        dlg.Destroy()

        if self.link_params is None:
            self.link_params = {}

        self.link_params['gltf file'] = fn
        self.link_params['gltf pos'] = fnpos

        return fn

    def set_fn_fnpos_gltf(self):
        """
        Définition ou récupération du nom de fichier GLTF/GLB à lire pour réaliser la comparaison

        Le nom de fichier est stocké dans la liste des paramètres partagés de façon à ce que l'appel de mise à jour puisse s'effectuer dans n'importe quel frame
        """
        fn = ''
        fnpos = ''
        if self.linked:
            for curgui in self.linkedList:
                if curgui.link_params is not None:
                    if 'gltf file' in curgui.link_params.keys():
                        fn = curgui.link_params['gltf file']
                        fnpos = curgui.link_params['gltf pos']
                        break
        elif self.link_params is None:
            self.link_params = {}
            fn = self._set_fn_fnpos_gltf()

        if fn == '':
            self._set_fn_fnpos_gltf()

    def read_last_result(self):
        """Lecture du dernier résultat pour les modèles ajoutés et plottés"""

        self.currently_readresults = True

        for curmodel in self.iterator_over_objects(draw_type.RES2D):
            curmodel: Wolfresults_2D
            logging.info(_('Updating {} - Last result'.format(curmodel.idx)))

            curmodel.read_oneresult()
            curmodel.set_currentview()

        self.Refresh()
        self.currently_readresults = False
        self._update_mytooltip()

    def read_one_result(self, which:int):
        """
        Lecture d'un résultat spécific pour les modèles ajoutés et plottés
        """
        self.currently_readresults = True

        for curmodel in self.iterator_over_objects(draw_type.RES2D):
            curmodel: Wolfresults_2D
            if curmodel.checked:
                logging.info(_('Updating {} - Specific result {}'.format(curmodel.idx, which)))

                curmodel.read_oneresult(which)
                curmodel.set_currentview()

        self.Refresh()
        self.currently_readresults = False
        self._update_mytooltip()


    def simul_previous_step(self):
        """
        Mise à jour au pas précédent
        """
        self.currently_readresults = True

        for curmodel in self.iterator_over_objects(draw_type.RES2D):
            curmodel: Wolfresults_2D
            logging.info(_('Updating {} - Previous result'.format(curmodel.idx)))

            curmodel.read_previous()
            curmodel.set_currentview()

        self.Refresh()
        self.currently_readresults = False
        self._update_mytooltip()

    def particle_next_step(self):
        """ Mise à jour au pas suivant """

        for curps in self.iterator_over_objects(draw_type.PARTICLE_SYSTEM):
            curps: Particle_system
            logging.info(_('Updating {} - Next result'.format(curps.idx)))

            curps.next_step()

        self._update_mytooltip()
        self.Refresh()

    def particle_previous_step(self):
        """ Mise à jour au pas précédent """

        for curps in self.iterator_over_objects(draw_type.PARTICLE_SYSTEM):
            curps: Particle_system
            logging.info(_('Updating {} - Next result'.format(curps.idx)))

            curps.previous_step()

        self._update_mytooltip()
        self.Refresh()

    def simul_next_step(self):
        """
        Mise à jour au pas suivant
        """
        self.currently_readresults = True

        for curmodel in self.iterator_over_objects(draw_type.RES2D):
            curmodel: Wolfresults_2D
            logging.info(_('Updating {} - Next result'.format(curmodel.idx)))

            curmodel.read_next()
            curmodel.set_currentview()

        self.Refresh()
        self.currently_readresults = False
        self._update_mytooltip()

    def OnMenuHighlight(self, event:wx.MenuEvent):

        id = event.GetId()
        item:wx.MenuItem
        item = self.menubar.FindItemById(event.GetId())

        if item is not None:
            self.set_statusbar_text(item.GetHelp())

    def _select_laz_source(self):
        """ Select laz source """

        if self.mylazdata is None and self.mylazgrid is None:
            logging.warning(_('No LAZ data loaded !'))
            return None
        elif self.mylazdata is None:
            laz_source = self.mylazgrid
        elif self.mylazgrid is None:
            laz_source = self.mylazdata
        else:
            choices = [_('From current loaded data'), _('From GridInfos')]
            dlg = wx.SingleChoiceDialog(None, _("Pick a data source"), "Choices", choices)
            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return None

            source = dlg.GetStringSelection()
            idx = choices.index(source)
            dlg.Destroy()

            if idx == 0:
                laz_source = self.mylazdata
            else:
                laz_source = self.mylazgrid

        return laz_source

    def OnMenubar(self, event: wx.MenuEvent):
        """
        Gestion des clicks sur le menu quel que soit le niveau

        Idée générale :
            - récupérer le label du menu sur base de l'id de l'event WX passé en argument --> itemlabel
            - tester le label du menu sur base de la chaîne traduite
            - a priori appeler une autre routine spécifique au traitement choisi
            - éviter autant que possible de coder des fonctions directement dans cette routine ce qui la rendrait complexe à lire

        AUTRE POSSIBILITE:
            - mettre en place un dictionnaire avec key==label, value==action qui se contenterait de tester la présence du label dans les clés et d'appeler l'action
            - dans ce dernier cas, il faudrait que les routines possèdent idéalement une interface unique
        """
        id = event.GetId()
        item = self.menubar.FindItemById(event.GetId())

        if item is None:
            return

        itemlabel = item.ItemLabel

        autoscale = True

        if id == wx.ID_OPEN:
            filterProject = "proj (*.proj)|*.proj|param (*.param)|*.param|all (*.*)|*.*"
            file = wx.FileDialog(self, "Choose file", wildcard=filterProject)
            if file.ShowModal() == wx.ID_CANCEL:
                file.Destroy()
                return
            else:
                # récuparétaion du nom de fichier avec chemin d'accès
                filename = file.GetPath()
                file.Destroy()

            os.chdir(os.path.dirname(filename))
            self.read_project(filename)

        elif itemlabel == _('Shortcuts'):
            # show shortcuts in log
            self.print_shortcuts(True)

        elif itemlabel == _('Show logs/informations'):
            self.check_logging()

        elif itemlabel == _('Show values'):
            self.check_tooltip()

        elif itemlabel == _('About'):
            #print About Frame
            self.print_About()

        elif itemlabel == _('Check for updates'):
            # check for new version

            self.check_for_updates()

        elif itemlabel == _("Integrate Q along active vector..."):
            """ Integrate Q along active vector """

            if self.active_vector is None:
                logging.warning(_('No active vector !'))
                return

            if self.active_res2d is None:
                logging.warning(_('No active 2D result !'))
                return

            self.active_res2d.plot_q(self.active_vector, 'border', toshow=True)

        elif itemlabel == _("Integrate Q along active zone..."):
            """ Integrate Q along active zone """

            if self.active_zone is None:
                logging.warning(_('No active zone !'))
                return

            if self.active_res2d is None:
                logging.warning(_('No active 2D result !'))
                return

            self.active_res2d.plot_q(self.active_zone.myvectors, ['border'] * self.active_zone.nbvectors, toshow=True)

        elif itemlabel == _("Plot stats unknown (selected nodes)..."):

            if self.active_res2d is None:
                logging.warning(_('No active 2D result !'))
                return

            dlg = wx.SingleChoiceDialog(None, _('Choose the unknown to plot'), _('Unknown'), [_('Water depth'), _('Water level'), _('Head')])

            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            unknown = dlg.GetStringSelection()
            dlg.Destroy()

            if unknown == _('Water depth'):
                unknown = 'h'
            elif unknown == _('Water level'):
                unknown = 'z'
            elif unknown == _('Head'):
                unknown = 'head'

            figax = None
            for curblock in self.active_res2d.myblocks.values():
                if curblock.SelectionData.nb > 0:

                    figax = self.active_res2d.plot_h(curblock.SelectionData.myselection,
                                                   unknown, toshow=False, figax=figax)
            if figax is not None:
                fig, ax = figax
                fig.show()

        elif itemlabel == _("Plot stats unknown (inside active vector)..."):

            if self.active_res2d is None:
                logging.warning(_('No active 2D result !'))
                return

            if self.active_vector is None:
                logging.warning(_('No active vector !'))
                return

            dlg = wx.SingleChoiceDialog(None, _('Choose the unknown to plot'), _('Unknown'), [_('Water depth'), ['Water level'], ['Head']])

            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            unknown = dlg.GetStringSelection()
            dlg.Destroy()

            if unknown == _('Water depth'):
                unknown = 'h'
            elif unknown == _('Water level'):
                unknown = 'z'
            elif unknown == _('Head'):
                unknown = 'head'

            fig, ax = self.active_res2d.plot_h(self.active_vector, unknown, toshow=True)

        elif itemlabel == _("Plot stats unknown (inside active zone)..."):

            if self.active_res2d is None:
                logging.warning(_('No active 2D result !'))
                return

            if self.active_zone is None:
                logging.warning(_('No active zone !'))
                return

            dlg = wx.SingleChoiceDialog(None, _('Choose the unknown to plot'), _('Unknown'), [_('Water depth'), ['Water level'], ['Head']])

            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            unknown = dlg.GetStringSelection()
            dlg.Destroy()

            if unknown == _('Water depth'):
                unknown = 'h'
            elif unknown == _('Water level'):
                unknown = 'z'
            elif unknown == _('Head'):
                unknown = 'head'

            for idx, curvect in enumerate(self.active_zone.myvectors):
                if idx ==0:
                    fig, ax = self.active_res2d.plot_h(curvect, unknown, toshow=False)
                else:
                    self.active_res2d.plot_h(curvect, unknown, toshow=False, figax = (fig, ax))

            fig.show()

        elif itemlabel == _("Plot active vector..."):
            """ Plot data along active vector """

            if self.active_vector is None:
                logging.warning(_('No active vector !'))
                return

            add_cloud = False
            if self.active_cloud is not None:
                dlg = wx.MessageDialog(self, _('Do you want to plot the cloud ?'), style=wx.YES_NO)

                if dlg.ShowModal() == wx.ID_YES:
                    add_cloud = True

                    prox = wx.TextEntryDialog(None,_('Proximity [m] ?'), value = '5.0')
                    ret = prox.ShowModal()
                    if ret == wx.ID_CANCEL:
                        prox.Destroy()
                        return
                    try:
                        proxval = float(prox.GetValue())
                    except:
                        prox.Destroy()
                        logging.warning(_('Bad value -- Rety'))
                        return

                    tol = wx.TextEntryDialog(None,_('Tolerance [m] ?'), value = '0.5')
                    ret = tol.ShowModal()
                    if ret == wx.ID_CANCEL:
                        tol.Destroy()
                        return
                    try:
                        tolval = float(tol.GetValue())
                    except:
                        tol.Destroy()
                        logging.warning(_('Bad value -- Rety'))
                        return

                else:
                    add_cloud = False

                dlg.Destroy()

            fig, ax = self.active_vector.plot_mpl(True, False)

            linkedarrays = self.get_linked_arrays()

            self.active_vector.plot_linked(fig, ax, linkedarrays)

            if add_cloud:
                s, z = self.active_cloud.projectontrace(self.active_vector, return_cloud=False, proximity= proxval)

                ax.scatter( s, z, c='black', s=1.0)

                for curs, curz in zip(s,z):
                    ax.plot([curs, curs], [curz-tolval, curz+tolval], 'k--', linewidth=0.5)
                    ax.plot([curs-.1, curs+.1], [curz+tolval, curz+tolval], c='black', linewidth=0.5)
                    ax.plot([curs-.1, curs+.1], [curz-tolval, curz-tolval], c='black', linewidth=0.5)

                fig.canvas.draw()
                fig.canvas.flush_events()

        elif itemlabel == _("Export arrays as Geotif..."):

            self.export_results_as('geotiff')

        elif itemlabel == _("Export arrays as Shapefile..."):

            self.export_results_as('shape')

        elif itemlabel == _("Compute and apply unique colormap on all..."):

            self.uniquecolormap()

        elif itemlabel == _("Load and apply unique colormap on all..."):

            self.uniquecolormap(True)

        elif itemlabel == _("Force uniform in parts on all..."):
            self.uniforminparts_all(True)

        elif itemlabel == _("Force linear interpolation on all..."):
            self.uniforminparts_all(False)

        elif itemlabel == _("Load and apply mask (nap)..."):

            self.loadnap_and_apply()

        elif itemlabel == _("Filter inundation arrays..."):

            self.filter_inundation()

        elif itemlabel == _("Plot active polygons..."):

            if self.active_zone is None:
                logging.warning(_('No active zone ! -- please select a zone containing polygons !'))
                return

            plotzone:list[zone]
            plotzone = []
            zonename = self.active_zone.myname
            if '_left_' in zonename or '_right_' in zonename:

                logging.info(_('Left and Right polygons are detected'))

                testname = zonename.replace('_left_', '')
                testname = testname.replace('_right_', '')

                for curzone in self.active_zones.myzones:
                    if testname == curzone.myname.replace('_left_', '').replace('_right_', ''):
                        plotzone.append(curzone)

                msg = wx.MessageDialog(self,
                                       _('Left and Right polygons are detected \nDo you want like to plot left and right polygons on the same plot ?'),
                                       style=wx.YES_NO | wx.YES_DEFAULT)
                ret = msg.ShowModal()
                msg.Destroy()
                if ret == wx.ID_NO:
                    plotzone = [self.active_zone]
            else:
                logging.info(_('Sole polygon detected'))
                plotzone = [self.active_zone]

            fig, ax = plt.subplots(1, 1)

            linkedarrays = {}

            # Matrices 2D
            for curarray in self.iterator_over_objects(draw_type.ARRAYS):
                curarray: WolfArray
                logging.info(_('Plotting array {}').format(curarray.idx))
                linkedarrays[curarray.idx] = curarray

            # Résultats 2D
            for curarray in self.iterator_over_objects(draw_type.RES2D):
                curarray: Wolfresults_2D
                logging.info(_('Plotting results {}').format(curarray.idx))
                linkedarrays[curarray.idx] = curarray

            linkedvecs={}
            for curvect in self.iterator_over_objects(draw_type.VECTORS):
                curvect: Zones
                logging.info(_('Plotting vector {}').format(curvect.idx))
                linkedvecs[curvect.idx] = curvect

            if len(plotzone) > 1:
                # left and right polygons
                for curzone in plotzone:
                    if '_left_' in curzone.myname:
                        locarrays = {}
                        for curkey, curarray in linkedarrays.items():
                            locarrays[curkey+ '_left'] = curarray

                        curzone.plot_linked_polygons(fig, ax, locarrays, linked_vec=linkedvecs, linestyle= '--')
                    elif '_right_' in curzone.myname:
                        locarrays = {}
                        for curkey, curarray in linkedarrays.items():
                            locarrays[curkey+ '_right'] = curarray

                        curzone.plot_linked_polygons(fig, ax, locarrays, linked_vec=linkedvecs, linestyle= '-.')
            else:
                # sole polygon
                plotzone[0].plot_linked_polygons(fig, ax, linkedarrays, linked_vec=linkedvecs)

            ax.grid()
            ax.legend()
            fig.show()

        # elif itemlabel == _("Change current view"):

        #     # Change view for results

        #     autoscale = False
        #     choices = [cur.value for cur in views_2D]
        #     dlg = wx.SingleChoiceDialog(None, _("Pick a view"), "Choices", choices)
        #     ret = dlg.ShowModal()
        #     if ret == wx.ID_CANCEL:
        #         dlg.Destroy()
        #         return

        #     method = dlg.GetStringSelection()

        #     method = list(views_2D)[choices.index(method)]

        #     dlg.Destroy()

        #     diamsize = None
        #     if method == views_2D.SHIELDS_NUMBER :

        #         if self.active_res2d is not None:
        #             sediment_diam = self.active_res2d.sediment_diameter
        #             sediment_density = self.active_res2d.sediment_density
        #         elif self.compare_results is not None:
        #             sediment_diam = 0.001
        #             sediment_density = 2.650
        #         else:
        #             logging.warning(_('No active 2D result or comparison !'))
        #             return

        #         dlg = wx.TextEntryDialog(None,_("Diameter grain size [m] ?"), value = str(sediment_diam))
        #         ret = dlg.ShowModal()
        #         if ret == wx.ID_CANCEL:
        #             dlg.Destroy()
        #             return
        #         try:
        #             diamsize = float(dlg.GetValue())
        #         except:
        #             dlg.Destroy()
        #             logging.warning(_("Bad value -- Rety"))
        #             return

        #         dlg = wx.TextEntryDialog(None,_("Density grain [-] ?"), value = str(sediment_density))
        #         ret = dlg.ShowModal()
        #         if ret == wx.ID_CANCEL:
        #             dlg.Destroy()
        #             return
        #         try:
        #             density = float(dlg.GetValue())
        #         except:
        #             dlg.Destroy()
        #             logging.warning(_("Bad value -- Rety"))
        #             return

        #     if len(self.myres2D)>1:

        #         dlg = wx.MessageDialog(None, _('Apply to all results?'), style=wx.YES_NO)
        #         ret = dlg.ShowModal()
        #         if ret == wx.ID_NO:
        #             if diamsize is not None:
        #                 self.active_res2d.sediment_diameter = diamsize
        #                 self.active_res2d.sediment_density = density
        #                 self.active_res2d.load_default_colormap('shields_cst')

        #             self.active_res2d.set_currentview(method, force_wx = True, force_updatepal = True)
        #         else:
        #             for curarray in self.iterator_over_objects(draw_type.RES2D):
        #                 curarray:Wolfresults_2D
        #                 if diamsize is not None:
        #                     curarray.sediment_diameter = diamsize
        #                     curarray.sediment_density  = density
        #                     curarray.load_default_colormap('shields_cst')

        #                 curarray.set_currentview(method, force_wx = True, force_updatepal = True)

        #     else:
        #         if self.active_res2d is not None:
        #             if diamsize is not None:
        #                 self.active_res2d.sediment_diameter = diamsize
        #                 self.active_res2d.sediment_density = density
        #                 self.active_res2d.load_default_colormap('shields_cst')
        #             self.active_res2d.set_currentview(method, force_wx = True, force_updatepal = True)

        #     if self.compare_results is not None:
        #         # update compare results
        #         if diamsize is not None:
        #             self.compare_results.set_shields_param(diamsize, density)
        #         self.compare_results.update_type_result(method)

        # elif itemlabel == _("Read last result"):

        #     self.read_last_result()

        # elif itemlabel == _("Filter independent"):

        #     self.menu_filter_independent.IsChecked = not self.menu_filter_independent.IsChecked

        #     for curmodel in self.iterator_over_objects(draw_type.RES2D):
        #         curmodel: Wolfresults_2D
        #         curmodel.to_filter_independent = not self.menu_filter_independent.IsChecked

        # elif itemlabel == _("Set epsilon water depth"):

        #     dlg = wx.TextEntryDialog(self, _('Enter an epsilon [m]'),value='0.0')

        #     ret = dlg.ShowModal()

        #     if ret == wx.ID_CANCEL:
        #         dlg.Destroy()
        #         return

        #     try:
        #         neweps = float(dlg.GetValue())
        #         dlg.Destroy()
        #     except:
        #         logging.error(_('Bad value -- retry !'))
        #         dlg.Destroy()
        #         return

        #     for curmodel in self.iterator_over_objects(draw_type.RES2D):
        #         curmodel: Wolfresults_2D
        #         curmodel.epsilon = neweps
        #         curmodel._epsilon_default = neweps
        #         curmodel.read_oneresult(curmodel.current_result)
        #         curmodel.set_currentview()

        # elif itemlabel == _("Manage boundary conditions..."):

        #     if self.active_res2d is not None:
        #         self.active_res2d.myparams.editing_bc(self.myres2D)

        # elif itemlabel ==_("Create video..."):
        #     if self.active_res2d is not None:
        #         self.create_video()

        elif itemlabel == _("Manage banks..."):
            if self.active_vector is None:
                msg = _('Active vector is None\nPlease activate the one desired')
                msg += _('\n')
                msg += _('Retry !\n')
                wx.MessageBox(msg)
                return

            self.managebanks()

        elif itemlabel == _("Create banks from vertices..."):

            self.active_cs.create_zone_from_banksbed()
            self.active_cs.linked_zones.showstructure()

        elif itemlabel == _("Link cross sections to active zones"):

            if self.active_cs is None:
                msg = _('Active cross sections is None\nPlease activate the one desired')
                msg += _('\n')
                msg += _('Retry !\n')
                wx.MessageBox(msg)
                return

            if self.active_zones is None:
                msg = _('Active zone is None\nPlease activate the one desired')
                msg += _('\n')
                msg += _('Retry !\n')
                wx.MessageBox(msg)
                return

            self.active_cs.link_external_zones(self.active_zones)


        elif itemlabel == _("Rename cross sections..."):

            dlg = wx.TextEntryDialog(None, _('Which starting point?'))
            ret = dlg.ShowModal()

            idxstart = dlg.GetValue()

            self.active_cs.rename(int(idxstart))

        elif itemlabel == _("Triangulate cross sections..."):
            self.triangulate_cs()

        # elif itemlabel == _("Import triangulation..."):
        #     self.import_3dfaces()

        elif itemlabel == _("Interpolate on active triangulation..."):
            self.interpolate_triangulation()

        elif itemlabel==_("Compare cloud to array..."):
            self.compare_cloud2array()

        elif itemlabel==_("Compare triangles to array..."):
            self.compare_tri2array()

        elif itemlabel ==  _("Create contour from checked arrays..."):

            # Create contour from checked arrays and add it to the list of objects
            newzones = self.create_Zones_from_arrays(self.get_list_objects(draw_type.ARRAYS, checked_state=True))
            self.add_object('vector', newobj=newzones, ToCheck=True, id='Contours from arrays')

        elif itemlabel == _("Calculator..."):

            if self.calculator is None:
                self.calculator = Calculator(mapviewer = self)
            else:
                self.calculator.Show()

        elif itemlabel == _("Create bridge and export gltf..."):

            if self.active_cs is None:
                msg = _('Active cross sections is None\nPlease activate the one desired')
                msg += _('\n')
                msg += _('Retry !\n')
                wx.MessageBox(msg)
                return

            self.start_action('bridge gltf', _('Create bridge and export gltf...'))

        elif itemlabel == _("Export cross sections to gltf..."):

            if self.active_cs is None:
                msg = _('Active cross sections is None\nPlease activate the one desired')
                msg += _('\n')
                msg += _('Retry !\n')
                wx.MessageBox(msg)
                return

            dlg = wx.TextEntryDialog(self, 'Z minimum ?', 'Choose an elevation as base')
            dlg.SetValue('')

            zmin = 0.
            if dlg.ShowModal() == wx.ID_OK:
                zmin = float(dlg.GetValue())
            dlg.Destroy()

            dlg = wx.FileDialog(None, _('Choose filename'),
                                wildcard='glb (*.glb)|*.glb|gltf2 (*.gltf)|*.gltf|All (*.*)|*.*', style=wx.FD_SAVE)
            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            fn = dlg.GetPath()
            dlg.Destroy()

            self.active_cs.export_gltf(zmin, fn)

        elif itemlabel == _("New cloud Viewer..."):
            if self.myinterp is not None:
                self.myinterp.viewer_interpolator()

        elif itemlabel == _("Interpolate on active array..."):
            if self.myinterp is not None:
                self.interpolate_cs()

        elif itemlabel == _("Interpolate active cloud on active array..."):
            self.interpolate_cloud()

        elif itemlabel == _('Save project'):
            filterProject = "proj (*.proj)|*.proj|param (*.param)|*.param|all (*.*)|*.*"
            file = wx.FileDialog(self, "Name your file", wildcard=filterProject, style=wx.FD_SAVE)
            if file.ShowModal() == wx.ID_CANCEL:
                file.Destroy()
                return
            else:
                # récuparétaion du nom de fichier avec chemin d'accès
                filename = file.GetPath()
                file.Destroy()

            self.save_project(filename)

        elif itemlabel == _('Initialize from npz'):
            autoscale=False
            self.init_laz_from_numpy()

        elif itemlabel == _('Initialize from GridInfos'):
            autoscale=False
            self.init_laz_from_gridinfos()

        elif itemlabel == _('Create cloud points from bridges'):
            autoscale=False
            if self.mylazdata is None:
                self.init_laz_from_numpy()

            mybridges = self.mylazdata[np.where(self.mylazdata[:, 3] == 10)]
            mycloud = cloud_vertices()
            mycloud.init_from_nparray(mybridges)
            mycloud.myprop.style = 2
            mycloud.myprop.color = getIfromRGB([255, 102, 102])
            mycloud.myprop.width = .5

            if self.linked:
                if len(self.linkedList) > 0:
                    for curframe in self.linkedList:
                        curframe.add_object('cloud', newobj=mycloud, ToCheck=True, id='Bridges')
            else:
                self.add_object('cloud', newobj=mycloud, ToCheck=True, id='Bridges')

        elif itemlabel == _('Create cloud points from buildings'):
            autoscale=False
            if self.mylazdata is None:
                self.init_laz_from_numpy()

            mybuildings = self.mylazdata[np.where(self.mylazdata[:, 3] == 1)]
            mycloud = cloud_vertices()
            mycloud.init_from_nparray(mybuildings)
            mycloud.myprop.style = 2
            mycloud.myprop.color = getIfromRGB([102, 102, 102])
            mycloud.myprop.width = .5
            if self.linked:
                if len(self.linkedList) > 0:
                    for curframe in self.linkedList:
                        curframe.add_object('cloud', newobj=mycloud, ToCheck=True, id='Buildings')
            else:
                self.add_object('cloud', newobj=mycloud, ToCheck=True, id='Buildings')

        elif itemlabel == _('Create LAZ viewer'):

            autoscale = False

            laz_source = self._select_laz_source()
            if laz_source is None:
                return

            choices, ass_values = choices_laz_colormap()
            dlg = wx.SingleChoiceDialog(None, _("Pick a colormap"), "Choices", choices)
            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            colormap = dlg.GetStringSelection()
            idx = choices.index(colormap)
            dlg.Destroy()

            if laz_source is self.mylazgrid:
                if self.mylazgrid is None:
                    logging.warning(_('No gridded LAZ data loaded !'))
                    return
                autoscale=False
                self.clip_laz_gridded()

                self.myviewer = myviewer(self.mylazdata, ass_values[idx], palette_classif = self.mylazgrid.colors)
            else:
                self.myviewer = myviewer(laz_source, ass_values[idx]) #, palette_classif = self.mylazdata_colors)

        elif itemlabel == _('Clip LAZ grid on current zoom'):
            if self.mylazgrid is None:
                logging.warning(_('No gridded LAZ data loaded !'))
                return
            autoscale=False
            self.clip_laz_gridded()

        elif itemlabel == _('Fill active array from LAZ data'):
            if self.mylazgrid is None:
                logging.warning('')
                return
            if self.active_array is None:
                logging.warning(_('No active array -- select an array first and retry!'))
                return

            autoscale = False
            self.fill_active_array_from_laz(self.active_array)

        elif itemlabel == _('Select cells in array from LAZ data'):
            if self.mylazgrid is None:
                logging.warning('')
                return
            if self.active_array is None:
                logging.warning(_('No active array -- select an array first and retry!'))
                return

            autoscale = False
            self.select_active_array_from_laz(self.active_array)

        elif itemlabel == _('Plot LAZ around active vector'):

            self.plot_laz_around_active_vec()

        elif itemlabel == _('Plot LAZ around temporary vector'):

            self.active_vector = vector()
            self.active_vector.add_vertex(wolfvertex(0.,0.))
            self.mimicme()

            self.start_action('laz tmp vector', _('LAZ tmp'))

        elif itemlabel == _('Change colors - Classification'):

            if self.mylazgrid is not None:
                self.mylazgrid.colors.interactive_update_colors()

        elif itemlabel == _('Multiviewer'):
            dlg = wx.NumberEntryDialog(self, _("Additional viewers"), _("How many?"), _("How many additional viewers?"),1, 0, 5)
            ret = dlg.ShowModal()

            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            nb = dlg.GetValue()
            dlg.Destroy()
            for i in range(nb):
                self.add_viewer_and_link()

        elif itemlabel == _('3D viewer'):

            self.active_viewer3d = Wolf_Viewer3D(self, _("3D Viewer"))
            self.active_viewer3d.Show()
            self.myviewers3d.append(self.active_viewer3d)

            for curarray in self.iterator_over_objects(draw_type.ARRAYS):
                curarray:WolfArray
                if curarray.checked:
                    if curarray._array3d is None:
                        curarray.prepare_3D()

                    if self.active_viewer3d not in curarray.viewers3d:
                        curarray.viewers3d.append(self.active_viewer3d)

                    self.active_viewer3d.add_array(curarray.idx, curarray._array3d)
                    self.active_viewer3d.autoscale()

            pass

        elif itemlabel == _('Create/Open multiblock model'):

            self.create_2D_MB_model()

        elif itemlabel == _('Open hydrological model'):

            self.open_hydrological_model()

        elif itemlabel == _('Check headers'):

            self.check_2D_MB_headers()

        elif itemlabel == _('Set comparison'):

            # Comparaison de deux résultats ou de deux matrices

            self.compare_results = Compare_Arrays_Results(self, True, True)

            add_elt = True
            while add_elt:
                add_elt = self.compare_results.add()

            if len(self.compare_results.paths) < 2 :
                logging.warning(_('Not enough elements to compare !'))
                self.compare_results = None
                return

            self.compare_results.bake()

        elif id == wx.ID_EXIT:
            dlg = wx.MessageDialog(None,_('Do you really want to quit?'), style = wx.YES_NO|wx.NO_DEFAULT)
            ret=dlg.ShowModal()
            if ret == wx.ID_YES:
                wx.Exit()
            else:
                dlg.Destroy()

        elif id == wx.ID_FILE1:
            self.add_object(which='array', ToCheck=True)

        elif itemlabel == _('Add view...'):
            self.add_object(which='views', ToCheck=True)

        elif itemlabel == _('Add tiles GPU...'):
            self.add_object(which='array_tiles', ToCheck=True)

        elif itemlabel == _('Add tiles...'):
            self.add_object(which='tiles', ToCheck=True)

        elif itemlabel == _('Add tiles comparator...'):
            self.add_object(which='tilescomp', ToCheck=True)

        elif id == wx.ID_FILE2:
            self.add_object(which='vector', ToCheck=True)

        elif id == wx.ID_FILE3:
            self.add_object(which='cloud', ToCheck=True)

        elif itemlabel == _('Add triangulation...'):
            self.add_object(which='triangulation', ToCheck=True)

        elif itemlabel == _('Add particle system...'):
            self.add_object(which = 'particlesystem', ToCheck = True)
            self.menu_particlesystem()

        elif itemlabel == _('Create particle system...'):
            self.active_particlesystem = newpart = Particle_system()
            self.add_object(which='particlesystem', newobj=newpart, ToCheck=True)
            self.menu_particlesystem()

        elif id == wx.ID_FILE4:
            self.add_object(which='cross_sections', ToCheck=True)

        elif itemlabel == _('Add Wolf2D results...'):
            self.add_object(which='res2d', ToCheck=True)
            self.menu_wolf2d()

        elif itemlabel == _('Add Wolf2D GPU results...'):
            self.add_object(which='res2d_gpu', ToCheck=True)
            self.menu_wolf2d()
            self.menu_2dgpu()

        elif itemlabel == _('Add bridges...'):
            self.add_object(which='bridges', ToCheck=True)

        elif itemlabel == _('Add weirs...'):
            self.add_object(which='weirs', ToCheck=True)

        elif itemlabel == _('Add array and crop...'):
            self.add_object(which='array_crop', ToCheck=True)

        elif itemlabel == _('Create array from bathymetry file...'):

            self.add_object(which='array_xyz', ToCheck=True)

        elif itemlabel == _('Create array from Lidar 2002...'):

            dlg = wx.SingleChoiceDialog(None, _('What source of data?'), _('Lidar 2002'),
                                        [_('First echo'), _('Second echo')])

            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                return

            sel = dlg.GetStringSelection()

            if sel == _('First echo'):
                self.add_object(which='array_lidar_first', ToCheck=True)
            elif sel == _('Second echo'):
                self.add_object(which='array_lidar_second', ToCheck=True)

        elif id == wx.ID_FILE5:
            def addscandir(mydir):
                for entry in scandir(mydir):
                    if entry.is_dir():
                        addscandir(entry)
                    elif entry.is_file():
                        if entry.name.endswith('.vec') or entry.name.endswith('.vecz'):

                            msg = wx.MessageDialog(self,
                                                   _(entry.name + ' found in ' + mydir + '\n\n Is it a "cross sections" file?'),
                                                   style=wx.YES_NO | wx.NO_DEFAULT)
                            ret = msg.ShowModal()
                            if ret == wx.ID_YES:
                                self.add_object(which='vector',
                                                filename=join(mydir, entry.name),
                                                ToCheck=True,
                                                id=join(mydir, entry.name))
                            else:
                                self.add_object(which='cross_sections',
                                                filename=join(mydir, entry.name),
                                                ToCheck=True,
                                                id=join(mydir, entry.name))

                        elif entry.name.endswith('.bin'):
                            self.add_object(which='array',
                                            filename=join(mydir, entry.name),
                                            ToCheck=True,
                                            id=join(mydir, entry.name))

            mydialog = wx.DirDialog(self, _("Choose directory to scan"))
            if mydialog.ShowModal() == wx.ID_CANCEL:
                mydialog.Destroy()
                return
            else:
                # récupération du nom de fichier avec chemin d'accès
                mydir = mydialog.GetPath()
                mydialog.Destroy()

            if exists(mydir):
                addscandir(mydir)

        elif id == wx.ID_FILE6:
            # Création d'une nouvelle matrice
            newarray = WolfArray(create=True, mapviewer=self)
            self.add_object('array', newobj=newarray)

        elif itemlabel == _('Create view...'):
            # Création d'une nouvelle vue
            newview = WolfViews(mapviewer=self)
            self.add_object('array', newobj=newarray)

        elif itemlabel==_('Create Wolf2D manager ...'):
            from .mesh2d.config_manager import config_manager_2D
            newmanager = config_manager_2D(mapviewer=self)

        elif itemlabel==_('Create scenarios manager ...'):
            from .scenario.config_manager import Config_Manager_2D_GPU
            newmanager = Config_Manager_2D_GPU(mapviewer=self)

        elif itemlabel == _('Create acceptability manager...'):

            from .acceptability.acceptability_gui import AcceptabilityGui
            newmanager = AcceptabilityGui()
            newmanager.mapviewer = self
            newmanager.Show()

        elif itemlabel==_('Create BC manager Wolf2D...'):

            if self.active_array is not None:

                choices = {'WOLF prev':1, 'WOLF OO':2, 'GPU':3}

                dlg = wx.SingleChoiceDialog(None,
                                            _("Which version of BC Manager"),
                                            _("Version"),
                                            ['WOLF prev', 'WOLF OO', 'GPU'])
                ret = dlg.ShowModal()
                if ret == wx.ID_CANCEL:
                    dlg.Destroy()
                    return

                method = dlg.GetStringSelection()
                dlg.Destroy()

                which_version = choices[method]

                self.mybc.append(BcManager(self,
                                           linked_array=self.active_array,
                                           version = which_version,
                                           DestroyAtClosing=False,
                                           Callback=self.pop_boundary_manager,
                                           mapviewer=self))
                ret = self.mybc[-1].FindBorders()
                if ret == -1:
                    self.mybc.pop(-1)
                    return
                self.active_bc = self.mybc[-1]

        elif itemlabel == _('Create Wolf1D...'):
            self.frame_create1Dfrom2D = GuiNotebook1D(mapviewer= self)
            logging.info(_(f'New window available - Wolf1D.'))

        elif id == wx.ID_FILE7:
            # Création de nouveaux vecteurs
            newzones = Zones(parent=self)
            self.add_object('vector', newobj=newzones)
        elif id == wx.ID_FILE8:
            # Création d'un nouveau nuage de point
            newcloud = cloud_vertices()
            self.add_object('cloud', newobj=newcloud)
        elif id in self.tools.keys():
            # gestion des actions
            self.ManageActions(id)
        elif id == wx.ID_SAVE:

            for obj in self.iterator_over_objects(draw_type.ARRAYS):
                obj: WolfArray

                if obj.filename == '':
                    filterArray = "bin (*.bin)|*.bin|Geotif (*.tif)|*.tif|Numpy (*.npy)|*.npy|all (*.*)|*.*"
                    fdlg = wx.FileDialog(self, "Choose file", wildcard=filterArray, style=wx.FD_SAVE)
                    fdlg.ShowModal()
                    if fdlg.ShowModal() == wx.ID_OK:
                        obj.filename = fdlg.GetPath()

                obj.write_all()

            for obj in self.iterator_over_objects(draw_type.VECTORS):
                obj:Zones
                obj.saveas()

        elif itemlabel == 'Save to image...':
            autoscale = False
            fn, ds = self.save_canvasogl()
            self.save_linked_canvas(fn[:-4], ds)

        elif itemlabel == _('Copy image...'):
            autoscale = False
            self.copy_canvasogl()

        elif itemlabel == _('Export...'):

            curarray: WolfArray
            curvec: vector

            msg = ''
            if self.active_array is None:
                msg += _('Active array is None\n')
            if self.active_vector is None:
                msg += _('Active vector is None\n')

            if msg != '':
                msg += _('\n')
                msg += _('Retry !\n')
                wx.MessageBox(msg)
                return

            curarray = self.active_array
            curvec = self.active_vector

            curvec.find_minmax()

            i1, j1 = curarray.get_ij_from_xy(curvec.xmin, curvec.ymin)
            x1, y1 = curarray.get_xy_from_ij(i1, j1)
            x1 -= curarray.dx / 2.
            y1 -= curarray.dy / 2.

            i2, j2 = curarray.get_ij_from_xy(curvec.xmax, curvec.ymax)
            x2, y2 = curarray.get_xy_from_ij(i2, j2)
            x2 += curarray.dx / 2.
            y2 += curarray.dy / 2.
            mybounds = [[x1, x2], [y1, y2]]

            dlg = wx.FileDialog(None, _('Choose filename'),
                                wildcard='glb (*.glb)|*.glb|gltf2 (*.gltf)|*.gltf|All (*.*)|*.*', style=wx.FD_SAVE)
            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            fn = dlg.GetPath()
            dlg.Destroy()

            with wx.lib.busy.BusyInfo(_('Export to gltf/glb')):
                wait = wx.BusyCursor()
                curarray.export_to_gltf(mybounds, fn)
                del wait

        elif itemlabel == _('Import...'):

            curarray: WolfArray

            msg = ''
            if self.active_array is None:
                msg += _('Active array is None\n')

            if msg != '':
                msg += _('\n')
                msg += _('Retry !\n')
                wx.MessageBox(msg)
                return

            curarray = self.active_array

            dlg = wx.FileDialog(None, _('Choose filename'),
                                wildcard='glb (*.glb)|*.glb|gltf2 (*.gltf)|*.gltf|All (*.*)|*.*', style=wx.FD_OPEN)
            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            fn = dlg.GetPath()
            dlg.Destroy()

            dlg = wx.FileDialog(None, _('Choose pos filename'), wildcard='pos (*.pos)|*.pos|All (*.*)|*.*',
                                style=wx.FD_OPEN)
            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            fnpos = dlg.GetPath()
            dlg.Destroy()

            choices = ["matplotlib", "scipy", "pyvista"]
            dlg = wx.SingleChoiceDialog(None, _("Pick an interpolation method"), _("Choices"), choices)
            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            method = dlg.GetStringSelection()
            dlg.Destroy()

            with wx.lib.busy.BusyInfo(_('Importing gltf/glb')):
                wait = wx.BusyCursor()
                try:
                    curarray.import_from_gltf(fn, fnpos, method)
                except:
                    pass
                del wait

        elif itemlabel == _('Compare...'):

            msg = ''
            if self.active_array is None:
                msg += _('Active array is None\n')

            if msg != '':
                msg += _('\n')
                msg += _('Retry !\n')
                wx.MessageBox(msg)
                return

            self.set_blender_sculpting()
            autoscale = False
            self.set_fn_fnpos_gltf()
            self.update_blender_sculpting()

        elif itemlabel == _('Update...'):
            autoscale = False

            msg = ''
            if self.active_array is None:
                msg += _('Active array is None\n')

            if msg != '':
                msg += _('\n')
                msg += _('Retry !\n')
                wx.MessageBox(msg)
                return

            self.set_fn_fnpos_gltf()
            self.update_blender_sculpting()

        elif id == wx.ID_SAVEAS:

            for obj in self.iterator_over_objects(draw_type.ARRAYS):
                obj: WolfArray

                filterArray = "bin (*.bin)|*.bin|Geotif (*.tif)|*.tif|Numpy (*.npy)|*.npy|all (*.*)|*.*"
                fdlg = wx.FileDialog(self, "Choose file name for Array : " + obj.idx, wildcard=filterArray,
                                        style=wx.FD_SAVE)
                ret = fdlg.ShowModal()
                if ret == wx.ID_OK:
                    obj.filename = fdlg.GetPath()
                    obj.write_all()

            for obj in self.iterator_over_objects(draw_type.VECTORS):
                obj:Zones
                if obj.idx=='grid':
                    pass
                else:
                    filterArray = "vec (*.vec)|*.vec|vecz (*.vecz)|*.vecz|Shapefile (*.shp)|*.shp|all (*.*)|*.*"
                    fdlg = wx.FileDialog(self, "Choose file name for Vector :" + obj.idx, wildcard=filterArray,
                                         style=wx.FD_SAVE)
                    ret = fdlg.ShowModal()
                    if ret == wx.ID_OK:
                        obj.saveas(fdlg.GetPath())

        # elif id == ID_RUN_SIMULATION:
        #     from ..debug.test_glsim.wolf_gpu import load_sim_to_gpu
        #     print("Simulation")
        #     load_sim_to_gpu()

        if len(self.myarrays) + len(self.myvectors) + len(self.myclouds) + len(self.mytri) + len(self.myres2D) + len(self.mytiles) + len(self.mypartsystems) == 2 and autoscale:
            # Trouve les bornzs si un seul élément est présent, sinon on conserve l'état du zoom
            self.Autoscale()

    def pop_boundary_manager(self, which:BcManager):
        """ Pop a boundary condition manager after Destroying """

        idx = self.mybc.index(which)
        if self.active_bc is which:
            self.active_bc = None
        self.mybc.pop(idx)

        self.Refresh()


    def get_boundary_manager(self, which:WolfArray):
        """ Get a boundary manager """

        for curbc in self.mybc:
            if curbc.linked_array is which:
                return curbc

        return None

    def uniquecolormap(self, loadfromfile = False):
        """ Compute unique colormap from all (arrays, 2D results) and apply it to all """

        workingarray=[]
        nbnotnull=0

        newpal = wolfpalette(self)

        if loadfromfile :
            newpal.readfile()
        else:
            with wx.lib.busy.BusyInfo(_('Compute unique colormap from all arrays')):
                wait = wx.BusyCursor()

                curarray:WolfArray
                curres2d:Wolfresults_2D

                for curarray in self.myarrays:
                    if curarray.plotted:
                        workingarray.append(curarray.get_working_array())
                        nbnotnull+=curarray.nbnotnull

                for curres2d in self.myres2D:
                    if curres2d.plotted:
                        workingarray.append(curres2d.get_working_array())
                        nbnotnull+=curres2d.nbnotnull

                workingarray = np.concatenate(workingarray)

                newpal.default16()
                newpal.isopop(workingarray, nbnotnull)
            del wait

        with wx.lib.busy.BusyInfo(_('Apply unique colormap to all arrays')):
            wait = wx.BusyCursor()
            for curarray in self.myarrays:
                if curarray.plotted:
                    curarray.mypal.automatic = False
                    curarray.myops.palauto.SetValue(0)
                    curarray.mypal.values = newpal.values.copy()
                    curarray.mypal.colors = newpal.colors.copy()
                    curarray.mypal.fill_segmentdata()
                    curarray.reset_plot()

            for curres2d in self.myres2D:
                if curres2d.plotted:
                    curres2d.mypal.automatic = False
                    curres2d.mypal.nb     = newpal.nb
                    curres2d.mypal.values = newpal.values.copy()
                    curres2d.mypal.colors = newpal.colors.copy()
                    curres2d.mypal.fill_segmentdata()
                    curres2d.reset_plot()
            del wait

    def loadnap_and_apply(self):

        dlg = wx.MessageDialog(self,_('Load mask for all?'),style=wx.YES_NO|wx.YES_DEFAULT)
        ret=dlg.ShowModal()

        if ret == wx.ID_NO:
            dlg.Destroy()
            return

        with wx.lib.busy.BusyInfo(_('Loading masks')):
            wait = wx.BusyCursor()
            curarray:WolfArray
            for curarray in self.myarrays:
                if curarray.plotted:
                    curarray.loadnap_and_apply()
            del wait

    def uniforminparts_all(self, TrueOrFalse:bool):

        for curarray in self.myarrays:
            curarray:WolfArray
            if curarray.plotted:
                curarray.mypal.interval_cst = TrueOrFalse
                curarray.reset_plot()

        for curarray in self.myres2D:
            curarray:Wolfresults_2D
            if curarray.plotted:
                curarray.mypal.interval_cst = TrueOrFalse
                curarray.link_palette()
                curarray.reset_plot()

    def filter_inundation(self):

        dlg = wx.TextEntryDialog(self,_('Upper bound \n\n All values strictly lower than the bound will not be extracted !'),value='.0005')
        ret=dlg.ShowModal()

        if ret == wx.ID_CANCEL:
            dlg.Destroy()
            return

        bound = float(dlg.GetValue())
        dlg.Destroy()

        logging.info(_('Filtering results'))

        curarray:WolfArray
        for curarray in self.myarrays:
            if curarray.plotted:
                curarray.filter_inundation(epsilon = bound)
                curarray.filter_independent_zones(n_largest = 1)

        curarray:Wolfresults_2D
        for curarray in self.myres2D:
            if curarray.plotted:
                curarray.filter_inundation(eps = bound)
                curarray.filter_independent_zones(n_largest = 1)

        logging.info(_('Filtering done !'))

    def export_results_as(self,which='geotiff'):

        """
        Export des résultats WOLF2D vers le format GeoTiff
        On attend que les matrices ".hbin" aient été chargées dans l'interface

        TODO : Modifier la routine pour prendre les classe de résultats 2D
        """

        dlg = wx.DirDialog(self,_('Choose output directory'),style = wx.DD_DIR_MUST_EXIST)
        ret=dlg.ShowModal()

        if ret == wx.ID_CANCEL:
            dlg.Destroy()
            return

        outdir = dlg.GetPath()
        dlg.Destroy()

        with wx.lib.busy.BusyInfo(_('Exporting arrays')):
            wait = wx.BusyCursor()
            curarray:WolfArray
            for curarray in self.myarrays:
                if curarray.plotted:

                    fn = os.path.splitext(curarray.filename)[0] #curarray.filename[:-4]
                    top = WolfArray(fn+'.topini_fine')
                    top.array.mask = curarray.array.mask
                    top.array.data[np.where(top.array.mask)] = 0.
                    qx = WolfArray(fn+'.qxbin')
                    qx.array.mask = curarray.array.mask
                    qx.array.data[np.where(qx.array.mask)] = 0.
                    qy = WolfArray(fn+'.qybin')
                    qy.array.mask = curarray.array.mask
                    qy.array.data[np.where(qy.array.mask)] = 0.

                    qnorm = (qx**2.+qy**2.)**.5
                    vnorm=qnorm/curarray
                    froude=vnorm/(curarray*9.81)**.5

                    frott = WolfArray(fn+'.frot')

                    def compute_critdiam(h:WolfArray,qnorm:WolfArray,n:WolfArray):

                        ij = np.argwhere(h.array>0.)

                        diamcrit_shields = WolfArray(mold=h)
                        diamcrit_izbach = WolfArray(mold=h)

                        diam = np.asarray([get_d_cr(qnorm.array[i,j],h.array[i,j],1./n.array[i,j]) for i,j in ij])

                        diamcrit_shields.array[ij[:,0],ij[:,1]] = diam[:,0]
                        diamcrit_izbach.array[ij[:,0],ij[:,1]] = diam[:,1]

                        return diamcrit_shields,diamcrit_izbach

                    shields,izbach = compute_critdiam(curarray,qnorm,frott)

                    myarrays=[top,curarray,qx,qy,vnorm,froude,shields,izbach]
                    mynames=['Z [mDNG]',
                             'H [m]',
                             'QX [m2s-1]',
                             'QY [m2s-1]',
                             'Un [ms-1]',
                             'Fr [-]',
                             'D_Sh [m]',
                             'D_Iz [m]']

                    if which =='geotiff':
                        self.export_geotif(outdir,curarray.idx,myarrays,mynames)

                    elif which=='shape':
                        self.export_shape(outdir,curarray.idx,myarrays,mynames,curarray)
            del wait

    def export_shape(self, outdir='', fn = '', myarrays=[], descr=[], mask:WolfArray=None):
        """ Export multiple arrays to shapefile

        :param outdir: output directory
        :param fn: filename -- .shp will be added if not present
        :param myarrays: list of Wolfarrays to export
        :param descr: list of descriptions
        :param mask: mask array -- export only where mask > 0
        """

        if len(myarrays)==0:
            logging.warning(_('No arrays provided for shapefile export'))
            return

        if mask is None:
            logging.warning(_('No mask provided for shapefile export'))
            return

        from osgeo import gdal, osr, gdalconst,ogr

        # create the spatial reference system, Lambert72
        srs = osr.SpatialReference()
        srs.ImportFromEPSG(31370)

        # create the data source
        driver: ogr.Driver
        driver = ogr.GetDriverByName("ESRI Shapefile")

        # create the data source
        filename = join(outdir,fn)
        if not filename.endswith('.shp'):
            filename+='.shp'

        ds = driver.CreateDataSource(filename)

        # create one layer
        layer = ds.CreateLayer("results", srs, ogr.wkbPolygon)

        # Add ID fields
        idFields=[]
        for curlab in descr:
            idFields.append(ogr.FieldDefn(curlab, ogr.OFTReal))
            layer.CreateField(idFields[-1])

        # Create the feature and set values
        featureDefn = layer.GetLayerDefn()
        feature = ogr.Feature(featureDefn)

        usednodes = np.argwhere(mask.array>0.)
        for i,j in usednodes:

            x,y = mask.get_xy_from_ij(i,j)
            # Creating a line geometry
            ring = ogr.Geometry(ogr.wkbLinearRing)
            ring.AddPoint(x-mask.dx/2,y-mask.dy/2)
            ring.AddPoint(x+mask.dx/2,y-mask.dy/2)
            ring.AddPoint(x+mask.dx/2,y+mask.dy/2)
            ring.AddPoint(x-mask.dx/2,y+mask.dy/2)
            ring.AddPoint(x-mask.dx/2,y-mask.dy/2)

            # Create polygon
            poly = ogr.Geometry(ogr.wkbPolygon)
            poly.AddGeometry(ring)

            feature.SetGeometry(poly)

            for arr, id in zip(myarrays,descr):

                feature.SetField(id, float(arr.array[i,j]))

            layer.CreateFeature(feature)

        feature = None

        # Save and close DataSource
        ds = None

    def export_geotif(self, outdir='', fn = '', myarrays=[], descr=[]):
        """ Export multiple arrays to geotiff

        :param outdir: output directory
        :param fn: filename -- .tif will be added if not present
        :param myarrays: list of Wolfarrays to export
        :param descr: list of descriptions -- Bands names

        """

        if len(myarrays)==0:
            logging.warning(_('No arrays provided for geotiff export'))
            return

        from osgeo import gdal, osr, gdalconst

        srs = osr.SpatialReference()
        srs.ImportFromEPSG(31370)

        filename = join(outdir,fn)
        if not filename.endswith('.tif'):
            filename+='.tif'

        arr=myarrays[0].array
        if arr.dtype == np.float32:
            arr_type = gdal.GDT_Float32
        else:
            arr_type = gdal.GDT_Int32

        driver: gdal.Driver
        out_ds: gdal.Dataset
        band: gdal.Band
        driver = gdal.GetDriverByName("GTiff")
        out_ds = driver.Create(filename, arr.shape[0], arr.shape[1], len(myarrays), arr_type, options=['COMPRESS=LZW'])
        out_ds.SetProjection(srs.ExportToWkt())
        out_ds.SetGeoTransform([myarrays[0].origx+myarrays[0].translx,
                                myarrays[0].dx,
                                0.,
                                myarrays[0].origy+myarrays[0].transly,
                                0.,
                                myarrays[0].dy])

        k=1
        for arr, name in zip(myarrays,descr):
            band = out_ds.GetRasterBand(k)
            band.SetNoDataValue(0.)
            band.SetDescription(name)
            band.WriteArray(arr.array.transpose())
            band.FlushCache()
            band.ComputeStatistics(True)
            k+=1

    def get_linked_arrays(self) -> dict:
        linkedarrays = {}

        for locarray in self.iterator_over_objects(draw_type.ARRAYS):
            linkedarrays[locarray.idx] = locarray

        for locarray in self.iterator_over_objects(draw_type.RES2D):
            linkedarrays[locarray.idx] = locarray

        return linkedarrays

    def save_linked_canvas(self, fn, mpl=True, ds=0.):
        if self.linked:
            for idx, curel in enumerate(self.linkedList):
                curel.save_canvasogl(fn + '_' + str(idx) + '.png', mpl, ds)

    def thread_update_blender(self):
        print("Update blender")
        if self.SetCurrentContext():
            self.update_blender_sculpting()
            t = threading.Timer(10.0, self.thread_update_blender)
            t.start()

    def add_object(self,
                   which:Literal['array','array_lidar_first','array_lidar_second','array_xyz','array_tiles',
                                 'bridges',
                                 'weirs',
                                 'vector',
                                 'tiles', 'tilescomp'
                                 'cloud',
                                 'triangulation',
                                 'cross_sections',
                                 'other',
                                 'views',
                                 'res2d',
                                 'res2d_gpu',
                                 'particlesystem',
                                 'wmsback',
                                 'wmsfore'] = 'array',
                   filename='',
                   newobj=None,
                   ToCheck=True,
                   id=''):

        """
        Add object to current Frame/Drawing area
        """

        filterArray = "All supported formats|*.bin;*.tif;*.tiff;*.top;*.flt;*.npy;*.npz|bin (*.bin)|*.bin|Elevation WOLF2D (*.top)|*.top|Geotif (*.tif)|*.tif|Float ESRI (*.flt)|*.flt|Numpy (*.npy)|*.npy|Numpy named arrays(*.npz)|*.npz|all (*.*)|*.*"
        filterjson = "json (*.json)|*.json|all (*.*)|*.*"
        filterall = "all (*.*)|*.*"
        filterres2d = "all (*.*)|*.*"
        filterVector = "vec (*.vec)|*.vec|vecz (*.vecz)|*.vecz|dxf (*.dxf)|*.dxf|shp (*.shp)|*.shp|all (*.*)|*.*"
        filterCloud = "xyz (*.xyz)|*.xyz|dxf (*.dxf)|*.dxf|text (*.txt)|*.txt|shp (*.shp)|*.shp|all (*.*)|*.*"
        filtertri = "tri (*.tri)|*.tri|text (*.txt)|*.txt|dxf (*.dxf)|*.dxf|gltf (*.gltf)|*.gltf|gltf binary (*.glb)|*.glb|*.*'all (*.*)|*.*"
        filterCs = "vecz WOLF (*.vecz)|*.vecz|txt 2022 (*.txt)|*.txt|WOLF (*.sxy)|*.sxy|text 2000 (*.txt)|*.txt|all (*.*)|*.*"
        filterimage = "Geotif (*.tif)|*.tif|all (*.*)|*.*"

        if filename == '' and newobj is None:
            # ouverture d'une boîte de dialogue
            if which.lower() == 'array' or which.lower() == 'array_crop':
                file = wx.FileDialog(self, "Choose file", wildcard=filterArray)
            elif which.lower() == 'particlesystem':
                file = wx.FileDialog(self, "Choose file", wildcard=filterjson)
            elif which.lower() == 'array_lidar_first' or which.lower() == 'array_lidar_second':
                file = wx.DirDialog(self, "Choose directory containing Lidar data")
            elif which.lower() == 'array_xyz':
                file = wx.DirDialog(self, "Choose directory containing XYZ files")
            elif which.lower() == 'array_tiles':
                file = wx.DirDialog(self, "Choose directory containing GPU results")
            elif which.lower() == 'bridges':
                file = wx.DirDialog(self, "Choose directory containing bridges")
            elif which.lower() == 'weirs':
                file = wx.DirDialog(self, "Choose directory containing weirs")
            elif which.lower() in ['vector', 'tiles', 'tilescomp']:
                file = wx.FileDialog(self, "Choose file", wildcard=filterVector)
            elif which.lower() == 'cloud':
                file = wx.FileDialog(self, "Choose file", wildcard=filterCloud)
            elif which.lower() == 'triangulation':
                file = wx.FileDialog(self, "Choose file", wildcard=filtertri)
            elif which.lower() == 'cross_sections':
                file = wx.FileDialog(self, "Choose file", wildcard=filterCs)
            elif which.lower() == 'other':
                file = wx.FileDialog(self, "Choose file", wildcard=filterall)
            elif which.lower() == 'views':
                file = wx.FileDialog(self, "Choose file", wildcard=filterall)
            elif which.lower() == 'res2d':
                file = wx.FileDialog(self, "Choose file", wildcard=filterres2d)
            elif which.lower() == 'res2d_gpu':
                file = wx.DirDialog(self, "Choose directory containging WolfGPU results")

            # FIXME : particularize filters for wmsback and wmsfore
            elif which.lower() == 'wmsback':
                file = wx.FileDialog(self, "Choose file", wildcard=filterimage)
            elif which.lower() == 'wmsfore':
                file = wx.FileDialog(self, "Choose file", wildcard=filterimage)

            if file.ShowModal() == wx.ID_CANCEL:
                file.Destroy()
                return -1
            else:
                # récuparétaion du nom de fichier avec chemin d'accès
                filename = file.GetPath()
                try:
                    curfilter = file.GetFilterIndex()
                except:
                    pass
                file.Destroy()

        if filename != '':
            if (not (os.path.exists(filename))):
                logging.warning("Warning : the following file is not present here : " + filename)
                return -1

        curtree = None
        if which.lower() == 'array' or which.lower() == 'array_crop':

            curdict = self.myarrays
            curtree = self.myitemsarray

            if newobj is None:

                if filename.endswith('.npz'):

                    wait = wx.BusyCursor()
                    logging.info(_('Start of importing arrays from npz file'))

                    with np.load(filename) as data:
                        if 'header' in data.keys():
                            header = data['header']

                            if len(header) == 6:
                                logging.info(_('Header found in npz file'))

                                origx, origy, dx, dy, nbx, nby = header

                                logging.info(_('Origin X : ') + str(origx))
                                logging.info(_('Origin Y : ') + str(origy))
                                logging.info(_('dx : ') + str(dx))
                                logging.info(_('dy : ') + str(dy))
                                logging.info(_('nbx : ') + str(nbx))
                                logging.info(_('nby : ') + str(nby))
                                nbx, nby = int(nbx), int(nby)
                            else:
                                logging.warning(_('Header found in npz file but not complete -- Only {} values found - Must be 6').format(len(header)))

                            for key, curarray in data.items():
                                if isinstance(curarray, np.ndarray):
                                    if curarray.shape == (nby, nbx):
                                        logging.info("Importing array : " + key)
                                        curhead = header_wolf()
                                        curhead.origx, curhead.origy, curhead.dx, curhead.dy, curhead.nbx, curhead.nby = origx, origy, dx, dy, nbx, nby
                                        newobj = WolfArray(srcheader=curhead, idx = key)
                                        newobj.set_array_from_numpy(curarray)
                                        self.add_object('array', newobj= newobj, id= key)
                        else:
                            origx, origy, dx, dy, nbx, nby = 0.,0.,1,1.,1,1
                            for key, curarray in data.items():
                                if isinstance(curarray, np.ndarray):
                                    logging.info(_('No header found in npz file - Using default values for header'))
                                    logging.info("Importing array : " + key)
                                    curhead = header_wolf()
                                    curhead.origx, curhead.origy, curhead.dx, curhead.dy, curhead.nbx, curhead.nby = 0., 0., 1., 1., curarray.shape[0], curarray.shape[1]
                                    newobj = WolfArray(srcheader=curhead, idx = key)
                                    newobj.set_array_from_numpy(curarray)
                                    self.add_object('array', newobj= newobj, id= key)

                    logging.info(_('End of importing arrays from npz file'))
                    del wait
                    return -1
                else:
                    testobj = WolfArray()
                    testobj.filename = filename
                    testobj.read_txt_header()

                    if testobj.wolftype in WOLF_ARRAY_MB:
                        with wx.lib.busy.BusyInfo(_('Importing array')):
                            wait = wx.BusyCursor()
                            newobj = WolfArrayMB(filename, mapviewer=self)
                            del wait
                    else:
                        if which.lower() == 'array_crop':
                            newobj = WolfArray(filename, mapviewer=self, crop='newcrop')
                        else:
                            with wx.lib.busy.BusyInfo(_('Importing array')):
                                wait = wx.BusyCursor()
                                newobj = WolfArray(filename, mapviewer=self)
                                del wait

                if newobj is not None:
                    if newobj.dx==0. or newobj.dy==0.:
                        dlg_pos = CropDialog(None)
                        dlg_pos.SetTitle(_('Choose informations'))

                        dlg_pos.ox.SetValue('99999.')
                        dlg_pos.oy.SetValue('99999.')

                        dlg_pos.ex.Hide()
                        dlg_pos.ey.Hide()

                        badvalues = True
                        while badvalues:
                            badvalues = False

                            ret = dlg_pos.ShowModal()
                            if ret == wx.ID_CANCEL:
                                newcrop.Destroy()
                                return -1
                            else:
                                cropini = [[float(dlg_pos.ox.Value), float(dlg_pos.ex.Value)],
                                                [float(dlg_pos.oy.Value), float(dlg_pos.ey.Value)]]
                                tmpdx = float(dlg_pos.dx.Value)
                                tmpdy = float(dlg_pos.dy.Value)

                            if tmpdx ==0. or tmpdy==0.:
                                badvalues = True

                        dlg_pos.Destroy()

                        newobj.dx = tmpdx
                        newobj.dy = tmpdy

                        # if newobj.SelectionData is not None:
                        #     newobj.SelectionData.dx = tmpdx
                        #     newobj.SelectionData.dy = tmpdy

                        if cropini[0][0] != 99999. and cropini[1][0]!=99999.:
                            newobj.origx = cropini[0][0]
                            newobj.origy = cropini[1][0]

            newobj.updatepalette(0)
            self.myarrays.append(newobj)
            newobj.change_gui(self)
            self.active_array = newobj
            self._set_active_bc()

        elif which.lower() == 'array_tiles':

            res = wolfres2DGPU(filename, plotted=False)

            tilesmap = res._result_store._tile_packer.tile_indirection_map()

            header = header_wolf()
            res_header = res[0].get_header()

            header.origx = res_header.origx
            header.origy = res_header.origy
            header.dx = res_header.dx * 16.
            header.dy = res_header.dy * 16.
            header.nbx = tilesmap.shape[1]
            header.nby = tilesmap.shape[0]

            newobj_i = WolfArray(mapviewer=self, srcheader=header, idx = 'tils_i')
            newobj_j = WolfArray(mapviewer=self, srcheader=header, idx = 'tils_j')

            newobj_i.array = np.ma.asarray(tilesmap[:,:,0].T.astype(np.float32))
            newobj_j.array = np.ma.asarray(tilesmap[:,:,1].T.astype(np.float32))

            newobj_i.mask_data(0.)
            newobj_j.mask_data(0.)

            self.add_object('array', newobj=newobj_i, id=newobj_i.idx)
            self.add_object('array', newobj=newobj_j, id=newobj_j.idx)

            return

        elif which.lower() == 'bridges':
            curdict = self.myvectors
            curtree = self.myitemsvector

            if newobj is None:
                with wx.lib.busy.BusyInfo(_('Importing files')):
                    wait = wx.BusyCursor()
                    newobj = Bridges(filename, mapviewer=self)
                    del wait
            self.myvectors.append(newobj)

        elif which.lower() == 'weirs':
            curdict = self.myvectors
            curtree = self.myitemsvector

            if newobj is None:
                with wx.lib.busy.BusyInfo(_('Importing files')):
                    wait = wx.BusyCursor()
                    newobj = Weirs(filename, mapviewer=self)
                    del wait
            self.myvectors.append(newobj)

        elif which.lower() in ['tiles', 'tilescomp']:
            curdict = self.mytiles
            curtree = self.myitemsvector

            if newobj is None:

                file = wx.DirDialog(self, "Choose directory containing data")
                if file.ShowModal() == wx.ID_CANCEL:
                    file.Destroy()
                    return -1
                else:
                    # récuparétaion du nom de fichier avec chemin d'accès
                    dirname = file.GetPath()
                    file.Destroy()

                if which.lower() == 'tilescomp':
                    file = wx.DirDialog(self, "Choose directory containing comparison data")
                    if file.ShowModal() == wx.ID_CANCEL:
                        file.Destroy()
                        return -1
                    else:
                        # récuparétaion du nom de fichier avec chemin d'accès
                        dirname_comp = file.GetPath()
                        file.Destroy()

                with wx.lib.busy.BusyInfo(_('Importing files')):
                    wait = wx.BusyCursor()
                    newobj = Tiles(filename, parent=self, linked_data_dir=dirname, mapviewer=self)
                    del wait

                    if which.lower() == 'tilescomp':
                        newobj.linked_data_dir_comp = dirname_comp

            self.mytiles.append(newobj)
            self.active_tile = newobj
            self.menu_tiles()

        elif which.lower() == 'array_xyz':

            curdict = self.myarrays
            curtree = self.myitemsarray

            msg = wx.MessageDialog(self, _('Do you want to crop the data?'), style=wx.YES_NO | wx.YES_DEFAULT)
            ret = msg.ShowModal()
            msg.Destroy()

            if ret == wx.ID_YES:

                newcrop = CropDialog(None)

                badvalues = True
                while badvalues:
                    badvalues = False

                    ret = newcrop.ShowModal()
                    if ret == wx.ID_CANCEL:
                        newcrop.Destroy()
                        return -1
                    else:
                        cropini = [[float(newcrop.ox.Value), float(newcrop.ex.Value)],
                                   [float(newcrop.oy.Value), float(newcrop.ey.Value)]]
                        tmpdx = float(newcrop.dx.Value)
                        tmpdy = float(newcrop.dy.Value)

                newcrop.Destroy()

                myxyz = xyz_scandir(filename, cropini)

                myhead = newcrop.get_header()
                # if min(myhead.dx, myhead.dy) != 1.:
                #     myhead.nbx = int(myhead.nbx * myhead.dx)
                #     myhead.nby = int(myhead.nby * myhead.dy)
                #     myhead.dx = 1.
                #     myhead.dy = 1.

            else:
                myxyz = xyz_scandir(filename, None)
                myhead = header_wolf()

                myhead.origx = np.min(myxyz[:, 0]) - .5
                myhead.origy = np.min(myxyz[:, 1]) - .5

                myhead.dx = 1.
                myhead.dy = 1.
                tmpdx = 1.
                tmpdy = 1.

                myhead.nbx = int(np.max(myxyz[:, 0]) - myhead.origx) + 1
                myhead.nby = int(np.max(myxyz[:, 1]) - myhead.origy) + 1

            if len(myxyz) == 0:
                return -1

            newobj = WolfArray()

            newobj.init_from_header(myhead)
            newobj.nullvalue = -99999.
            newobj.array.data[:, :] = -99999.

            newobj.fillin_from_xyz(myxyz)

            newobj.mask_data(newobj.nullvalue)

            # if min(tmpdx, tmpdy) != 1.:
            #     newobj.rebin(min(tmpdx, tmpdy))
            #     newobj.mask_data(newobj.nullvalue)

            newobj.change_gui(self)
            newobj.updatepalette(0)
            self.myarrays.append(newobj)
            self.active_array = newobj
            self._set_active_bc()

        elif which.lower() == 'array_lidar_first' or which.lower() == 'array_lidar_second':

            curdict = self.myarrays
            curtree = self.myitemsarray

            newcrop = CropDialog(None)

            badvalues = True
            while badvalues:
                badvalues = False

                ret = newcrop.ShowModal()
                if ret == wx.ID_CANCEL:
                    newcrop.Destroy()
                    return -1
                else:
                    cropini = [[float(newcrop.ox.Value), float(newcrop.ex.Value)],
                               [float(newcrop.oy.Value), float(newcrop.ey.Value)]]
                    tmpdx = float(newcrop.dx.Value)
                    tmpdy = float(newcrop.dy.Value)

            newcrop.Destroy()

            first, sec = Lidar2002.lidar_scandir(filename, cropini)

            if which.lower() == 'array_lidar_first':
                if len(first) == 0:
                    return -1

                newobj = Lidar2002.create_wolfarray(first, bounds=cropini)

                if min(tmpdx, tmpdy) != 1.:
                    newobj.rebin(min(tmpdx, tmpdy))

                newobj.change_gui(self)
                newobj.updatepalette(0)
                self.myarrays.append(newobj)
                self.active_array = newobj
                self._set_active_bc()

                id = 'lidar2002_firstecho'
            else:
                if len(sec) == 0:
                    return -1
                newobj = Lidar2002.create_wolfarray(sec, bounds=cropini)
                if min(tmpdx, tmpdy) != 1.:
                    newobj.rebin(min(tmpdx, tmpdy))

                newobj.change_gui(self)
                newobj.updatepalette(0)
                self.myarrays.append(newobj)
                self.active_array = newobj
                self._set_active_bc()
                id = 'lidar2002_secondecho'

        elif which.lower() == 'res2d':

            curdict = self.myres2D
            curtree = self.myitemsres2d

            if newobj is None:
                with wx.lib.busy.BusyInfo(_('Importing 2D model')):
                    wait = wx.BusyCursor()
                    newobj = Wolfresults_2D(filename, mapviewer=self)
                    del wait

            newobj.get_nbresults(True)
            newobj.updatepalette()
            self.myres2D.append(newobj)
            self.active_res2d = newobj

        elif which.lower() == 'res2d_gpu':

            curdict = self.myres2D
            curtree = self.myitemsres2d

            if newobj is None:
                with wx.lib.busy.BusyInfo(_('Importing 2D GPU model')):
                    wait = wx.BusyCursor()
                    newobj = wolfres2DGPU(filename, mapviewer=self)
                    del wait

            newobj.get_nbresults(True)
            newobj.read_oneresult(-1)
            newobj.updatepalette()
            self.myres2D.append(newobj)
            self.active_res2d = newobj

        elif which.lower() == 'vector':
            curdict = self.myvectors
            curtree = self.myitemsvector
            if newobj is None:
                with wx.lib.busy.BusyInfo(_('Importing file')):
                    wait = wx.BusyCursor()
                    newobj = Zones(filename, parent=self)
                    del wait
            self.myvectors.append(newobj)

        elif which.lower() == 'cross_sections':

            curdict = self.myvectors
            curtree = self.myitemsvector

            if newobj is None:

                dlg = wx.MessageDialog(None, 'Load LAZ data?', style=wx.YES_NO | wx.NO_DEFAULT)
                ret = dlg.ShowModal()
                dlg.Destroy()
                dirlaz = ''

                if ret == wx.ID_YES:
                    if self.mylazgrid is not None:
                        dlg = wx.MessageDialog(None, 'Gridded LAZ data exist - use them ?', style=wx.YES_NO | wx.YES_DEFAULT)
                        ret = dlg.ShowModal()
                        dlg.Destroy()

                        if ret == wx.ID_YES:
                            dirlaz = self.mylazgrid
                        else:
                            dlg = wx.DirDialog(None, 'If exist, where are the LAZ data?')
                            ret = dlg.ShowModal()
                            if ret == wx.ID_OK:
                                dirlaz = dlg.GetPath()
                    else:
                        dlg = wx.DirDialog(None, 'If exist, where are the LAZ data?')
                        ret = dlg.ShowModal()
                        if ret == wx.ID_OK:
                            dirlaz = dlg.GetPath()

                with wx.lib.busy.BusyInfo(_('Importing cross sections')):
                    wait = wx.BusyCursor()
                    if curfilter == 1:  # txt 2022
                        newobj = crosssections(filename, format='2022', dirlaz=dirlaz, mapviewer=self)
                    if curfilter == 0:  # vecz
                        newobj = crosssections(filename, format='vecz', dirlaz=dirlaz, mapviewer=self)
                    elif curfilter == 2:  # sxy
                        newobj = crosssections(filename, format='sxy', dirlaz=dirlaz, mapviewer=self)
                    else:  # txt 2000
                        newobj = crosssections(filename, format='2000', dirlaz=dirlaz, mapviewer=self)
                    del wait
            self.myvectors.append(newobj)
            newobj.mapviewer = self

        elif which.lower() == 'cloud':

            curdict = self.myclouds
            curtree = self.myitemscloud
            if newobj is None:

                loadhead = False
                if not filename.endswith('.dxf') and not filename.endswith('.shp'):
                    with open(filename,'r') as f:
                        text=f.read().splitlines()
                        tmphead=''
                        for i in range(4):
                            tmphead += text[i].replace('\t','\\t') +'\n'

                    dlg = wx.MessageDialog(None,_('Is there a file header (one upper line containing column names)?') + '\n\n' + tmphead,style=wx.YES_NO|wx.NO_DEFAULT)
                    ret=dlg.ShowModal()

                    if ret == wx.ID_YES:
                        loadhead = True

                with wx.lib.busy.BusyInfo(_('Importing cloud points')):
                    wait = wx.BusyCursor()
                    newobj = cloud_vertices(filename, header=loadhead, mapviewer=self)
                    del wait

            self.myclouds.append(newobj)
            self.active_cloud = newobj

            self.create_cloud_menu()

        elif which.lower() == 'triangulation':

            curdict = self.mytri
            curtree = self.myitemstri
            if newobj is None:
                with wx.lib.busy.BusyInfo(_('Importing triangulation')):
                    wait = wx.BusyCursor()
                    newobj = Triangulation(filename, mapviewer=self)
                    del wait

            self.mytri.append(newobj)
            self.active_tri = newobj

        elif which.lower() == 'other':

            if not newobj is None:
                curdict = self.myothers
                curtree = self.myitemsothers
                self.myothers.append(newobj)
                newobj.mapviewer = self
            else:
                logging.warning('No object to add in "Other" category -- Please provide an object to add or check your code')

        elif which.lower() == 'views':

            if newobj is None:
                newobj = WolfViews(plotted=ToCheck, mapviewer=self)
                newobj.read_from_file(filename)

            curdict = self.myviews
            curtree = self.myitemsviews
            self.myviews.append(newobj)

        elif which.lower() == 'wmsback':

            if not newobj is None:
                curdict = self.mywmsback
                curtree = self.myitemswmsback
                self.mywmsback.append(newobj)
            else:
                logging.warning('No object to add in "WMS background" category -- Please provide an object to add or check your code')

        elif which.lower() == 'wmsfore':

            if not newobj is None:
                curdict = self.mywmsfore
                curtree = self.myitemswmsfore
                self.mywmsfore.append(newobj)
            else:
                logging.warning('No object to add in "WMS foreground" category -- Please provide an object to add or check your code')

        elif which.lower() == 'particlesystem':

            curdict = self.mypartsystems
            curtree = self.myitemsps
            if newobj is None:
                    newobj = Particle_system(mapviewer=self)
                    newobj.load(filename)

            self.mypartsystems.append(newobj)
            self.active_particle_system = newobj

        if id == '':
            dlg = wx.TextEntryDialog(self, 'ID ? (case insensitive)', 'Choose an identifier', '')
            if filename != '':
                dlg.SetValue((Path(filename).with_suffix('')).name)
            else:
                dlg.SetValue('')

            if len(curdict) == 0:
                if dlg.ShowModal() == wx.ID_OK:
                    id = dlg.GetValue()
                    if id =='':
                        id = '001'
            else:
                ids = [cur.idx for cur in curdict]
                while id.lower() in ids:
                    if dlg.ShowModal() == wx.ID_OK:
                        id = dlg.GetValue()
                        if id =='':
                            id = '001'
            dlg.Destroy()

        ids = [cur.idx for cur in curdict]
        if id.lower() in ids:
            # endid = '_'
            endid = 1
            while (id + str(endid).zfill(3)).lower() in ids:
                endid += 1

        newobj.idx = id.lower()

        if curtree is not None:
            myitem = self.treelist.AppendItem(curtree, id, data=newobj)

            if ToCheck:
                self.treelist.CheckItem(myitem)
                self.treelist.CheckItem(self.treelist.GetItemParent(myitem))

                newobj.check_plot()
        else:
            logging.info(f'No tree item for this object {newobj.idx}')

        # curdict[id.lower()] = newobj
        if filename != '':
            newobj._filename_vector = filename.lower() # FIXME useful ??
        newobj.checked = ToCheck

        if isinstance(newobj,crosssections):
            self.add_object('cloud',newobj=newobj.cloud    ,id=newobj.idx+'_intersect',ToCheck=False)
            self.add_object('cloud',newobj=newobj.cloud_all,id=newobj.idx+'_all',      ToCheck=False)

        elif type(newobj) == WolfArray:
            if self.active_cs is None:
                self.active_cs = self.get_cross_sections()

        return 0

    def get_obj_from_treeitem(self, treeitem):
        """ Find the object associated with treeitem """

        return self.treelist.GetItemData(treeitem)

    def getobj_from_id(self, id: str):
        """ Find the object associated with id """

        for curdict in draw_type:
            keys = self.get_list_keys(curdict, checked_state=None)
            if id.lower() in keys:
                try:
                    idx = keys.index(id.lower())
                    return self.get_list_objects(curdict, checked_state=None)[idx]
                except:
                    return None

    def get_obj_from_id(self, id: str, drawtype: draw_type):
        """ Find the object associated with id in a specifid drawtype

        If you want to search in all drawtypes, use getobj_from_id instead.

        :param id: str : id of the object
        :param drawtype: draw_type : type of object to search

        """

        keys = self.get_list_keys(drawtype, checked_state=None)
        if id.lower() in keys:
            try:
                idx = keys.index(id.lower())
                return self.get_list_objects(drawtype, checked_state=None)[idx]
            except:
                return None

    def _get_list(self, drawing_type:draw_type):
        """ return the list of objects of type drawing_type """

        # ARRAYS = 'arrays'
        # BRIDGES= 'bridges'
        # WEIRS = 'weirs'
        # VECTORS = 'vectors'
        # CLOUD = 'clouds'
        # TRIANGULATION = 'triangulations'
        # PARTICLE_SYSTEM = 'particle systems'
        # CROSS_SECTIONS = 'cross_sections'
        # OTHER = 'others'
        # VIEWS = 'views'
        # RES2D = 'wolf2d'
        # WMSBACK = 'wms-background'
        # WMSFORE = 'wms-foreground'

        if drawing_type == draw_type.ARRAYS:
            return self.myarrays
        elif drawing_type == draw_type.VECTORS or drawing_type == draw_type.BRIDGES or drawing_type == draw_type.WEIRS or drawing_type == draw_type.CROSS_SECTIONS :
            return self.myvectors
        elif drawing_type == draw_type.TILES:
            return self.mytiles
        elif drawing_type == draw_type.CLOUD:
            return self.myclouds
        elif drawing_type == draw_type.TRIANGULATION:
            return self.mytri
        elif drawing_type == draw_type.RES2D:
            return self.myres2D
        elif drawing_type == draw_type.PARTICLE_SYSTEM:
            return self.mypartsystems
        elif drawing_type == draw_type.OTHER:
            return self.myothers
        elif drawing_type == draw_type.VIEWS:
            return self.myviews
        elif drawing_type == draw_type.WMSBACK:
            return self.mywmsback
        elif drawing_type == draw_type.WMSFORE:
            return self.mywmsfore
        else:
            logging.error('Unknown drawing type : ' + drawing_type)
            return None


    def get_list_keys(self, drawing_type:draw_type, checked_state:bool=True):
        """ Create a list of keys of type draw_type """

        if checked_state is None:
            return [curobj.idx for curobj in self._get_list(drawing_type)]
        else:
            return [curobj.idx for curobj in self._get_list(drawing_type) if curobj.plotted == checked_state]

    def get_list_objects(self, drawing_type:draw_type, checked_state:bool=True):
        """ Create a list of objects of type draw_type """

        if checked_state is None:
            return [curobj for curobj in self._get_list(drawing_type)]
        else:
            return [curobj for curobj in self._get_list(drawing_type) if curobj.plotted == checked_state]

    def single_choice_key(self, draw_type:draw_type, checked_state:bool=True, message:str=_('Make a choice'), title:str=_('Choice')):
        """ Create wx dialog to choose a key object of type draw_type """

        keys = self.get_list_keys(draw_type, checked_state)
        dlg = wx.SingleChoiceDialog(None, message, title, keys, style=wx.CHOICEDLG_STYLE)
        ret = dlg.ShowModal()

        if ret != wx.ID_OK:
            dlg.Destroy()
            return None

        idx = dlg.GetSelection()
        dlg.Destroy()

        return keys[idx]

    def single_choice_object(self, draw_type:draw_type, checked_state:bool=True, message:str=_('Make a choice'), title:str=_('Choice')):
        """ Create wx dialog to choose an object of type draw_type """

        keys = self.get_list_keys(draw_type, checked_state)
        obj = self.get_list_objects
        dlg = wx.SingleChoiceDialog(None, message, title, keys, style=wx.CHOICEDLG_STYLE)
        ret = dlg.ShowModal()

        if ret != wx.ID_OK:
            dlg.Destroy()
            return None

        idx = dlg.GetSelection()
        dlg.Destroy()

        return obj[idx]

    def multiple_choice_key(self, draw_type:draw_type, checked_state:bool=True, message:str=_('Make a choice'), title:str=_('Choice')):
        """ Create wx dialog to choose multiple keys object of type draw_type """

        keys = self.get_list_keys(draw_type, checked_state)
        dlg = wx.MultiChoiceDialog(None, message, title, keys, style=wx.CHOICEDLG_STYLE)
        ret = dlg.ShowModal()

        if ret != wx.ID_OK:
            dlg.Destroy()
            return None

        idx = dlg.GetSelections()
        dlg.Destroy()

        return [keys[i] for i in idx]

    def multiple_choice_object(self, draw_type:draw_type, checked_state:bool=True, message:str=_('Make a choice'), title:str=_('Choice')):
        """ Create wx dialog to choose multiple objects of type draw_type """

        keys = self.get_list_keys(draw_type, checked_state)
        obj = self.get_list_objects
        dlg = wx.MultiChoiceDialog(None, message, title, keys, style=wx.CHOICEDLG_STYLE)
        ret = dlg.ShowModal()

        if ret != wx.ID_OK:
            dlg.Destroy()
            return None

        idx = dlg.GetSelections()
        dlg.Destroy()

        return [obj[i] for i in idx]

    def iterator_over_objects(self, drawing_type:draw_type, checked_state:bool=True):
        """ Create iterator over objects of type draw_type """

        for obj in self.get_list_objects(drawing_type, checked_state):
            yield obj

    def gettreeitem(self, obj):

        """ Find the tree item associated with obj """

        up = self.treelist.GetFirstItem()
        updata = self.treelist.GetItemData(up)

        while updata is not obj:
            up = self.treelist.GetNextItem(up)
            updata = self.treelist.GetItemData(up)

        return up

    def removeobj(self):

        """Remove selected item from general tree"""

        if self.selected_treeitem is None:
            return

        id = self.treelist.GetItemText(self.selected_treeitem).lower()

        self.removeobj_from_id(id)

    def removeobj_from_id(self, id:str):

        """ Remove object from id """

        myobj = self.getobj_from_id(id)
        if myobj is not None:
            self.treelist.DeleteItem(self.gettreeitem(myobj))

            for curlist in self.all_lists:
                if myobj in curlist:
                    curlist.pop(curlist.index(myobj))

            myobj.hide_properties()

    def upobj(self):

        """Up selected item into general tree"""

        if self.selected_treeitem is None:
            return

        id:str
        id = self.treelist.GetItemText(self.selected_treeitem).lower()
        myobj = self.getobj_from_id(id)
        ischecked = self.treelist.GetCheckedState(self.selected_treeitem)

        assert self.selected_object is myobj, 'selected_object is not myobj'

        if myobj is not None:

            down = self.treelist.GetNextItem(self.selected_treeitem)
            up = self.treelist.GetFirstItem()
            up2 = up

            while self.treelist.GetNextItem(up) != self.selected_treeitem:
                up2= up
                up = self.treelist.GetNextItem(up)

            parent = self.treelist.GetItemParent(self.selected_treeitem)
            parentup = self.treelist.GetItemParent(up)
            parentup2 = self.treelist.GetItemParent(up2)

            if parent == parentup2:
                # up n'est pas le premier élément de la liste
                myitem = self.treelist.InsertItem(parent,up2,id,data=myobj)
                self.treelist.CheckItem(myitem,ischecked)
            elif parentup == parent:
                # up est le premier élément de la liste
                myitem = self.treelist.PrependItem(parent,id,data=myobj)
                self.treelist.CheckItem(myitem,ischecked)
            else:
                # nothing to do
                return

            self.treelist.DeleteItem(self.selected_treeitem)
            self.selected_treeitem = myitem

            # mouvement dans les listes pour garder l'ordre identique à l'arbre
            for curlist in self.all_lists:
                if myobj in curlist:
                    idx = curlist.index(myobj)
                    if idx>0:
                        curlist.pop(idx)
                        curlist.insert(idx-1,myobj)

    def downobj(self):

        """Down selected item into general tree"""

        if self.selected_treeitem is None:
            return

        id = self.treelist.GetItemText(self.selected_treeitem).lower()
        myobj = self.getobj_from_id(id)
        ischecked = self.treelist.GetCheckedState(self.selected_treeitem)

        if myobj is not None:

            down = self.treelist.GetNextItem(self.selected_treeitem)
            down2 = self.treelist.GetNextItem(down)

            parent = self.treelist.GetItemParent(self.selected_treeitem)
            parentdown = self.treelist.GetItemParent(down)
            parentdown2 = self.treelist.GetItemParent(down2)

            if parent == parentdown:
                # on n'est pas sur le dernoier élément
                myitem = self.treelist.InsertItem(parent,down,id,data=myobj)
                self.treelist.CheckItem(myitem,ischecked)
            else:
                # nothing to do
                return

            self.treelist.DeleteItem(self.selected_treeitem)
            self.selected_treeitem = myitem

            for curlist in self.all_lists:
                if myobj in curlist:
                    idx = curlist.index(myobj)
                    if idx==len(curlist)-2:
                        curlist.append(myobj)
                        curlist.pop(idx)
                    elif idx<len(curlist)-1:
                        curlist.insert(idx+1,myobj)
                        curlist.pop(idx)

    def OnShowPopup(self, event):
        pos = event.GetPosition()
        if pos == (-1, -1):
            width, height = self.GetSize()
            pos = (width / 2., height / 2.)
        # else:
        #     pos = pos - self.GetPosition()

        self.PopupMenu(self.popupmenu, pos)

    def OnPopupItemSelected(self, event):
        item = self.popupmenu.FindItemById(event.GetId())
        text = item.ItemLabel

        if text == _('Save'):
            if self.selected_object is not None:
                if issubclass(type(self.selected_object), WolfArray):
                    self.selected_object.write_all()
                elif type(self.selected_object) is Zones:
                    self.selected_object.saveas()
                elif type(self.selected_object) is Triangulation:
                    self.selected_object.saveas()
                elif isinstance(self.selected_object, Particle_system):
                    self.selected_object.save()
        elif text==_('Up'):
            self.upobj()
        elif text == _('Down'):
            self.downobj()
        elif text == _('Rename'):

            #Modification du nom de l'objet sélectionné
            if self.selected_object is not None:
                #récupération de l'id courant
                label = self.selected_object.idx
                dlg = wx.TextEntryDialog(self, message=_('Chose a new label :'), value=label)
                ret=dlg.ShowModal()
                newlab = dlg.GetValue()
                dlg.Destroy()

                #MAJ de l'id dans l'objet
                self.selected_object.idx = newlab
                #MAJ de l'arbre
                self.treelist.SetItemText(self.selected_treeitem, newlab)

        elif text ==  _('Duplicate'):

            # Duplication de l'objet sélectionné
            if self.selected_object is not None:
                #récupération de l'id courant
                label = self.selected_object.idx + '_copy'
                dlg = wx.TextEntryDialog(self, message=_('Chose a label for the copy:'), value=label)
                ret=dlg.ShowModal()
                newlab = dlg.GetValue()
                dlg.Destroy()

                if isinstance(self.selected_object, WolfArray) and (not type(self.selected_object) in [WolfArrayMB, WolfArrayMNAP]):

                    curtype = self.selected_object.dtype

                    if curtype == np.float64:
                        curtype = 'float64'
                    elif curtype == np.float32:
                        curtype = 'float32'
                    elif curtype == np.int32:
                        curtype = 'int32'
                    elif curtype == np.int16:
                        curtype = 'int16'
                    elif curtype == np.int8:
                        curtype = 'int8'

                    dlg = wx.MessageDialog(None, _('The type of the data is {}.\nDo you want to change this type?'.format(curtype)), style=wx.YES_NO | wx.NO_DEFAULT)
                    ret = dlg.ShowModal()
                    dlg.Destroy()

                    if ret == wx.ID_YES:
                        dlg = wx.SingleChoiceDialog(None, _('Choose a type'), _('Type'), ['float32','float64','int32','int16','int8'], style=wx.CHOICEDLG_STYLE)
                        ret = dlg.ShowModal()

                        if ret != wx.ID_OK:
                            dlg.Destroy()
                            return

                        idx = dlg.GetSelection()
                        dlg.Destroy()

                        if idx == 0:
                            curtype = WOLF_ARRAY_FULL_SINGLE # np.float32
                        elif idx == 1:
                            curtype = WOLF_ARRAY_FULL_DOUBLE # np.float64
                        elif idx == 2:
                            curtype = WOLF_ARRAY_FULL_INTEGER # np.int32
                        elif idx == 3:
                            curtype = WOLF_ARRAY_FULL_INTEGER16 #np.int16
                        elif idx == 4:
                            curtype = WOLF_ARRAY_FULL_INTEGER8 #np.int8

                        newarray = WolfArray(srcheader=self.selected_object.get_header(), whichtype=curtype, nullvalue=self.selected_object.nullvalue)

                        newarray.allocate_ressources()
                        asnewtype = self.selected_object.array.data.astype(newarray.dtype)
                        newarray.array.data[:,:] = asnewtype[:,:]
                        newarray.copy_mask(self.selected_object, forcenullvalue=True, link=False)
                    else:
                        newarray = WolfArray(mold=self.selected_object)

                    self.add_object('array', newobj=newarray, id=newlab)
                    self.Refresh()
                else:
                    logging.warning(_('Not yet implemented'))

        elif text == _('Save as'):
            # save objet to file, choosing the file name

            if self.selected_object is not None:
                if issubclass(type(self.selected_object), WolfArray):
                    filterArray = "bin (*.bin)|*.bin|Geotif (*.tif)|*.tif|Numpy (*.npy)|*.npy|all (*.*)|*.*"
                    fdlg = wx.FileDialog(self, "Choose file name for Array : " + self.selected_object.idx, wildcard=filterArray,
                                         style=wx.FD_SAVE)
                    ret = fdlg.ShowModal()
                    if ret == wx.ID_OK:
                        curfil = fdlg.GetFilterIndex()

                        self.selected_object.filename = fdlg.GetPath()

                        self.selected_object.write_all()
                    fdlg.Destroy()
                elif type(self.selected_object) is Zones:
                    filterArray = "vec (*.vec)|*.vec|vecz (*.vecz)|*.vecz|Shapefile (*.shp)|*.shp|all (*.*)|*.*"
                    fdlg = wx.FileDialog(self, "Choose file name for Vector :" + self.selected_object.idx, wildcard=filterArray,
                                         style=wx.FD_SAVE)
                    ret = fdlg.ShowModal()
                    if ret == wx.ID_OK:
                        self.selected_object.saveas(fdlg.GetPath())
                    fdlg.Destroy()
                elif type(self.selected_object) is Triangulation:
                    filterArray = "tri (*.tri)|*.tri|all (*.*)|*.*"
                    fdlg = wx.FileDialog(self, "Choose file name for triangulation :" + self.selected_object.idx, wildcard=filterArray,
                                         style=wx.FD_SAVE)
                    ret = fdlg.ShowModal()
                    if ret == wx.ID_OK:
                        self.selected_object.saveas(fdlg.GetPath())

                    fdlg.Destroy()
                elif isinstance(self.selected_object, Particle_system):
                    filterArray = "json (*.json)|*.json|all (*.*)|*.*"
                    fdlg = wx.FileDialog(self, "Choose file name for particle system :" + self.selected_object.idx, wildcard=filterArray,
                                         style=wx.FD_SAVE)
                    ret = fdlg.ShowModal()
                    if ret == wx.ID_OK:
                        self.selected_object.save(fdlg.GetPath())

                    fdlg.Destroy()

        elif text == _('Properties'):

            myobj = self.selected_object
            if type(myobj) in [WolfArray, WolfArrayMB, WolfArrayMNAP, Zones, Wolfresults_2D, wolfres2DGPU, Particle_system, Picc_data, Cadaster_data, hydrometry_wolfgui]:
                myobj.show_properties()

        elif text == _('Boundary conditions'):
            bc = self.get_boundary_manager(self.selected_object)
            if bc is not None:
                bc.Show()

        elif _('Convert to mono-block') in text:

            if isinstance(self.selected_object, WolfArrayMB):
                mono = self.selected_object.as_WolfArray()
                self.add_object('array', newobj=mono, id=self.selected_object.idx + '_mono')
                logging.info(_('Mono-block created and added to the viewer'))

            elif isinstance(self.selected_object, Wolfresults_2D):
                mono = self.selected_object.as_WolfArray()

                if isinstance(mono, WolfArrayMB):
                    mono = mono.as_WolfArray()

                self.add_object('array', newobj=mono, id=self.selected_object.idx + '_mono')
                logging.info(_('Mono-block created and added to the viewer'))

            else:
                logging.warning(_('Convert to mono-blocks not yet implemented for this type of object'))

        elif _('Convert to multi-blocks') in text:

            if isinstance(self.selected_object, Wolfresults_2D):
                mb = self.selected_object.as_WolfArray(force_mb=True)

                if isinstance(mb, WolfArrayMB):
                    logging.info(_('Multi-blocks created and added to the viewer'))

                elif isinstance(mb, WolfArray):
                    logging.warning(_('Mono-blocks created and added to the viewer -- Instead of multi-blocks as only one block was found'))

                self.add_object('array', newobj=mb, id=self.selected_object.idx + '_mb')
            else:
                logging.warning(_('Convert to multi-blocks not yet implemented for this type of object'))

        elif _('Export to Shape file') in text:

            if isinstance(self.selected_object, Zones):
                filterArray = "Shapefile (*.shp)|*.shp"
                fdlg = wx.FileDialog(self, "Choose file name for Zones :" + self.selected_object.idx, wildcard=filterArray,
                                        style=wx.FD_SAVE)
                ret = fdlg.ShowModal()
                if ret == wx.ID_OK:
                    self.selected_object.export_to_shapefile(fdlg.GetPath())
                fdlg.Destroy()

        elif _('Export active zone to Shape file') in text:

            if isinstance(self.selected_object, Zones):

                filterArray = "Shapefile (*.shp)|*.shp"
                fdlg = wx.FileDialog(self, "Choose file name for Vector :" + self.selected_object.idx, wildcard=filterArray,
                                        style=wx.FD_SAVE)
                ret = fdlg.ShowModal()
                if ret == wx.ID_OK:
                    self.selected_object.export_active_zone_to_shapefile(fdlg.GetPath())
                fdlg.Destroy()


    def OnClose(self, event):
        nb = 0
        if self.linked:
            if self.linkedList is not None:
                if self in self.linkedList:
                    id = self.linkedList.index(self)
                    self.linkedList.pop(id)
                    nb = len(self.linkedList)

        if nb == 0:
            if self.wxlogging is not None:
                dlg = wx.MessageDialog(None, _('Do you want to quit Wolf ?'), _('Quit Wolf'), wx.YES_NO | wx.NO_DEFAULT)
                ret = dlg.ShowModal()
                dlg.Destroy()
                if ret == wx.ID_YES:
                    self.Destroy()
                    #FIXME : It is not a really proper way to quit the application
                    wx.Exit()
                    return
                else:
                    return
        self.Destroy()

    def onselectitem(self,event):
        myitem = event.GetItem()

        self.selected_treeitem = myitem

        nameitem = self.treelist.GetItemText(myitem).lower()
        curobj = self.getobj_from_id(nameitem)

        if isinstance(curobj, PlansTerrier):
            self.active_landmap = curobj

        elif isinstance(curobj, Zones):
            self.active_zones = curobj


        self.treelist.SetToolTip(self.treelist.GetItemText(myitem))

    def OnCheckItem(self, event:TreeListEvent):

        myitem = event.GetItem()
        myparent = self.treelist.GetItemParent(myitem)
        check = self.treelist.GetCheckedState(myitem)
        nameparent = self.treelist.GetItemText(myparent).lower()
        nameitem = self.treelist.GetItemText(myitem).lower()

        ctrl = wx.GetKeyState(wx.WXK_CONTROL)

        # ctrl = event.ControlDown()

        if nameparent != '':
            curobj = self.getobj_from_id(nameitem)
            if curobj is None:
                return

            if bool(check):
                try:
                    curobj.check_plot()

                    if isinstance(curobj, PlansTerrier):
                        if curobj.initialized:
                            self.menu_landmaps()
                            logging.info(_('Landmap initialized'))
                        else:
                            logging.warning(_('Landmap not initialized'))

                except Exception as ex:
                    wx.LogMessage(str(ex))
                    wx.MessageBox(str(ex), _("Error"), wx.ICON_ERROR)
            else:
                if issubclass(type(curobj), WolfArray):
                    curobj.uncheck_plot(not ctrl,ctrl)
                else:
                    curobj.uncheck_plot()

            # if nameparent == 'vectors' or nameparent == 'cross_sections':
            #     if wx.GetKeyState(wx.WXK_CONTROL):
            #         curobj.showstructure(self)

            if curobj.idx == 'grid' and check:
                dlg = wx.TextEntryDialog(self, 'Size of the Grid ? (float)', 'Choose an size')
                dlg.SetValue('1000.')
                size = 1000.
                if dlg.ShowModal() == wx.ID_OK:
                    size = float(dlg.GetValue())
                curobj.creategrid(size, self.xmin, self.ymin, self.xmax, self.ymax)

            self.Refresh()

    def getXY(self, pospix):

        width, height = self.canvas.GetSize()
        X = float(pospix[0]) / self.sx + self.xmin
        Y = float(height - pospix[1]) / self.sy + self.ymin
        return X, Y

    def OnZoomGesture(self, e):
        pass

    def OnLeave(self, e):
        if e.ControlDown():
            self.mytooltip.Show(False)

    def get_cross_sections(self):
        """
        Récupération du premier objet crosssections disponible
        """
        for obj in self.iterator_over_objects(draw_type.VECTORS):
            if isinstance(obj,crosssections):
                return obj

        return None

    def set_active_profile(self, active_profile: profile):
        """
        This method sets the active profile in Pydraw (useful for interfaces communication).
        """
        self.active_profile = active_profile

    def set_active_vector(self, active_vector: vector):
        """
        This method sets the active vector in Pydraw (useful for interfaces communication).
        """
        self.active_vector = active_vector

    def get_active_profile(self):
        """
        This methods returns the  active profile in pydraw (useful for interfaces communication).
        """
        return self.active_profile

    def plot_cross(self, x:float, y:float):

        # Search for cross sections (List of profiles)
        if self.active_cs is None:
            self.active_cs = self.get_cross_sections()
            if self.active_cs is None:
                logging.warning(_('No cross sections available -- Please load a file or create data !'))
            return

        # Initialisation of the notebook where the active profile is plotted.
        if self.notebookprof is None:
            self.notebookprof = ProfileNotebook(mapviewer=self)
            # self.myfigprof = self.notebookprof.add('Figure 1', which= "all")
            self.myfigprof = self.notebookprof.add('Reference') # FIXME Updated add method

        else:
            try:
                self.notebookprof.Show()
            except:
                self.notebookprof = ProfileNotebook(mapviewer=self)
                # self.myfigprof = self.notebookprof.add('Figure 1', which= "all")
                self.myfigprof = self.notebookprof.add('Figure 1') # FIXME updated  add method

        # Initialisation of the active profile
        # 1. We uncolor the active profile in wolf GUI.
        self.active_profile: profile
        if self.active_profile is not None:
            self.active_profile.uncolor_active_profile()

        # 2. We select the closest profile corresponding to the user's right click in the GUI.
        self.active_profile = self.active_cs.select_profile(x, y)

        #Finally, we set the profile and the cross section (list of profiles) in the notebook.
        #FIXME Iden establishes the communications between pydraw and the notebook (to avoid circular information).
        self.myfigprof.cs_setter(mycross = self.active_cs, active_profile= self.active_profile, mapviewer = self)


    def OnRightDown(self, e: wx.MouseEvent):
        """
        Event when the right button of the mouse is pressed.

        We use this event to manage "action" set by others objects.

        """

        pos = e.GetPosition()
        x, y = self.getXY(pos)

        alt = e.AltDown()
        ctrl = e.ControlDown()
        shiftdown = e.ShiftDown()

        if self.action is None:
            if self.active_bc is not None:

                self.start_action('select bc', _('Select a boundary condition'))
                tmpvec = vector()

                self.last_active_vector = self.active_vector
                self.active_vector = tmpvec
                tmpvec.add_vertex(wolfvertex(x, y))

            self.rightdown = (x, y)

        elif 'pick landmap' in self.action:
            # Pick a landmap if loaded

            if self.active_landmap is None:
                logging.warning(_('No landmap available -- Please activate the data and retry !'))
                return

            if 'full' in self.action:
                self.active_landmap.load_texture(x,y, which='full')
            else:
                self.active_landmap.load_texture(x,y, which='low')
            self.Refresh()

        elif self.action == 'bridge gltf':
            # Create a bridge in gltf format

            self.bridgepar = (x, y)

            dlg = wx.TextEntryDialog(self, 'Z maximum ?', 'Choose an elevation as top')
            dlg.SetValue('')

            zmax = 0.
            if dlg.ShowModal() == wx.ID_OK:
                zmax = float(dlg.GetValue())
            dlg.Destroy()

            dlg = wx.FileDialog(None, _('Choose filename'),
                                wildcard='glb (*.glb)|*.glb|gltf2 (*.gltf)|*.gltf|All (*.*)|*.*', style=wx.FD_SAVE)
            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            fn = dlg.GetPath()
            dlg.Destroy()

            points, triangles = self.active_vector.triangulation_ponts(self.bridgepar[0], self.bridgepar[1], zmax)
            self.active_cs.export_gltf_gen(points, triangles, fn)

            self.start_action('', 'None')

        elif self.action == 'Plot cross section':
            # Plot cross section
            self.plot_cross(x, y)

        elif self.action == 'Set 1D profile':
            # Set 1D profile

            if self.active_cs is None:
                self.active_cs = self.get_cross_sections()

            if self.active_cs is None:
                logging.warning(_('No cross sections available -- Please load a file or create data !'))
                return

            if self.notebookcs is None:
                self.notebookcs = PlotNotebook()
                self.myfigcs = self.notebookcs.add(_("Cross section"), "CS")
            else:
                try:
                    self.notebookcs.Show()
                except:
                    self.notebookcs = PlotNotebook()
                    self.myfigcs = self.notebookcs.add(_("Cross section"), "CS")

            self.active_profile: profile

            #on met l'ancien profil actif en noir
            if self.active_profile is not None:
                self.active_profile.uncolor_active_profile()
            if self.myfigcs.mycs is not None:
                self.myfigcs.mycs.uncolor_active_profile()

            # self.active_profile = self.active_cs.select_profile(x, y)
            self.active_profile = self.frame_create1Dfrom2D.active_profile

            self.myfigcs.set_linked_arrays(self.get_linked_arrays())
            self.myfigcs.set_cs(self.active_profile)

            #on met le profil en rouge et plus épais
            self.active_profile.color_active_profile()
            self.zoom_on_active_profile()

            self.Paint()

        elif self.action == 'Select nearest profile':
            # Select nearest profile

            if self.active_cs is None:
                self.active_cs = self.get_cross_sections()

            if self.active_cs is None:
                logging.warning(_('No cross sections available -- Please load a file or create data !'))
                return

            if self.notebookcs is None:
                self.notebookcs = PlotNotebook()
                self.myfigcs = self.notebookcs.add(_("Cross section"), "CS")
            else:
                try:
                    self.notebookcs.Show()
                except:
                    self.notebookcs = PlotNotebook()
                    self.myfigcs = self.notebookcs.add(_("Cross section"), "CS")

            self.active_profile: profile

            #on met l'ancien profil actif en noir
            if self.active_profile is not None:
                self.active_profile.uncolor_active_profile()
            if self.myfigcs.mycs is not None:
                self.myfigcs.mycs.uncolor_active_profile()

            self.active_profile = self.active_cs.select_profile(x, y)

            self.myfigcs.set_linked_arrays(self.get_linked_arrays())
            self.myfigcs.set_cs(self.active_profile)

            #on met le profil en rouge et plus épais
            self.active_profile.color_active_profile()

            self.Refresh()

        elif self.action == 'select active tile':
            # Select active tile from Lidar data

            self.active_tile.select_vectors_from_point(x, y, True)
            self.active_vector = self.active_tile.get_selected_vectors()

            tilearray = self.active_tile.get_array(self.active_vector)
            if tilearray is not None:
                if self.active_vector.myname =='':
                    bbox = self.active_vector.get_bounds()
                    id_label = '{}-{}'.format(bbox[0][0], bbox[1][1])
                else:
                    id_label = self.active_vector.myname

                self.add_object('array', newobj = tilearray, ToCheck=True, id=id_label)

        elif self.action.find('select active vector') > -1:
            # Select active vector

            inside = self.action.find('inside') > -1 # only polygons/closed polyline if 'inside' is in the action name
            onlyonezone = self.action.find('2') > -1 # only the active zone if '2' is in the action name, all zones otherwise

            if onlyonezone:
                self.active_zone.select_vectors_from_point(x, y, inside)
                self.active_vector = self.active_zone.get_selected_vectors()[0]

                if self.active_vector is not None:
                    self.active_zone.parent.Activate_vector(self.active_vector)
                    self.active_zone.active_vector = self.active_vector
                    self.active_zones.active_zone = self.active_vector.parentzone
            else:
                self.active_zones.select_vectors_from_point(x, y, inside)
                self.active_vector = self.active_zones.get_selected_vectors()

            if self.active_vector is not None:
                self.active_zones.Activate_vector(self.active_vector)
                self.active_zone = self.active_vector.parentzone
                self.active_zones.expand_tree(self.active_zone)

        elif 'select node by node' in self.action:
            # Select node by node

            if 'results' in self.action:
                curobj:Wolfresults_2D
                curobj = self.active_res2d.SelectionData
            else:
                curobj: WolfArray
                curobj = self.active_array.SelectionData

            if curobj.myselection == 'all':
                logging.warning(_('All nodes are selected !!'))
                logging.warning(_('Selecting node by node will force to reset the selection'))
                logging.warning(_('and start from scratch'))

            curobj.add_node_to_selection(x, y)
            curobj.update_nb_nodes_selection()
            self.Paint()

        elif 'select by tmp vector' in self.action or 'select by vector' in self.action:
            # Select nodes by vector or temporary vector
            self.active_vector.add_vertex(wolfvertex(x, y))

        elif 'laz tmp vector' == self.action:
            self.active_vector.add_vertex(wolfvertex(x, y))
            self.active_vector.find_minmax()

        elif self.action == 'create polygon - tiles':
            self.active_vector.add_vertex(wolfvertex(x, y))
            self.active_vector.find_minmax()

        elif self.action == 'capture vertices':

            if ctrl:
                if self.active_array is not None:
                    z = self.active_array.get_value(x, y)
                    self.active_vector.myvertices[-1].z = z
                else:
                    logging.warning(_('No array available and ctrl is pressed -- Please load a file or create data !'))

            self.active_vector.add_vertex(wolfvertex(x, y))

            self.active_vector.find_minmax()
            self.active_zone.find_minmax()

        elif self.action == 'dynamic parallel':
            # Create a dynamic parallel line
            if ctrl:
                if self.active_array is not None:
                    z = self.active_array.get_value(x, y)
                    self.active_vector.myvertices[-1].z = z
                else:
                    logging.warning(_('No array available and ctrl is pressed -- Please load a file or create data !'))

            self.active_vector.add_vertex(wolfvertex(x, y))

            self.active_zone.parallel_active(self.dynapar_dist)

        elif self.action == 'modify vertices':

            if self.active_vertex is None:
                self.active_vertex = self.active_vector.find_nearest_vert(x, y)
            else:
                self.active_vertex.limit2bounds(self.active_vector.mylimits)

                if ctrl:
                    if self.active_array is not None:
                        z = self.active_array.get_value(x, y)
                        self.active_vertex.z = z
                    else:
                        logging.warning(_('No array available and ctrl is pressed -- Please load a file or create data !'))

                self.active_vertex = None

        elif self.action == 'insert vertices':

            if self.active_vertex is None:
                self.active_vertex = self.active_vector.insert_nearest_vert(x, y)
            else:

                if ctrl:
                    if self.active_array is not None:
                        z = self.active_array.get_value(x, y)
                        self.active_vertex.z = z
                    else:
                        logging.warning(_('No array available and ctrl is pressed -- Please load a file or create data !'))

                self.active_vertex = None

        else:
            self.rightdown = (x, y)


    def OnRightUp(self, e):
        pos = e.GetPosition()
        x, y = self.getXY(pos)

        if self.active_bc is not None:
            if self.action == 'select bc':
                try:
                    minx = min(self.rightdown[0], x)
                    miny = min(self.rightdown[1], y)
                    maxx = max(self.rightdown[0], x)
                    maxy = max(self.rightdown[1], y)

                    if minx != maxx and maxy != miny:
                        self.active_bc.ray_tracing_numpy([[minx, miny], [maxx, miny], [maxx, maxy], [minx, maxy]], 'X')
                        self.active_bc.ray_tracing_numpy([[minx, miny], [maxx, miny], [maxx, maxy], [minx, maxy]], 'Y')
                    else:
                        self.active_bc.query_kdtree((x, y))

                    self.active_bc.update_selection()
                    self.Refresh()

                    self.active_vector = self.last_active_vector

                    self.end_action(_('End selection BC'))

                except:
                    pass

    def OnButton(self, e: wx.MouseEvent):
        d = e.GetWheelDelta()
        r = e.GetWheelRotation()
        a = e.GetWheelAxis()

        altdown = e.AltDown()
        ctrldown = e.ControlDown()
        shiftdown = e.ShiftDown()
        spacedown = wx.GetKeyState(wx.WXK_SPACE)

        if self.action == 'dynamic parallel' and shiftdown and not ctrldown:
            self.dynapar_dist *= (1 - .1 * (r / max(d, 1)))
            self.dynapar_dist = max(self.dynapar_dist, .01)

            self.active_zone.parallel_active(self.dynapar_dist)
            self.Refresh()
            return
        elif self.action == 'dynamic parallel' and shiftdown and ctrldown:
            dlg = wx.NumberEntryDialog(None,
                                       _('What is the desired size [cm] ?'),
                                       'ds',
                                       'ds size',
                                       int(self.dynapar_dist * 100.),
                                       1,
                                       100000)
            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            self.dynapar_dist = float(dlg.GetValue()) / 100.
            self.dynapar_dist = max(self.dynapar_dist, .01)
            dlg.Destroy()

            self.active_zone.parallel_active(self.dynapar_dist)
            self.Refresh()
            return

        # Allow the user to zoom onto the pixel where the
        # mouse cursor is

        # Step1: move the map so that the pixem under the mouse cursor
        # ends up right in the middle of the screen (this move is
        # not visible from the end user point of view, it's just
        # here to make computation seasier)
        if spacedown:
            self.center_view_on( *self.getXY( e.GetPosition()))

        # Zoom/dezoom, center pf the tranfromation is the center of the screen
        self.width = self.width * (1 - .1 * (r / max(d, 1)))
        self.height = self.height * (1 - .1 * (r / max(d, 1)))

        if spacedown:
            self.updatescalefactors() # not base on mousex

            # Translate back the pixel at the center of the screen to where the
            # mouse cursor is. For that we measure the delta in screen coordinates
            # and transform it to map space coordinates.
            x_mid, y_mid = self.canvas.GetSize()
            x_mid, y_mid = self.getXY((0.5*x_mid, y_mid*0.5))
            x, y = self.getXY( e.GetPosition())
            dx, dy = x_mid - x, y_mid - y
            self.mousex +=  dx
            self.mousey +=  dy

        # will translate and rescale the map view so that it fits the window.
        self.setbounds()

    def OnRDClick(self, e):
        self._endactions()

    def OnLDClick(self, e:wx.MouseEvent):
        pos = e.GetPosition()

        ctrldown = e.ControlDown()

        x, y = self.getXY(pos)

        self.mousex = self.mousedown[0]
        self.mousey = self.mousedown[1]
        self.mousedown = (0., 0.)
        self.oneclick = False
        self.setbounds()

        if ctrldown:
            if self.active_viewer3d is not None:
                self.active_viewer3d.force_view(self.mousex, self.mousey, self.active_array.get_value(self.mousex, self.mousey))
                self.Refresh()

    def OnLDown(self, e):

        if not self.move:
            pos = e.GetPosition()
            x, y = self.getXY(pos)
            self.mousedown = (x, y)
            self.move=True

    def _set_active_bc(self):
        """Search and activate BCManager according to active_array"""
        if self.active_bc is not None:
            if self.active_array != self.active_bc.linked_array:
                # it is not the good one -> Hide
                self.active_bc.Hide()
            else:
                return
        # searching if bcmanager is attached to active_array
        self.active_bc = None
        for curbc in self.mybc:
            if self.active_array == curbc.linked_array:
                self.active_bc = curbc
                self.active_bc.Show()
                return

    def set_statusbar_text(self, txt:str):
        """ Set the status bar text """
        self.StatusBar.SetStatusText(txt)

    def set_label_selecteditem(self, nameitem:str):
        """ Set the label of the selected item in the tree list """
        self._lbl_selecteditem.SetLabel(nameitem)

    def OnActivateTreeElem(self, e): #:dataview.TreeListEvent ):
        """ Activate the selected item in the tree list """
        curzones: Zones
        curzone: zone
        curvect: vector

        myitem = e.GetItem()
        ctrl = wx.GetKeyState(wx.WXK_CONTROL)
        alt = wx.GetKeyState(wx.WXK_ALT)

        myparent = self.treelist.GetItemParent(myitem)
        check = self.treelist.GetCheckedState(myitem)


        nameparent = self.treelist.GetItemText(myparent).lower()
        nameitem = self.treelist.GetItemText(myitem).lower()

        myobj = self.treelist.GetItemData(myitem)
        self.selected_object = myobj
        self.set_label_selecteditem(nameitem)

        #FIXME : To generalize using draw_type
        if type(myobj) == Zones:
            if ctrl:
                myobj.show_properties()

        elif type(myobj) == hydrometry_wolfgui:
            if ctrl:
                myobj.show_properties()

        elif type(myobj) in [Picc_data, Cadaster_data]:
            if ctrl:
                myobj.show_properties()

        elif type(myobj) == Particle_system:
            if ctrl:
                myobj.show_properties()

        elif type(myobj) == Tiles:
            self.active_tile= myobj

        elif issubclass(type(myobj), WolfArray):
            if ctrl:
                myobj.show_properties()
                # myobj.myops.SetTitle(_('Operations on array: ')+myobj.idx)
                # myobj.myops.Show()

            logging.info(_('Activating array : ' + nameitem))
            self.active_array = myobj

            # If BC maneger is attached to the array, we activate it
            self._set_active_bc()

            #Print info in the status bar
            txt  = 'Dx : {:.4f} ; Dy : {:.4f}'.format(self.active_array.dx, self.active_array.dy)
            txt += ' ; Xmin : {:.4f} ; Ymin : {:.4f}'.format(self.active_array.origx, self.active_array.origy)
            txt += ' ; Xmax : {:.4f} ; Ymax : {:.4f}'.format(self.active_array.origx + self.active_array.dx * float(self.active_array.nbx),
                                                           self.active_array.origy + self.active_array.dy * float(self.active_array.nby))
            txt += ' ; Nx : {:d} ; Ny : {:d}'.format(self.active_array.nbx, self.active_array.nby)

            if self.active_array.nb_blocks > 0:
                txt += ' ; Nb blocks : {:d}'.format(self.active_array.nb_blocks)

            txt += ' ; Type : ' + self.active_array.dtype_str

            self.set_statusbar_text(txt)


        elif type(myobj) in [WolfViews]:
            logging.info(_('Activating view : ' + nameitem))
            self.active_view = myobj

        elif type(myobj) == cloud_vertices:
            self.active_cloud = myobj
            if ctrl:
                myobj.myprop.show()

        elif type(myobj) == crosssections:
            if ctrl:
                myobj.showstructure()
            logging.info(_('Activating cross sections : ' + nameitem))
            self.active_cs = myobj

        elif type(myobj) == Triangulation:
            self.active_tri = myobj

        elif type(myobj) == Wolfresults_2D:
            logging.info(_('Activating Wolf2d results : ' + nameitem))
            self.active_res2d = myobj

            if ctrl:
                myobj.show_properties()

            if alt:
                dlg = wx.MessageDialog(self,_('Do you want to open the 2D model?'),style=wx.YES_NO|wx.NO_DEFAULT)
                ret=dlg.ShowModal()
                dlg.Destroy()
                if ret == wx.ID_NO:
                    return

                from .PyGui import Wolf2DModel
                mywolf = Wolf2DModel(dir=os.path.dirname(self.active_res2d.filenamegen), splash=False)

        elif type(myobj) == wolfres2DGPU:
            logging.info(_('Activating Wolf2d results : ' + nameitem))
            self.active_res2d = myobj

            if ctrl:
                myobj.show_properties()

    def _update_mytooltip(self):
        """ Update the tooltip with the values of the active arrays and results at position x,y """

        x,y,pos = self._last_mouse_pos

        self.mytooltip.myparams.clear()


        curgroup = 'Position'
        self.mytooltip.myparams[curgroup] = {}

        curpar = _('Pixel (col,row)')
        self.mytooltip.add_param(groupname = curgroup,
                                 name = curpar,
                                 value = '(' + str(pos[0]) + ' ; ' + str(pos[1]) + ')',
                                 type = Type_Param.String,
                                 comment = '')

        curpar = _('Coordinate X [m]')
        self.mytooltip.add_param(groupname = curgroup,
                                 name = curpar,
                                 value = '{:3f}'.format(x),
                                 type = Type_Param.String,
                                 comment = '')

        curpar = _('Coordinate Y [m]')
        self.mytooltip.add_param(groupname = curgroup,
                                 name = curpar,
                                 value = '{:3f}'.format(y),
                                 type = Type_Param.String,
                                 comment = '')

        for locarray in self.myres2D:
            locarray:Wolfresults_2D
            curgroup = locarray.idx
            if locarray.checked:
                try:
                    vals,labs = locarray.get_values_labels(x,y)

                    i, j, curbloc = locarray.get_blockij_from_xy(x, y)

                    if i != '-':
                        curpar = 'Indices (i;j;bloc) (1-based)'

                        self.mytooltip.add_param(groupname = curgroup,
                                                 name = curpar,
                                                 value =  '(' + str(i) + ';' + str(j) + ';' + str(curbloc) + ')',
                                                 type = Type_Param.String,
                                                 comment = '')

                        for val,curpar in zip(vals,labs):

                            if val is np.nan:
                                self.mytooltip.add_param(groupname = curgroup,
                                                        name = 'Value',
                                                        value =  "Nan",
                                                        type = Type_Param.String,
                                                        comment = '')

                            elif np.ma.is_masked(val):

                                self.mytooltip.add_param(groupname = curgroup,
                                                        name = 'Value',
                                                        value =  "Masked",
                                                        type = Type_Param.String,
                                                        comment = '')

                            elif isinstance(val, str):
                                self.mytooltip.add_param(groupname = curgroup,
                                                        name = curpar,
                                                        value =  val,
                                                        type = Type_Param.String,
                                                        comment = '')

                            elif isinstance(val, int):
                                self.mytooltip.add_param(groupname = curgroup,
                                                        name = curpar,
                                                        value =  int(val),
                                                        type = Type_Param.Integer,
                                                        comment = '')

                            else:
                                self.mytooltip.add_param(groupname = curgroup,
                                                        name = curpar,
                                                        value =  float(val),
                                                        type = Type_Param.Float,
                                                        comment = '')

                except:
                    pass

        for locarray in self.myarrays:
            if locarray.checked:
                curgroup = locarray.idx

                try:
                    val = locarray.get_value(x, y)

                    if val != -99999.:

                        if locarray.wolftype in WOLF_ARRAY_MB:
                            i, j, curbloc = locarray.get_blockij_from_xy(x, y)
                            curpar = 'Indices (i;j;bloc) (1-based)'

                            self.mytooltip.add_param(groupname = curgroup,
                                                        name = curpar,
                                                        value =  '(' + str(i+1) + ';' + str(j+1) + ';' + str(curbloc) + ')',
                                                        type = Type_Param.String,
                                                        comment = '')

                        else:
                            i, j = locarray.get_ij_from_xy(x, y)
                            curpar = 'Indices (i;j) (1-based)'

                            self.mytooltip.add_param(groupname = curgroup,
                                                        name = curpar,
                                                        value =  '(' + str(i+1) + ';' + str(j+1) + ')',
                                                        type = Type_Param.String,
                                                        comment = '')

                        curpar = 'Value'

                        if val is np.nan:
                            self.mytooltip.add_param(groupname = curgroup,
                                                    name = 'Value',
                                                    value =  "Nan",
                                                    type = Type_Param.String,
                                                    comment = '')
                        elif np.ma.is_masked(val):

                            self.mytooltip.add_param(groupname = curgroup,
                                                    name = 'Value',
                                                    value =  "Masked",
                                                    type = Type_Param.String,
                                                    comment = '')

                        elif isinstance(val, str):
                            self.mytooltip.add_param(groupname = curgroup,
                                                    name = curpar,
                                                    value =  val,
                                                    type = Type_Param.String,
                                                    comment = '')

                        elif isinstance(val, int):
                            self.mytooltip.add_param(groupname = curgroup,
                                                    name = curpar,
                                                    value =  int(val),
                                                    type = Type_Param.Integer,
                                                    comment = '')

                        else:
                            self.mytooltip.add_param(groupname = curgroup,
                                                    name = curpar,
                                                    value =  float(val),
                                                    type = Type_Param.Float,
                                                    comment = '')

                except:
                    pass

        if self.linked:
            for curFrame in self.linkedList:
                if not curFrame is self:
                    for locarray in curFrame.myarrays:
                        curgroup = locarray.idx
                        if locarray.plotted:

                            try:
                                val = locarray.get_value(x, y)

                                if val != -99999.:
                                    if locarray.wolftype in WOLF_ARRAY_MB:
                                        i, j, curbloc = locarray.get_blockij_from_xy(x, y)
                                        curpar = 'Indices (i;j;bloc) (1-based)'

                                        self.mytooltip.add_param(groupname = locarray.idx,
                                                                    name = curpar,
                                                                    value =  '(' + str(i+1) + ';' + str(j+1) + ';' + str(curbloc) + ')',
                                                                    type = Type_Param.String,
                                                                    comment = '')

                                    else:
                                        i, j = locarray.get_ij_from_xy(x, y)
                                        curpar = 'Indices (i;j) (1-based)'

                                        self.mytooltip.add_param(groupname = locarray.idx,
                                                                    name = curpar,
                                                                    value =  '(' + str(i+1) + ';' + str(j+1) + ')',
                                                                    type = Type_Param.String,
                                                                    comment = '')

                                    curpar = 'Value'

                                    if val is np.nan:
                                        self.mytooltip.add_param(groupname = locarray.idx,
                                                                name = 'Value',
                                                                value =  "Nan",
                                                                type = Type_Param.String,
                                                                comment = '')

                                    elif np.ma.is_masked(val):

                                        self.mytooltip.add_param(groupname = locarray.idx,
                                                                name = 'Value',
                                                                value =  "Masked",
                                                                type = Type_Param.String,
                                                                comment = '')

                                    elif isinstance(val, str):
                                        self.mytooltip.add_param(groupname = locarray.idx,
                                                                name = curpar,
                                                                value =  val,
                                                                type = Type_Param.String,
                                                                comment = '')
                                    elif isinstance(val, int):
                                        self.mytooltip.add_param(groupname = locarray.idx,
                                                                name = curpar,
                                                                value =  int(val),
                                                                type = Type_Param.Integer,
                                                                comment = '')
                                    else:
                                        self.mytooltip.add_param(groupname = locarray.idx,
                                                                name = curpar,
                                                                value =  float(val),
                                                                type = Type_Param.Float,
                                                                comment = '')
                            except:
                                logging.warning(_('Error in linked frame -- Please check !'))

        for loc_ps in self.mypartsystems:
            if loc_ps.checked:
                curgroup = loc_ps.idx
                try:
                    self.mytooltip.add_param(groupname = curgroup,
                                            name = _('Step [s]'),
                                            value =  loc_ps.current_step,
                                            type = Type_Param.Float,
                                            comment = 'Step in seconds')
                    self.mytooltip.add_param(groupname = curgroup,
                                            name = _('Step [-]'),
                                            value =  loc_ps.current_step_idx,
                                            type = Type_Param.Integer,
                                            comment = 'Step index')
                except:
                    pass

        self.mytooltip.PopulateOnePage()


    def OnMotion(self, e: wx.MouseEvent):
        """ Mouse move event """

        # Déplacement de la souris sur le canvas OpenGL
        posframe = self.GetPosition()
        pos = e.GetPosition()
        x, y = self.getXY(pos)
        altdown = e.AltDown()
        shiftdown = e.ShiftDown()

        if e.LeftIsDown() or e.MiddleIsDown():
            self.mousex -= x - self.mousedown[0]
            self.mousey -= y - self.mousedown[1]

            self.setbounds(False)
            return
        elif e.RightIsDown():
            if self.action == 'select bc':
                if self.active_vector is None:

                    self.end_action(_('None because no active vector'))
                    return

                self.active_vector.myvertices=[wolfvertex(self.rightdown[0],self.rightdown[1]),
                                               wolfvertex(self.rightdown[0],y),
                                               wolfvertex(x,y),
                                               wolfvertex(x,self.rightdown[1]),
                                               wolfvertex(self.rightdown[0],self.rightdown[1])]
        else:
            self.move=False

        if self.action is not None:

            if 'select by tmp vector' in self.action or \
               'select by vector' in self.action or \
               self.action == 'capture vertices' or \
               self.action == 'dynamic parallel' or \
               self.action == 'laz tmp vector' or \
               self.action == 'create polygon - tiles':

                if self.active_vector.nbvertices > 0:
                    self.active_vector.myvertices[-1].x = x
                    self.active_vector.myvertices[-1].y = y

            if self.action == 'modify vertices' or \
               self.action == 'insert vertices':
                if self.active_vertex is not None:
                    if shiftdown:
                        ox = self.active_vector.myvertices[0].x
                        oy = self.active_vector.myvertices[0].y

                        dirx = self.active_vector.myvertices[-1].x - ox
                        diry = self.active_vector.myvertices[-1].y - oy
                        normdir = np.sqrt(dirx ** 2. + diry ** 2.)

                        dirx /= normdir
                        diry /= normdir

                        vecx = x - ox
                        vecy = y - oy

                        norm = np.sqrt(vecx ** 2. + vecy ** 2.)

                        self.active_vertex.x = ox + np.inner([dirx, diry], [vecx, vecy]) * dirx
                        self.active_vertex.y = oy + np.inner([dirx, diry], [vecx, vecy]) * diry

                    else:
                        self.active_vertex.x = x
                        self.active_vertex.y = y

                    self.active_vertex.limit2bounds(self.active_vector.mylimits)

            if self.action == 'dynamic parallel':
                self.active_zone.parallel_active(self.dynapar_dist)

            self.Paint()

        self._last_mouse_pos = (x,y,pos)
        self._update_mytooltip()

        if e.ControlDown():
            if self._oldpos_tooltip is None:
                # Store the position of the tooltip
                # Useful to restore it after CTRL is released
                self._oldpos_tooltip = self.mytooltip.GetPosition()

            self.mytooltip.SetWindowStyle(wx.STAY_ON_TOP) # Just on top, without Title bar
            ttsize = self.mytooltip.GetSize()
            self.mytooltip.position(pos + posframe + (ttsize[0] / 2. + 15, 15))
        else:

            width, height = self.GetSize()

            if self.IsMaximized():
                # Frame is maximized -> tooltip must be on the Screen
                self.mytooltip.SetWindowStyle(wx.STAY_ON_TOP)
            else:

                if self._oldpos_tooltip is None:
                    # No old position stored -> tooltip does not move
                    pos_tooltip = self.mytooltip.GetPosition()
                else:
                    # Restore the position of the tooltip
                    pos_tooltip = self._oldpos_tooltip

                # Reset the old position, so when CTRL is pressed again, the memory will be updated
                self._oldpos_tooltip = None

                if shiftdown or (pos_tooltip[0] == 0 and pos_tooltip[1] == 0):
                    # SHIFT is pressed or tooltip is at the top right corner of the Frame
                    # or it is the first time the tooltip is displayed
                    posframe[0] += width
                    posframe[1] -= 50
                    self.mytooltip.position(posframe)
                    w, h = self.mytooltip.GetSize()
                    self.mytooltip.SetSize((w, height))

                else:
                    # Force the position
                    self.mytooltip.SetPosition(pos_tooltip)

                self.mytooltip.SetIcon(self.GetIcon()) # update icon
                self.mytooltip.SetWindowStyle(wx.DEFAULT_FRAME_STYLE | wx.STAY_ON_TOP) # on top, with Title bar

        # self.mytooltip.Show(True)

    def Autoscale(self, update_backfore=True):
        """ Redimensionnement de la fenêtre pour afficher tous les objets """

        self.findminmax()
        self.width = self.xmax - self.xmin
        self.height = self.ymax - self.ymin

        centerx = self.xmin + self.width / 2.
        centery = self.ymin + self.height / 2.

        iwidth = self.width * self.sx
        iheight = self.height * self.sy

        width, height = self.canvas.GetSize()

        if iwidth == 0 or iheight == 0:
            logging.error(_('Width or height of the canvas is null -- Please check the "findminmax" routine in "Autoscale" !'))
            iwidth = 1
            iheight = 1

        sx = float(width) / float(iwidth)
        sy = float(height) / float(iheight)

        if sx == 0 or sy == 0:
            logging.error(_('At least one scale factor is null -- Please check the "Autoscale" routine !'))
            sx = 1.
            sy = 1.

        if sx > sy:
            self.xmax = self.xmin + self.width * sx / sy
            self.width = self.xmax - self.xmin
        else:
            self.ymax = self.ymin + self.height * sy / sx
            self.height = self.ymax - self.ymin

        self.mousex = centerx
        self.mousey = centery

        if update_backfore:
            # dessin du background
            for obj in self.iterator_over_objects(draw_type.WMSBACK):
                obj.reload()

            # dessin du foreground
            for obj in self.iterator_over_objects(draw_type.WMSFORE):
                obj.reload()

        self.setbounds()

    def _endactions(self):
        """
        End of actions

        Call when the user double click on the right button of the mouse or press return.

        Depending on the action, the method will call differnt routines and refresh the figure.

        Each action must call self.end_action() to nullify the action and print a message.
        """

        if self.action is not None:
            locaction = self.action
            if 'select by tmp vector' in self.action or 'select by vector' in self.action:
                inside_under = 'inside' in self.action

                self.end_action(_('End of vector selection'))

                self.active_vector.myvertices.pop(-1)

                if inside_under:
                    self.active_vector.close_force()
                    self.active_array.SelectionData.select_insidepoly(self.active_vector)
                else:
                    self.active_array.SelectionData.select_underpoly(self.active_vector)

                if 'tmp' in locaction:
                    # we must reset the temporary vector
                    self.active_vector.reset()

            elif self.action == 'pick landmap':

                self.end_action(_('End of landmap picking'))

            elif self.action == 'laz tmp vector':
                self.end_action(_('End of LAZ selection'))
                self.active_vector.myvertices.pop(-1)
                self.plot_laz_around_active_vec()
                self.active_vector.reset()

            elif self.action == 'create polygon - tiles':
                self.end_action(_('End of polygon creation'))
                self.active_vector.myvertices.pop(-1)
                self.active_vector.close_force()

                dlg = wx.MessageDialog(None, _('Do you want to align vertices on magnetic grid ?'), _('Confirm'), wx.YES_NO | wx.YES_DEFAULT | wx.ICON_QUESTION)
                ret = dlg.ShowModal()
                dlg.Destroy()

                if ret == wx.ID_YES:
                    dlg = wx.NumberEntryDialog(None, _('Which is the sptial step size [m] ?'), _('Size'),  _('Spatial grid size'), 50, 1, 10000)
                    ret = dlg.ShowModal()
                    if ret == wx.ID_CANCEL:
                        dlg.Destroy()
                        return
                    ds = dlg.GetValue()
                    dlg.Destroy()

                    vertices = self.active_vector.myvertices.copy()
                    self.active_vector.myvertices.clear()

                    x_aligned = np.asarray([(curvert.x // ds)*ds for curvert in vertices])
                    y_aligned = np.asarray([(curvert.y // ds)*ds for curvert in vertices])

                    if (x_aligned.min() == x_aligned.max()) and (y_aligned.min() == y_aligned.max()):
                        logging.error(_('All vertices are aligned on the same point -- Choose another step size'))
                        return

                    self.active_vector.add_vertex(wolfvertex(x_aligned.min(), y_aligned.min()))
                    self.active_vector.add_vertex(wolfvertex(x_aligned.max(), y_aligned.min()))
                    self.active_vector.add_vertex(wolfvertex(x_aligned.max(), y_aligned.max()))
                    self.active_vector.add_vertex(wolfvertex(x_aligned.min(), y_aligned.max()))
                    self.active_vector.add_vertex(wolfvertex(x_aligned.min(), y_aligned.min()))
                    self.active_vector.close_force()
                    self.active_vector.find_minmax()

                self._create_data_from_tiles_common()

            elif self.action == 'capture vertices':
                self.end_action(_('End of points capturing'))
                self.active_vector.myvertices.pop(-1)
                r = wx.MessageDialog(
                    None,
                    _('End of points capturing') + '\n' +
                    _('Force to close the vector ?'),
                    _('Confirm'),
                    wx.YES_NO | wx.YES_DEFAULT | wx.ICON_QUESTION
                ).ShowModal()
                if r == wx.ID_YES:
                    self.active_vector.close_force()

                # force to prepare OpenGL to accelerate the plot
                # Le test not(self in self.linkedList) permet de ne pas créer le liste OpenGL en cas de multi-viewers
                # car une liste OpenGL ne sera pas tracée sur les autres fenêtres
                # C'est donc plus lent mais plus sûr pour que l'affichage dynamique soit correct
                self.active_vector.parentzone.plot(prep = self.linkedList is None or not(self in self.linkedList))

            elif self.action == 'modify vertices':

                # end of vertices modification
                self.end_action(_('End of vertices modification'))

                # force to prepare OpenGL to accelerate the plot
                # Le test not(self in self.linkedList) permet de ne pas créer le liste OpenGL en cas de multi-viewers
                # car une liste OpenGL ne sera pas tracée sur les autres fenêtres
                # C'est donc plus lent mais plus sûr pour que l'affichage dynamique soit correct
                self.active_vector.parentzone.plot(prep = self.linkedList is None or not(self in self.linkedList))
                self.active_zones.find_minmax(True)

                self.active_vertex = None

            elif self.action == 'insert vertices':
                self.end_action(_('End of vertices insertion'))

                # force to prepare OpenGL to accelerate the plot
                # Le test not(self in self.linkedList) permet de ne pas créer le liste OpenGL en cas de multi-viewers
                # car une liste OpenGL ne sera pas tracée sur les autres fenêtres
                # C'est donc plus lent mais plus sûr pour que l'affichage dynamique soit correct
                self.active_vector.parentzone.plot(prep = self.linkedList is None or not(self in self.linkedList))
                self.active_zones.find_minmax(True)

                self.active_vertex = None

            elif self.action == 'dynamic parallel':
                self.active_vector.myvertices.pop(-1)
                self.active_zone.parallel_active(self.dynapar_dist)

                self.active_zones.fill_structure()
                self.active_zones.find_minmax(True)

                # force to prepare OpenGL to accelerate the plot
                # Le test not(self in self.linkedList) permet de ne pas créer le liste OpenGL en cas de multi-viewers
                # car une liste OpenGL ne sera pas tracée sur les autres fenêtres
                # C'est donc plus lent mais plus sûr pour que l'affichage dynamique soit correct
                self.active_vector.parentzone.plot(prep = self.linkedList is None or not(self in self.linkedList))

                self.active_vertex = None

                self.end_action(_('End of dynamic parallel'))

            elif 'select active vector' in self.action:

                self.end_action(_('End of vector selection'))

            elif 'select node by node' in self.action:
                self.end_action(_('End of node by node selection'))

        self.copyfrom = None

        self.Refresh()
        self.mimicme()

    def print_About(self):
        """ Print the About window """
        from .apps.version import WolfVersion

        version = WolfVersion()
        dlg = wx.MessageDialog(None, _('Wolf - Version {}\n\n'.format(str(version))) + _('Developed by : ') + 'HECE ULiège\n' + _('Contact : pierre.archambeau@uliege.be'), _('About'), wx.OK | wx.ICON_INFORMATION)
        dlg.ShowModal()
        dlg.Destroy()

    def check_for_updates(self):
        """ Check for updates """
        from .apps.version import WolfVersion
        import requests

        current_version = str(WolfVersion())
        package_name = "wolfhece"

        try:
            available_version = requests.get(f"https://pypi.org/pypi/{package_name}/json").json()["info"]["version"]

            if available_version > current_version:
                with wx.MessageDialog(None, _("A new version is available: {}\n\nYour version is {}\n\nIf you want to upgrade, 'pip install wolfhece --upgrade' from your Python environment.").format(available_version, current_version), _("Upgrade"), wx.OK | wx.ICON_INFORMATION) as dlg:
                    dlg.ShowModal()
            else:
                with wx.MessageDialog(None, _("You have the latest version."), _("Upgrade"), wx.OK | wx.ICON_INFORMATION) as dlg:
                    dlg.ShowModal()

        except Exception as e:
            logging.error("Package not found on PyPI. -- {}".format(e))

    def print_shortcuts(self, inframe:bool = None):
        """ Print the list of shortcuts into logging """

        # shortcuts = "F1 : mise à jour du dernier pas de résultat\n \
        #     F2 : mise à jour du résultat pas suivant\n \
        #     F4 : mise à jour du particle system au pas suivant\n \
        #     Shift+F2 : mise à jour du résultat pas précédent\n \
        #     Shift+F4 : mise à jour du particle system au pas précédent\n \
        #     CTRL+F2 : choix du pas\n \
        #     CTRL+F4 : choix du pas (particle system)\n \
        #     CTRL+Shift+F2 : choix du pas sur base du temps\n \
        #     CTRL+Shift+F4 : choix du pas sur base du temps (particle system)\n \
        #     F5 : autoscale\n \
        #     F7 : refresh\n \
        #     F9 : sélection de toutes les mailles dans la matrice courante\n \
        #     F11 : sélection sur matrice courante\n \
        #     F12 : opération sur matrice courante\n \
        #     \n \
        #     ESPACE : pause/resume animation\n \
        #     \n \
        #     Z  : zoom avant\n \
        #     z  : zoom artrière\n \
        #     Flèches : déplacements latéraux\n \
        #     P : sélection de profil\n \
        #     1,2 : Transfert de la sélection de la amtrice courante vers le dictionnaire\n \
        #     F, CTRL+F : recherche de la polyligne dans la zone courante ou dans toutes les zones\n \
        #     i : interpolation2D sur base de la sélection sur la matrice courante\n \
        #     +,- (numpad) : augmente ou diminue la taille des flèches de resultats 2D\n \
        #     \n \
        #     o, O : Gestion de la transparence de la matrice courante\n \
        #     CTRL+o, CTRL+O : Gestion de la transparence du résultat courant\n \
        #     \n \
        #     !! ACTIONs !!\n \
        #     N : sélection noeud par noeud de la matrice courante\n \
        #     B : sélection par vecteur temporaire de la matrice courante\n \
        #     V : sélection par vecteur activé de la matrice courante - zone intérieure\n \
        #     r : reset de la sélection de la matrice courante\n \
        #     R : reset de toutes les sélections de la matrice courante\n \
        #     P : sélection de la section transversale par click souris\n \
        #     \n \
        #     RETURN : end current action (cf aussi double clicks droit 'OnRDClick')\n \
        #     DELETE : remove item\n \
        #     \n \
        #     CTRL+Q : Quit application\n \
        #     CTRL+U : Import GLTF/GLB\n \
        #     CTRL+C : Set copy source / Copy canvas to Clipboard\n \
        #     CTRL+V : Paste selected values\n \
        #     CTRL+ALT+V ou ALTGr+V : Paste/Recopy selection\n \
        #     CTRL+L : chargement d'une matrice sur base du nom de fichier de la tile\n \
        #     \n \
        #     ALT+C : Copy image"

        groups = ['Results', 'Particle system', 'Drawing', 'Arrays', 'Cross sections', 'Zones', 'Action', 'Tree', 'Tiles', 'GLTF/GLB', 'App']

        shortcuts = {'F1': _('Results : read the last step'),
                     'F2': _('Results : read the next step'),
                     'Shift+F2': _('Results : read the previous step'),
                     'CTRL+F2': _('Results : choose the step'),
                     'CTRL+Shift+F2': _('Results : choose the step based on time'),
                     '+,- (numpad)': _('Results : increase or decrease the size of 2D result arrows'),

                     'F4': _('Particle system : update to the next step'),
                     'Shift+F4': _('Particle system : update to the previous step'),
                     'CTRL+F4': _('Particle system : choose the step'),
                     'CTRL+Shift+F4': _('Particle system : choose the step based on time'),
                     'SPACE': _('Particle system : pause/resume animation'),

                     'LMB double clicks': _('Drawing : center the view on the clicked point -- future zoom will be centered on the point'),
                     'LMB and move': _('Drawing : translate the view'),
                     'Mouse wheel click and move': _('Drawing : translate the view'),
                     'Mouse wheel': _('Drawing : zoom in/out - centered on the middle of the canvas'),
                     'Mouse wheel + Space Bar': _('Drawing : zoom in/out - centered on the mouse position'),
                     'z, Z': _('Drawing : zoom out/in - centered on the middle of the canvas'),
                     'Touchpad 2 fingers': _('Drawing : zoom in/out - centered on the middle of the canvas'),
                     'CTRL + z': _('Drawing : Autoscale only on active array'),
                     'CTRL + Z': _('Drawing : Autoscale only on active vector'),

                     'F5': _('Drawing : autoscale'),
                     'F7': _('Drawing : refresh'),
                     'Arrows': _('Drawing : lateral movements'),
                     'c or C': _('Drawing : copy canvas to Clipboard wo axes'),
                     'CTRL+C': _('Drawing : copy canvas to Clipboard as Matplotlib image'),

                     'CTRL+o': _('Results : increase transparency of the current result'),
                     'CTRL+O': _('Results : decrease transparency of the current result'),
                     'o': _('Arrays : increase transparency of the current array'),
                     'O': _('Arrays : decrease transparency of the current array'),

                     'F9': _('Arrays : select all cells'),
                     'F11': _('Arrays : select by criteria'),
                     'F12': _('Arrays : operations'),
                     'n or N': _('Arrays : node-by-node selection'),
                     'b or B': _('Arrays : temporary vector selection'),
                     'v or V': _('Arrays : activated vector selection - inner zone'),

                     'r': _('Arrays : reset the selection'),
                     'R': _('Arrays : reset the selection and the associated dictionnary'),

                     '1,2...9': _('Arrays : transfer the selection to the associated dictionary - key 1 to 9'),

                     'i': _('Arrays : 2D interpolation based on the selection on the current matrix'),
                     'CTRL+C': _('Arrays : Set copy source and current selection to clipboard as string'),
                     'CTRL+X': _('Arrays : Crop the active array using the active vector and make a copy'),
                     'CTRL+V': _('Arrays : paste selected values'),
                     'CTRL+ALT+C or ALTGr+C': _('Arrays : Set copy source and current selection to clipboard as script'),
                     'CTRL+ALT+X or ALTGr+X': _('Arrays : Crop the active array using the active vector without masking the values outside the vector'),
                     'CTRL+ALT+V or ALTGr+V': _('Arrays : paste selection to active array'),

                     'p or P': _('Cross sections : Pick a profile/cross section'),

                     'f or F, CTRL+F': _('Zones : search for the polyline in the current zone or in all zones'),

                     'RETURN': _('Action : End the current action (see also right double-click -- OnRDClick)'),
                     'Press and Hold CTRL': _('Action : Data Frame follows the mouse cursor'),

                     'DELETE': _('Tree : Remove item'),

                     'CTRL+L': _('Tiles: Pick a tile by clicking on it'),

                     'CTRL+U': _('GLTF/GLB : import/update GLTF/GLB'),

                     'CTRL+Q': _('App : Quit application'),}

        def gettxt():
            txt = ''
            for curgroup in groups:
                txt += curgroup + '\n'
                for curkey, curval in shortcuts.items():
                    if curgroup in curval:
                        txt += '\t' + curkey + ' : ' + curval.split(':')[1] + '\n'
                txt += '\n'
            return txt

        logging.info(gettxt())

        if inframe :
            frame = wx.Frame(None, -1, _('Shortcuts'), size=(500, 800))
            # panel = wx.Panel(frame, -1)

            sizer = wx.BoxSizer(wx.VERTICAL)

            multiline = wx.TextCtrl(frame, -1, '', style=wx.TE_MULTILINE | wx.TE_READONLY | wx.TE_RICH2)

            multiline.SetValue(gettxt())

            sizer.Add(multiline, 1, wx.EXPAND)

            frame.SetSizer(sizer)

            icon = wx.Icon()
            icon_path = Path(__file__).parent / "apps/wolf_logo2.bmp"
            icon.CopyFromBitmap(wx.Bitmap(str(icon_path), wx.BITMAP_TYPE_ANY))
            frame.SetIcon(icon)

            frame.SetAutoLayout(True)
            frame.Layout()

            frame.Show()

    def msg_action(self, which:int = 0):
        """ Message to end action """

        if which == 0:
            self.set_statusbar_text(_('Action in progress... -- To quit, press "RETURN" or "double clicks RIGHT" or press "ESC"'))
        else:
            self.set_statusbar_text('')

    def start_action(self, action:str, message:str=''):
        """ Message to start action """

        assert isinstance(action, str), 'action must be a string'
        if action == '':
            self.action = None
        else:
            self.action = action.lower()
        logging.info(_('ACTION : ') + _(message) if message != '' else _('ACTION : ') + _(action))
        self.msg_action(0)

    def end_action(self, message:str=''):
        """ Message to end action """

        self.action = None
        logging.info(_('ACTION : ') + _(message) if message != '' else _('ACTION : End of action') )
        self.msg_action(1)

    def OnHotKey(self, e: wx.KeyEvent):
        """
        Gestion des touches clavier -- see print_shortcuts for more details
        """

        key = e.GetKeyCode()
        ctrldown = e.ControlDown()
        altdown = e.AltDown()
        shiftdown = e.ShiftDown()

        myobj = e.EventObject

        logging.debug(_('You are pressing key code : ') + str(key))
        if ctrldown:
            logging.debug(_('Ctrl is down'))
        if altdown:
            logging.debug(_('Alt is down'))

        if ctrldown or altdown:
            if key == wx.WXK_F2 and not shiftdown:

                if self.active_res2d is not None:
                    nb = self.active_res2d.get_nbresults()
                    dlg = wx.NumberEntryDialog(None,_('Please choose a step (1 -> {})'.format(nb)),'Step :', _('Select a specific step'), nb, min=1, max=nb)
                    ret = dlg.ShowModal()

                    nb = dlg.GetValue()
                    dlg.Destroy()

                    self.active_res2d.read_oneresult(nb-1)
                    self.active_res2d.set_currentview()
                    self.Refresh()

                else:
                    logging.info(_('Please activate a simulation before search a specific result'))

            elif key == wx.WXK_F2 and shiftdown:

                if self.active_res2d is not None:
                    nb = self.active_res2d.get_nbresults()

                    choices = ['{:.3f} [s] - {} [h:m:s]'.format(cur, timedelta(seconds=int(cur),
                                                                   milliseconds=int(cur-int(cur))*1000))
                               for cur in self.active_res2d.times]

                    dlg = wx.SingleChoiceDialog(None,
                                                _('Please choose a time step'),
                                                _('Select a specific step'),
                                                choices)
                    ret = dlg.ShowModal()
                    if ret == wx.ID_CANCEL:
                        dlg.Destroy()
                        return

                    keyvalue = dlg.GetStringSelection()
                    dlg.Destroy()

                    self.active_res2d.read_oneresult(choices.index(keyvalue))
                    self.active_res2d.set_currentview()
                    self.Refresh()

                else:
                    logging.info(_('Please activate a simulation before searching a specific result'))

            if key == wx.WXK_F4 and not shiftdown:

                if self.active_particle_system is not None:

                    nb = self.active_particle_system.nb_steps
                    dlg = wx.NumberEntryDialog(None,_('Please choose a step (1 -> {})'.format(nb)),'Step :', _('Select a specific step'), nb, min=1, max=nb)
                    ret = dlg.ShowModal()

                    nb = dlg.GetValue()
                    dlg.Destroy()

                    self.active_particle_system.current_step = nb-1
                    self.Refresh()
                    self._update_mytooltip()

                else:
                    logging.info(_('Please activate a particle system before searching a specific result'))

            elif key == wx.WXK_F4 and shiftdown:

                if self.active_particle_system is not None:

                    choices = ['{:.3f} [s] - {} [h:m:s]'.format(cur, timedelta(seconds=int(cur),
                                                                   milliseconds=int(cur-int(cur))*1000))
                               for cur in self.active_particle_system.get_times()]

                    dlg = wx.SingleChoiceDialog(None,
                                                _('Please choose a time step'),
                                                _('Select a specific step'),
                                                choices)
                    ret = dlg.ShowModal()
                    if ret == wx.ID_CANCEL:
                        dlg.Destroy()
                        return

                    keyvalue = dlg.GetStringSelection()
                    dlg.Destroy()

                    self.active_particle_system.current_step = choices.index(keyvalue)
                    self.Refresh()
                    self._update_mytooltip()

                else:
                    logging.info(_('Please activate a simulation before search a specific result'))

            elif key == wx.WXK_NUMPAD_ADD: #+ from numpad
                if self.active_res2d is not None:
                    self.active_res2d.update_zoom_2(1.1)
                    self.Refresh()

            elif key == wx.WXK_NUMPAD_SUBTRACT: #- from numpad
                if self.active_res2d is not None:
                    self.active_res2d.update_zoom_2(1./1.1)
                    self.Refresh()

            elif key == ord('X'):
                # Create a new array from the active array and the active vector
                # Node outside the vector are set to NullValue
                if self.active_array is not None and self.active_vector is not None:

                    bbox = self.active_vector.get_bounds_xx_yy()
                    newarray = self.active_array.crop_array(bbox)

                    if not altdown:
                        newarray.mask_outsidepoly(self.active_vector)

                    newarray.nullify_border(width=1)

                    self.add_object('array', newobj = newarray, id = self.active_array.idx + '_crop')

                    self.Refresh()

            elif key == ord('Q'):
                # If Ctrl-Q is hit, then we must *not* handle it
                # because it is tied to the Ctrl-Q accelerator
                # of the "quit" menu...
                e.Skip()
                return

            if key == ord('U'):
                # CTRL+U
                # Mise à jour des données par import du fichier gtlf2
                msg = ''
                if self.active_array is None:
                    msg += _('Active array is None\n')

                if msg != '':
                    msg += _('\n')
                    msg += _('Retry !\n')
                    wx.MessageBox(msg)
                    return

                self.set_fn_fnpos_gltf()
                self.update_blender_sculpting()

            elif key == ord('F'):
                if self.active_zones is not None:
                    self.start_action('select active vector all', _('Select active vector all'))

            elif key == ord('L'):
                if self.active_tile is not None:
                    self.start_action('select active tile', _('Select active tile'))

            elif key == wx.WXK_UP:
                self.upobj()

            elif key == wx.WXK_DOWN:
                self.downobj()

            elif key == ord('C') and altdown and not ctrldown:
                # ALT+C
                #Copie du canvas dans le clipboard pour transfert vers autre application
                self.copy_canvasogl()

            elif key == ord('C') and ctrldown and not altdown:
                # CTRL+C
                if self.active_array is None:
                    dlg = wx.MessageDialog(self,
                                        _('The active array is None - Please active an array from which to copy the values !'),
                                        style=wx.OK)
                    dlg.ShowModal()
                    dlg.Destroy()
                    return

                logging.info(_('Start copying values / Current selection to clipboard'))
                self.copyfrom = self.active_array
                self.mimicme_copyfrom()  # force le recopiage de copyfrom dans les autres matrices liées

                if len(self.active_array.SelectionData.myselection) > 5000:
                    dlg = wx.MessageDialog(self, _('The selection is large, copy to clipboard may be slow ! -- Continue?'), style=wx.OK | wx.CANCEL)
                    ret = dlg.ShowModal()

                    if ret == wx.ID_CANCEL:
                        logging.info(_('Copy to clipboard cancelled -- But source array is well defined !'))
                        dlg.Destroy()
                        return

                    dlg.Destroy()

                self.active_array.SelectionData.copy_to_clipboard()

                logging.info(_('Values copied to clipboard'))

            elif key == ord('C') and ctrldown and altdown:
                if self.active_array is None:
                    dlg = wx.MessageDialog(self,
                                        _('The active array is None - Please active an array from which to copy the selection !'),
                                        style=wx.OK)
                    dlg.ShowModal()
                    dlg.Destroy()
                    return

                logging.info(_('Start copying selection / Current selection to clipboard as script (Python)'))
                self.copyfrom = self.active_array
                self.mimicme_copyfrom()  # force le recopiage de copyfrom dans les autres matrices liées

                if len(self.active_array.SelectionData.myselection) > 5000:
                    dlg = wx.MessageDialog(self, _('The selection is large, copy to clipboard may be slow ! -- Continue?'), style=wx.OK | wx.CANCEL)
                    ret = dlg.ShowModal()

                    if ret == wx.ID_CANCEL:
                        logging.info(_('Copy script to clipboard cancelled -- But source array is well defined !'))
                        dlg.Destroy()
                        return

                    dlg.Destroy()

                self.active_array.SelectionData.copy_to_clipboard(typestr='script')

                logging.info(_('Selection copied to clipboard as script (Python)'))

            elif key == ord('V') and ctrldown:
                # CTRL+V
                # CTRL+ALT+V ou Alt Gr + V

                if self.active_array is None:
                    if e.AltDown():
                        # CTRL+ALT+V
                        logging.warning(_('The active array is None - Please active an array into which to paste the selection !'))
                    else:
                        logging.warning(_('The active array is None - Please active an array into which to paste the values !'))

                    return

                fromarray = self.copyfrom
                if fromarray is None:
                    if self.linked:
                        if not self.linkedList is None:
                            for curFrame in self.linkedList:
                                if curFrame.copyfrom is not None:
                                    fromarray = curFrame.copyfrom
                                    break

                if fromarray is None:
                    logging.warning(_('No selection to be pasted !'))
                    return

                cursel = fromarray.SelectionData.myselection

                if e.AltDown():
                    logging.info(_('Paste selection position'))

                    if cursel == 'all':
                        self.active_array.SelectionData.OnAllSelect(0)
                    elif len(cursel) > 0:
                        self.active_array.SelectionData.myselection = cursel.copy()
                        self.active_array.SelectionData.update_nb_nodes_selection()

                else:
                    logging.info(_('Paste selection values'))
                    if cursel == 'all':
                        self.active_array.paste_all(fromarray)

                    elif len(cursel) > 0:
                        z = fromarray.SelectionData.get_values_sel()
                        self.active_array.set_values_sel(cursel, z)

                self.Refresh()

                logging.info(_('Selection/Values pasted'))

            elif key == ord('Z'):

                if ctrldown:
                    if shiftdown:
                        if self.active_vector is not None:
                            self.zoom_on_vector(self.active_vector, canvas_height= self.canvas.GetSize()[1])
                        else:
                            logging.warning(_('No active vector to zoom on !'))
                    else:
                        if self.active_array is not None:
                            self.zoom_on_array(self.active_array, canvas_height= self.canvas.GetSize()[1])
                        else:
                            logging.warning(_('No active array to zoom on !'))
        else:
            if key == wx.WXK_DELETE:
                self.removeobj()

            elif key == wx.WXK_ESCAPE:

                logging.info(_('Escape key pressed -- Set all active objects and "action" to None'))

                self.action = None
                self.active_array = None
                self.active_vector = None
                self.active_zone = None
                self.active_zones = None
                self.active_res2d = None
                self.active_tile = None
                self.active_particle_system = None
                self.active_vertex = None
                self.active_cloud = None

                self.set_statusbar_text(_('Esc pressed - No more action in progress - No more active object'))
                self.set_label_selecteditem('')

            elif key == ord('C'):

                 self.copy_canvasogl(mpl = False)

            elif key == wx.WXK_SPACE:
                if self.timer_ps is not None and self.active_particle_system is not None :
                    if self.timer_ps.IsRunning():
                        self.timer_ps.Stop()
                    else:
                        if self.active_particle_system.current_step_idx == self.active_particle_system.nb_steps-1:
                            self.active_particle_system.current_step_idx = 0
                            self.active_particle_system.current_step = 0
                        self.timer_ps.Start(1000. / self.active_particle_system.fps)

            elif key == 388: #+ from numpad
                if self.active_res2d is not None:
                    self.active_res2d.update_arrowpixelsize_vectorfield(-1)
                    self.Refresh()

            elif key == 390: #- from numpad
                if self.active_res2d is not None:
                    self.active_res2d.update_arrowpixelsize_vectorfield(1)
                    self.Refresh()

            elif key == 13 or key==370 or key == wx.WXK_RETURN or key == wx.WXK_NUMPAD_ENTER:
                # 13 = RETURN classic keyboard
                # 370 = RETURN NUMPAD
                self._endactions()

            elif key == ord('I'):
                if self.active_array is not None :
                    self.active_array.interpolation2D()

            elif key == ord('F'):
                if self.active_zone is not None:
                    self.start_action('select active vector2 all', _('Select active vector2 all'))

            elif key in LIST_1TO9:

                if self.active_array is not None:
                    colors = [(0, 0, 255, 255), (0, 255, 0, 255), (0, 128, 255, 255), (255, 255, 0, 255), (255, 165, 0, 255), (128, 0, 128, 255), (255, 192, 203, 255), (165, 42, 42, 255), (128, 128, 128, 255)]
                    idx = LIST_1TO9.index(key)
                    if idx > 8:
                        idx -= 9

                    self.active_array.SelectionData.move_selectionto(str(idx+1), colors[idx])

            elif key == wx.WXK_F1:
                self.read_last_result()

            elif key == wx.WXK_F2 and shiftdown:
                self.simul_previous_step()

            elif key == wx.WXK_F4 and shiftdown:
                self.particle_previous_step()

            elif key == wx.WXK_F4:
                self.particle_next_step()

            elif key == wx.WXK_F2:
                self.simul_next_step()

            elif key == wx.WXK_F5:
                # Autoscale
                self.Autoscale()

            elif key == wx.WXK_F7:
                self.update()

            elif key == wx.WXK_F12 or key == wx.WXK_F11:
                if self.active_array is not None:
                    self.active_array.myops.SetTitle(_('Operations on array: ')+self.active_array.idx)
                    self.active_array.myops.Show()
                    self.active_array.myops.array_ops.SetSelection(1)
                    self.active_array.myops.Center()

            elif key == wx.WXK_F9:

                if self.active_array is not None:
                    if self.active_array.SelectionData is not None:
                        self.active_array.SelectionData.myselection = 'all'
                        logging.info(_('Selecting all nodes in the active array !'))
                    else:
                        logging.warning(_('No selection manager for this array !'))

                    if self.active_array.myops is not None:
                        self.active_array.myops.nbselect.SetLabelText('All')
                    else:
                        logging.warning(_('No operations manager for this array !'))

            elif key == ord('N'):  # N
                if self.active_array is not None:
                    self.active_array.myops.select_node_by_node()

                if self.active_res2d is not None:
                    self.active_res2d.properties.select_node_by_node()

            elif key == ord('V'):  # V
                if self.active_array is not None:
                    self.active_array.myops.select_vector_inside_manager()

            elif key == ord('B'):  # B
                if self.active_array is not None:
                    self.active_array.myops.select_vector_inside_tmp()

            elif key == ord('P'):  # P
                self.start_action('Select nearest profile', _('Select nearest profile'))

            elif key == ord('Z') and shiftdown:  # Z
                self.width = self.width / 1.1
                self.height = self.height / 1.1
                self.setbounds()

            elif key == ord('Z'):  # z
                self.width = self.width * 1.1
                self.height = self.height * 1.1
                self.setbounds()

            elif key == ord('R') and shiftdown:  # R
                if self.active_array is not None:
                    self.active_array.myops.reset_all_selection()
                    self.Refresh()

            elif key == ord('R'):  # r
                if self.active_array is not None:
                    self.active_array.myops.reset_selection()
                    self.Refresh()

            elif key == ord('O'):
                # Active Opacity for the active array

                if ctrldown:
                    if self.active_res2d is None:
                        logging.warning(_('No active result 2D to change the opacity !'))
                        return

                    if shiftdown:
                        self.active_res2d.set_opacity(self.active_res2d.alpha + 0.25)
                    else:
                        self.active_res2d.set_opacity(self.active_res2d.alpha - 0.25)

                else:
                    if self.active_array is None:
                        logging.warning(_('No active array to change the opacity !'))
                        return

                    if shiftdown:
                        self.active_array.set_opacity(self.active_array.alpha + 0.25)
                    else:
                        self.active_array.set_opacity(self.active_array.alpha - 0.25)

            elif key == wx.WXK_UP:
                self.mousey = self.mousey + self.height / 10.
                self.setbounds()

            elif key == wx.WXK_DOWN:
                self.mousey = self.mousey - self.height / 10.
                self.setbounds()

            elif key == wx.WXK_LEFT:
                self.mousex = self.mousex - self.width / 10.
                self.setbounds()

            elif key == wx.WXK_RIGHT:
                self.mousex = self.mousex + self.width / 10.
                self.setbounds()

    def paste_values(self,fromarray:WolfArray):
        """ Paste selected values from a WolfArray to the active array """

        if self.active_array is None:
            logging.warning(_('The active array is None - Please active an array into which to paste the values !'))
            return

        logging.info(_('Paste selection values'))
        cursel = fromarray.SelectionData.myselection
        if cursel == 'all':
            self.active_array.paste_all(fromarray)
        elif len(cursel) > 0:
            z = fromarray.SelectionData.get_values_sel()
            self.active_array.set_values_sel(cursel, z)

    def paste_selxy(self,fromarray:WolfArray):
        """ Paste selected nodes from a WolfArray to the active array """

        if self.active_array is None:
            logging.warning(_('The active array is None - Please active an array into which to paste the selection !'))
            return

        logging.info(_('Paste selection position'))
        cursel = fromarray.SelectionData.myselection
        if cursel == 'all':
            self.active_array.SelectionData.OnAllSelect(0)
        elif len(cursel) > 0:
            self.active_array.SelectionData.myselection = cursel.copy()
            self.active_array.SelectionData.update_nb_nodes_selection()

    def OntreeRight(self, e: wx.MouseEvent):
        """ Gestion du menu contextuel sur l'arbre des objets """

        if self.selected_object is not None:

            # On va nettoyer le menu contextuel car certaines entrées ne sont
            # pas nécessairement pertinentes

            # Chaînes à supprimer
            tracks=[]
            tracks.append(_('Boundary conditions'))
            tracks.append(_('Convert to mono-block'))
            tracks.append(_('Convert to mono-block (result)'))
            tracks.append(_('Convert to multi-blocks (result)'))
            tracks.append(_('Export to Shape file'))
            tracks.append(_('Export active zone to Shape file'))

            # Récupération des items du menu contextuel
            menuitems = self.popupmenu.GetMenuItems()
            text = [cur.GetItemLabelText() for cur in menuitems]

            # Liste des indices à supprimer
            # Pas possible de supprimer à la volée car cela modifie la liste
            to_delete = []
            for track in tracks:
                if track in text:
                    to_delete.append(text.index(track))

            # Suppression des items
            if len(to_delete) > 0:
                # Suppression en ordre décroissant pour ne pas décaler les indices
                to_delete.sort(reverse=True)
                for idx in to_delete:
                    self.popupmenu.Remove(menuitems[idx])

            # Add specific menu items for WolfArray
            if isinstance(self.selected_object, WolfArray):
                bc = self.get_boundary_manager(self.selected_object)
                if bc is not None:
                    self.popupmenu.Append(wx.ID_ANY, _('Boundary conditions'), _('Boundary conditions'))

            # Add specific menu items for WolfArrayMB
            if isinstance(self.selected_object, WolfArrayMB):
                self.popupmenu.Append(wx.ID_ANY, _('Convert to mono-block'), _('Convert to mono-block'))

            # Add specific menu items for Wolfresults_2D
            if isinstance(self.selected_object, Wolfresults_2D):
                self.popupmenu.Append(wx.ID_ANY, _('Convert to mono-block (result)'), _('Convert to mono-block'))
                self.popupmenu.Append(wx.ID_ANY, _('Convert to multi-blocks (result)'), _('Convert to multi-blocks'))

            if isinstance(self.selected_object, Zones):
                self.popupmenu.Append(wx.ID_ANY, _('Export to Shape file'), _('Export to Shape file'))
                self.popupmenu.Append(wx.ID_ANY, _('Export active zone to Shape file'), _('Export active zone to Shape file'))

            self.treelist.PopupMenu(self.popupmenu)

    def update(self):
        """
        Update backgournd et foreground elements and arrays if local minmax is checked.

        """

        # dessin du background
        for obj in self.iterator_over_objects(draw_type.WMSBACK):
            obj.reload()

        # dessin du foreground
        for obj in self.iterator_over_objects(draw_type.WMSFORE):
            obj.reload()

        if self.locminmax.IsChecked() or self.update_absolute_minmax:
            for curarray in self.iterator_over_objects(draw_type.ARRAYS):
                curarray: WolfArray
                if self.update_absolute_minmax:
                    curarray.updatepalette()
                    self.update_absolute_minmax = False
                else:
                    curarray.updatepalette(onzoom=[self.xmin, self.xmax, self.ymin, self.ymax])
                curarray.delete_lists()

        self.Paint()

    def _plotting(self, drawing_type: draw_type, checked_state: bool = True):
        """ Drawing objets on canvas"""

        try:
            for curobj in self.iterator_over_objects(drawing_type, checked_state=checked_state):
                if not curobj.plotting:
                    curobj.plotting = True
                    curobj.plot(sx = self.sx, sy=self.sy, xmin=self.xmin, ymin=self.ymin, xmax=self.xmax, ymax=self.ymax, size = (self.xmax - self.xmin) / 100.)
                    curobj.plotting = False
        except Exception as ex:
            curobj.plotting = False
            logging.error(_('Error while plotting objects of type {}').format(drawing_type.name))

            traceback.print_exc()
            logging.error(ex)

    def get_MVP_Viewport_matrix(self):
        """ Get the modelview projection matrix """

        if self.SetCurrentContext():
            modelview = glGetFloatv(GL_MODELVIEW_MATRIX)
            projection = glGetFloatv(GL_PROJECTION_MATRIX)
            viewport = glGetIntegerv(GL_VIEWPORT)

            return modelview, projection, viewport
        else:

            return None, None, None

    def SetCurrentContext(self):
        """ Set the current OGL context if exists otherwise return False """

        if self.context is None:
            return False

        return self.canvas.SetCurrent(self.context)

    def Paint(self):
        """ Dessin des éléments ajoutés au viewer """

        if self.currently_readresults:
            return

        width, height = self.canvas.GetSize()

        # C'est bien ici que la zone de dessin utile est calculée sur base du centre et de la zone en coordonnées réelles
        # Les commandes OpenGL sont donc traitées en coordonnées réelles puisque la commande glOrtho définit le cadre visible
        self.xmin = self.mousex - self.width / 2.
        self.ymin = self.mousey - self.height / 2.
        self.xmax = self.mousex + self.width / 2.
        self.ymax = self.mousey + self.height / 2.

        if self.SetCurrentContext():

            bkg_color = self.bkg_color

            glClearColor(bkg_color[0]/255., bkg_color[1]/255., bkg_color[2]/255., bkg_color[3]/255.)
            # glClearColor(0., 0., 1., 0)
            glClear(GL_COLOR_BUFFER_BIT | GL_DEPTH_BUFFER_BIT)
            glViewport(0, 0, int(width), int(height))

            glMatrixMode(GL_PROJECTION)
            glLoadIdentity()
            glOrtho(self.xmin, self.xmax, self.ymin, self.ymax, -99999, 99999)

            glMatrixMode(GL_MODELVIEW)
            glLoadIdentity()

            # dessin du background
            self._plotting(draw_type.WMSBACK)

            # Dessin des matrices
            self._plotting(draw_type.ARRAYS)

            # Dessin des résultats 2D
            self._plotting(draw_type.RES2D)

            # Dessin des vecteurs
            self._plotting(draw_type.VECTORS)

            # Dessin des tuiles
            self._plotting(draw_type.TILES)

            if self.active_vector is not None:
                if self.active_vector.parentzone is None:
                    # we must plot this vector because it is a temporary vector outside any zone
                    self.active_vector.plot()

            # Dessin des triangulations
            self._plotting(draw_type.TRIANGULATION)

            # Dessin des nuages
            self._plotting(draw_type.CLOUD)

            # Dessin des vues
            self._plotting(draw_type.VIEWS)

            # Dessin des "particule systems"
            self._plotting(draw_type.PARTICLE_SYSTEM)

            # Dessin du reste
            self._plotting(draw_type.OTHER)

            # Dessin du Front
            self._plotting(draw_type.WMSFORE)

            # Gestion des BC (si actif)
            if self.active_bc is not None:
                self.active_bc.plot()
            # try:
            #     if self.active_bc is not None:
            #         self.active_bc.plot()
            # except:
            #     pass

            # glFlush()
            self.canvas.SwapBuffers()
        else:
            raise NameError(
                'Opengl setcurrent -- maybe a conflict with an existing opengl32.dll file - please rename the opengl32.dll in the libs directory and retry')

    def OnPaint(self, e):
        """ event handler for paint event"""

        self.Paint()
        if e is not None:
            e.Skip()

    def findminmax(self, force=False):
        """ Find min/max of all objects """

        # FIXME : use iterator

        xmin = 1.e30
        ymin = 1.e30
        xmax = -1.e30
        ymax = -1.e30

        k = 0
        for locarray in self.myarrays:
            if locarray.plotted or force:
                xmin = min(locarray.origx + locarray.translx, xmin)
                xmax = max(locarray.origx + locarray.translx + float(locarray.nbx) * locarray.dx, xmax)
                ymin = min(locarray.origy + locarray.transly, ymin)
                ymax = max(locarray.origy + locarray.transly + float(locarray.nby) * locarray.dy, ymax)
                k += 1

        for locvector in self.myvectors:
            if locvector.plotted or force:
                if locvector.idx != 'grid':
                    locvector.find_minmax()
                    if isinstance(locvector,Zones):
                        xmin = min(locvector.xmin, xmin)
                        xmax = max(locvector.xmax, xmax)
                        ymin = min(locvector.ymin, ymin)
                        ymax = max(locvector.ymax, ymax)
                    elif isinstance(locvector,Bridges):
                        xmin = min(locvector.xmin, xmin)
                        xmax = max(locvector.xmax, xmax)
                        ymin = min(locvector.ymin, ymin)
                        ymax = max(locvector.ymax, ymax)
                    elif isinstance(locvector,crosssections):
                        xmin = min(locvector.xmin, xmin)
                        xmax = max(locvector.xmax, xmax)
                        ymin = min(locvector.ymin, ymin)
                        ymax = max(locvector.ymax, ymax)
                    k += 1

        for locvector in self.mytiles:
            if locvector.plotted or force:
                if locvector.idx != 'grid':
                    locvector.find_minmax()
                    if isinstance(locvector,Zones):
                        xmin = min(locvector.xmin, xmin)
                        xmax = max(locvector.xmax, xmax)
                        ymin = min(locvector.ymin, ymin)
                        ymax = max(locvector.ymax, ymax)
                    elif isinstance(locvector,Bridges):
                        xmin = min(locvector.xmin, xmin)
                        xmax = max(locvector.xmax, xmax)
                        ymin = min(locvector.ymin, ymin)
                        ymax = max(locvector.ymax, ymax)
                    elif isinstance(locvector,crosssections):
                        xmin = min(locvector.xmin, xmin)
                        xmax = max(locvector.xmax, xmax)
                        ymin = min(locvector.ymin, ymin)
                        ymax = max(locvector.ymax, ymax)
                    k += 1

        for loccloud in self.myclouds:
            if loccloud.plotted or force:
                loccloud.find_minmax(force)
                xmin = min(loccloud.xbounds[0], xmin)
                xmax = max(loccloud.xbounds[1], xmax)
                ymin = min(loccloud.ybounds[0], ymin)
                ymax = max(loccloud.ybounds[1], ymax)
                k += 1

        for loctri in self.mytri:
            if loctri.plotted or force:
                loctri.find_minmax(force)
                xmin = min(loctri.minx, xmin)
                xmax = max(loctri.maxx, xmax)
                ymin = min(loctri.miny, ymin)
                ymax = max(loctri.maxy, ymax)
                k += 1

        for locres2d in self.myres2D:
            locres2d:Wolfresults_2D
            if locres2d.plotted or force:
                locres2d.find_minmax(force)
                xmin = min(locres2d.xmin, xmin)
                xmax = max(locres2d.xmax, xmax)
                ymin = min(locres2d.ymin, ymin)
                ymax = max(locres2d.ymax, ymax)
                k += 1

        for locps in self.mypartsystems:
            locps:Particle_system
            if locps.plotted or force:
                locps.find_minmax(force)
                xmin = min(locps.xmin, xmin)
                xmax = max(locps.xmax, xmax)
                ymin = min(locps.ymin, ymin)
                ymax = max(locps.ymax, ymax)
                k += 1

        for locview in self.myviews:
            locview.find_minmax(force)
            xmin = min(locview.xmin, xmin)
            xmax = max(locview.xmax, xmax)
            ymin = min(locview.ymin, ymin)
            ymax = max(locview.ymax, ymax)
            k += 1

        for locothers in self.myothers:
            if type(locothers) in [genericImagetexture]: #, hydrometry_wolfgui]:
                xmin = min(locothers.xmin, xmin)
                xmax = max(locothers.xmax, xmax)
                ymin = min(locothers.ymin, ymin)
                ymax = max(locothers.ymax, ymax)
                k += 1
            elif type(locothers) in [PlansTerrier]: #, hydrometry_wolfgui]:
                if locothers.initialized:
                    xmin = min(locothers.xmin, xmin)
                    xmax = max(locothers.xmax, xmax)
                    ymin = min(locothers.ymin, ymin)
                    ymax = max(locothers.ymax, ymax)
                    k += 1

        if k > 0:
            self.xmin = xmin
            self.xmax = xmax
            self.ymin = ymin
            self.ymax = ymax

    def resizeFrame(self, w:int, h:int):
        """ Resize the frame

        :param w: width in pixels
        :param h: height in pixels
        """

        self.SetClientSize(w, h)

    def mimicme(self):
        """
        Report des caractéristiques de la fenêtre sur les autres éléments liés
        """

        if self.linked and self.forcemimic:
            if not self.linkedList is None:
                width, height = self.GetClientSize()

                curFrame: WolfMapViewer
                for curFrame in self.linkedList:
                    curFrame.forcemimic = False

                for curFrame in self.linkedList:
                    if curFrame != self:
                        curFrame.resizeFrame(width, height)
                        curFrame.mousex = self.mousex
                        curFrame.mousey = self.mousey
                        curFrame.sx = self.sx
                        curFrame.sy = self.sy
                        curFrame.width = self.width
                        curFrame.height = self.height
                        curFrame.setbounds()

                        if curFrame.link_shareopsvect:
                            curFrame.Active_vector(self.active_vector)
                            curFrame.active_array.myops.Active_vector(self.active_vector, False)
                            curFrame.action = self.action

                for curFrame in self.linkedList:
                    curFrame.forcemimic = True

    def mimicme_copyfrom(self):
        if self.linked and self.forcemimic:
            if not self.linkedList is None:
                width, height = self.GetClientSize()

                curFrame: WolfMapViewer
                for curFrame in self.linkedList:
                    curFrame.forcemimic = False

                for curFrame in self.linkedList:
                    if curFrame != self:
                        curFrame.copyfrom = self.copyfrom

                for curFrame in self.linkedList:
                    curFrame.forcemimic = True

    def Active_vector(self, vect):
        """ Active un vecteur et sa zone parent si existante """

        self.active_vector = vect

        if vect is not None:
            logging.info(_('Activating vector : ' + vect.myname))
            if vect.parentzone is not None:
                self.Active_zone(vect.parentzone)

        self.mimicme()

    def Active_zone(self, zone: zone):
        """ Active une zone et son parent si existant """

        self.active_zone = zone
        self.active_zones = zone.parent
        logging.info(_('Activating zone : ' + zone.myname))

    def list_background(self):
        return [cur.idx for cur in self.mywmsback]

    def list_foreground(self):
        return [cur.idx for cur in self.mywmsfore]

    def check_id(self, id=str, gridsize = 100.):
        """ Check an element from its id """

        curobj = self.getobj_from_id(id)

        if curobj is None:
            logging.warning('Bad id')
            return

        curobj.check_plot()

        curitem = self.gettreeitem(curobj)
        self.treelist.CheckItem(curitem, True)

        if id == 'grid':
            curobj.creategrid(gridsize, self.xmin, self.ymin, self.xmax, self.ymax)

    def uncheck_id(self, id=str, unload=True, forceresetOGL=True, askquestion=False):
        """ Uncheck an element from its id """

        curobj = self.getobj_from_id(id)

        if curobj is None:
            logging.warning('Bad id')
            return

        if issubclass(type(curobj), WolfArray):
            curobj.uncheck_plot(unload, forceresetOGL, askquestion)
        else:
            curobj.uncheck_plot()

        curitem = self.gettreeitem(curobj)
        self.treelist.UncheckItem(curitem)

    def get_current_zoom(self):
        """
        Get the current zoom

        :return: dict with keys 'center', 'xmin', 'xmax', 'ymin', 'ymax', 'width', 'height'

        """

        return {'center': (self.mousex, self.mousey),
                'xmin' :  self.xmin,
                'xmax' :  self.xmax,
                'ymin' :  self.ymin,
                'ymax' :  self.ymax,
                'width' : self.xmax-self.xmin,
                'height' : self.ymax-self.ymin}

    def save_current_zoom(self, filepath):
        """ Save the current zoom in a json file """

        zoom = self.get_current_zoom()
        with open(filepath, 'w') as fp:
            json.dump(zoom, fp)

    def read_current_zoom(self, filepath):
        """ Read the current zoom from a json file """

        if exists(filepath):
            with open(filepath, 'r') as fp:
                zoom = json.load(fp)

            self.zoom_on(zoom)


class Comp_Type(Enum):
    ARRAYS = 1
    ARRAYS_MB = 2
    RES2D = 3
    RES2D_GPU = 4
class Compare_Arrays_Results():

    def __init__(self, parent:WolfMapViewer = None, share_cmap_array:bool = False, share_cmap_diff:bool = False):

        self.parent = parent

        self.paths = []
        self.elements = []
        self.linked_elts = []
        self.diff = []
        self.mapviewers = []
        self.mapviewers_diff = []

        self.times = None

        self.share_cmap_array = share_cmap_array
        self.share_cmap_diff  = share_cmap_diff

        self.type = Comp_Type.ARRAYS

        self._initialized_viewers = False
        self.independent = True

    def _check_type(self, file:Path):
        """
        Check the type of the file/directory

        If it is a file and suffix is empty, it is considered as RES2D.
        If it is a directory and contains a simul_gpu_results, it is considered as RES2D_GPU.
        If it is a file and suffix is not empty, it is considered as ARRAYS. A check is done to see if it is a multi-block array.

        """

        file = Path(file)

        if file.suffix == '' and not file.is_dir():

            return Comp_Type.RES2D, file

        elif (file.parent / 'simul_gpu_results').exists():
            file = file.parent / 'simul_gpu_results'
            return Comp_Type.RES2D_GPU, file

        elif (file.parent.parent / 'simul_gpu_results').exists():
            file = file.parent.parent / 'simul_gpu_results'
            return Comp_Type.RES2D_GPU, file

        else:
            if file.suffix in ('.bin', '.tif', '.tiff', '.npy', '.npz', '.top', '.frott', '.nap', '.hbin', '.hbinb', '.qxbin', '.qxbinb', '.qybin', '.qybinb', '.inf') :

                if file.suffix in ('.bin', '.top', '.frott', '.nap', '.hbin', '.hbinb', '.qxbin', '.qxbinb', '.qybin', '.qybinb', '.inf'):
                    if file.with_suffix(file.suffix + '.txt').exists():

                        test = WolfArray(file, preload=False)
                        test.read_txt_header()

                        mb = test.nb_blocks > 0

                        if mb:
                            return Comp_Type.ARRAYS_MB, file

                return Comp_Type.ARRAYS, file
            else:
                return None, None

    def add(self, file_or_dir:Union[str, Path] = None):

        if file_or_dir is None:

            filterProject = "all (*.*)|*.*"
            file = wx.FileDialog(None, "Choose array/model", wildcard=filterProject)

            if file.ShowModal() == wx.ID_CANCEL:
                file.Destroy()
                return False
            else:
                filename = Path(file.GetPath())
                file.Destroy()

        self.paths.append(self._check_type(filename))

        if self.paths[-1][0] is None:
            logging.warning(_('File type not recognized -- Retry !'))
            self.paths.pop()
            return False

        return True

    def check(self):
        """ Check the consystency of the elements to compare """

        reftype = self.paths[0][0]
        for cur in self.paths:
            if cur[0] != reftype:
                logging.warning(_('Inconsistency in the type of the elements to compare'))
                return False

        return True

    def update_comp(self, idx=list[int]):
        """
        Update Arrays from 2D modellings

        :param idx: indexes of the time step to update --> steps to read

        """
        assert self.type in (Comp_Type.RES2D, Comp_Type.RES2D_GPU), 'This method is only for 2D results'

        self.linked_elts = []

        for curelt, curstep in zip(self.elements, idx):

            curelt.read_oneresult(curstep)

            self.linked_elts.append(curelt.as_WolfArray())

            for curelt, curlink in zip(self.elements, self.linked_elts):
                curlink.idx = curelt.idx + ' ' + curelt.get_currentview().value

        self.set_diff()

        if self._initialized_viewers:
            self.update_viewers()

    def update_type_result(self, newtype):
        """
        Update the result type for each element

        """
        assert newtype in views_2D, 'This type is not a 2D result'
        assert self.type in (Comp_Type.RES2D, Comp_Type.RES2D_GPU), 'This method is only for 2D results'

        for curelt in self.elements:
            curelt.set_currentview(newtype, force_updatepal = True)

        # remove elements
        for baselt, curelt, curmap in zip(self.elements, self.linked_elts, self.mapviewers):

            curmap.removeobj_from_id(curelt.idx)

        for curdiff, curmap in zip(self.diff, self.mapviewers_diff):

            curmap.removeobj_from_id(curdiff.idx)


        self.update_comp(self.times.get_times_idx())

    def set_elements(self):
        """ Set the elements to compare with the right type """
        from .ui.wolf_times_selection_comparison_models import Times_Selection

        if self.check():
            self.type = self.paths[0][0]

            if self.type == Comp_Type.RES2D_GPU:
                self.parent.menu_wolf2d()
                self.elements = [wolfres2DGPU(cur[1], plotted=False, idx = cur[1].name + '_' + str(idx)) for idx, cur in enumerate(self.paths)]

                times = [curmod.get_times_steps()[0] for curmod in self.elements]

                self.times = Times_Selection(self, wx.ID_ANY, _("Times"), size=(400,400), times = times, callback = self.update_comp)
                self.times.Show()

            elif self.type == Comp_Type.RES2D:
                self.parent.menu_wolf2d()
                self.elements = [Wolfresults_2D(cur[1], plotted=False, idx = cur[1].name + '_' + str(idx)) for idx, cur in enumerate(self.paths)]

                times = [curmod.get_times_steps()[0] for curmod in self.elements]

                self.times = Times_Selection(self, wx.ID_ANY, _("Times"), size=(400,400), times = times, callback = self.update_comp)
                self.times.Show()

            elif self.type == Comp_Type.ARRAYS:
                self.elements = [WolfArray(cur[1], plotted=False, idx = cur[1].name + '_' + str(idx)) for idx, cur in enumerate(self.paths)]

            elif self.type == Comp_Type.ARRAYS_MB:
                self.elements = [WolfArrayMB(cur[1], plotted=False, idx = cur[1].name + '_' + str(idx)) for idx, cur in enumerate(self.paths)]

    def set_diff(self):
        """ Set the differential between the elements and the first one, which is the reference """

        if self.type in (Comp_Type.ARRAYS, Comp_Type.ARRAYS_MB):

            ref = self.elements[0]

            # Recherche d'un masque union des masques partiels
            ref.mask_unions(self.elements[1:])

            # Création du différentiel -- Les opérateurs mathématiques sont surchargés
            self.diff = [cur - ref for cur in self.elements[1:]]

            for curdiff, cur in zip(self.diff, self.elements[1:]):
                curdiff.idx = _('Difference') + cur.idx +' - ' + ref.idx

        elif self.type in (Comp_Type.RES2D, Comp_Type.RES2D_GPU):

            if len(self.linked_elts) == 0:
                self.update_comp([-1] * len(self.elements))

            elif len(self.linked_elts) == len(self.elements):
                ref = self.linked_elts[0]

                self.diff = [cur - ref for cur in self.linked_elts[1:]]

                for curdiff, cur in zip(self.diff, self.linked_elts[1:]):
                    curdiff.idx = _('Difference') + cur.idx +' - ' + ref.idx

    def set_viewers(self, independent:bool = None):
        """
        Set viewers

        """

        if independent is None:
            dlg = wx.MessageDialog(None, _("Create a viewer for each element ?"), _("Viewers"), style = wx.YES_NO|wx.YES_DEFAULT)
            ret = dlg.ShowModal()

            self.independent = ret == wx.ID_YES
        else:
            self.independent = independent

        if not self.independent:
            self.mapviewers = [self.parent] * len(self.elements)
            self.mapviewers_diff = self.mapviewers
        else:
            # Création de plusieurs fenêtres de visualisation basées sur la classe "WolfMapViewer"
            self.mapviewers = []

            self.mapviewers.append(self.parent) # parent as viewer for first element

            for id, file in enumerate(self.elements[1:]):

                self.mapviewers.append(WolfMapViewer(None, file.idx, w=600, h=600, wxlogging=self.parent.wxlogging, wolfparent = self.parent.wolfparent))
                self.mapviewers_diff.append(WolfMapViewer(None, 'Difference' + file.idx, w=600, h=600, wxlogging=self.parent.wxlogging, wolfparent = self.parent.wolfparent))

            for curviewer in self.mapviewers[1:] + self.mapviewers_diff:
                curviewer.add_grid()
                curviewer.add_WMS()

        for curviewer in self.mapviewers + self.mapviewers_diff:
            curviewer.linked = True
            curviewer.linkedList = self.mapviewers + self.mapviewers_diff

        self._initialized_viewers = True

        self.update_viewers()

    def set_shields_param(self, diamsize:float = .001, graindensity:float = 2.65):
        """ Set the parameters for the shields diagram """

        for curelt in self.elements:
            curelt.sediment_diameter = diamsize
            curelt.sediment_density = graindensity
            curelt.load_default_colormap('shields_cst')

    def update_viewers(self):
        """ Update the viewers with the new elements """

        if self.type in (Comp_Type.ARRAYS, Comp_Type.ARRAYS_MB):
            elts = self.elements
        elif self.type in (Comp_Type.RES2D, Comp_Type.RES2D_GPU):
            elts = self.linked_elts

        # on attribue une matrice par interface graphique
        ref = elts[0]
        for baselt, curelt, curmap in zip(self.elements, elts, self.mapviewers):

            # if self.type in (Comp_Type.RES2D, Comp_Type.RES2D_GPU):
            #     curmap.active_res2d = baselt

            curmap.removeobj_from_id(curelt.idx)

            curelt.change_gui(curmap)
            curmap.active_array = curelt
            curelt.myops.myzones = ref.myops.myzones

        # diff = self.diff[0]
        for curdiff, curmap in zip(self.diff, self.mapviewers_diff):

            curmap.removeobj_from_id(curdiff.idx)

            curdiff.change_gui(curmap)
            curmap.active_array = curdiff
            curdiff.myops.myzones = ref.myops.myzones

        # on partage la palette de couleurs
        ref.mypal.automatic = False
        ref.myops.palauto.SetValue(0)

        if self.share_cmap_array:
            for curelt in elts[1:]:
                ref.add_crosslinked_array(curelt)
                ref.share_palette()
        else:
            for curelt in elts[1:]:
                curelt.mypal.automatic = False
                curelt.myops.palauto.SetValue(0)
                curelt.mypal.updatefrompalette(ref.mypal)

        #palette de la différence
        diff = self.diff[0]
        diff.mypal = wolfpalette()
        if isinstance(diff, WolfArrayMB):
            diff.link_palette()

        path = os.path.dirname(__file__)
        fn = join(path, 'models\\diff16.pal')

        diff.mypal.readfile(fn)
        diff.mypal.automatic = False
        diff.myops.palauto.SetValue(0)

        if self.share_cmap_diff:
            for curelt in self.diff[1:]:
                diff.add_crosslinked_array(curelt)
                diff.share_palette()
        else:
            for curelt in self.diff[1:]:
                curelt.mypal.automatic = False
                curelt.myops.palauto.SetValue(0)
                curelt.mypal.updatefrompalette(diff.mypal)

        # Ajout des matrices dans les fenêtres de visualisation
        for curelt, curmap in zip(elts, self.mapviewers):
            curmap.add_object('array', newobj = curelt, ToCheck = True, id = curelt.idx)

        for curdiff, curmap in zip(self.diff, self.mapviewers_diff):
            curmap.add_object('array', newobj = curdiff, ToCheck = True, id = curdiff.idx)

        if self.independent:
            for curmap in self.mapviewers + self.mapviewers_diff:
                curmap.Refresh()
        else:
            self.mapviewers[0].Refresh()

    def bake(self):

        self.set_elements()
        self.set_diff()
        self.set_viewers()